"""
Document Tools MCP Server
Exposes Word, PowerPoint, and Excel document creation as MCP tools.
Generated files are saved to OUTPUT_DIR with unique IDs.

Requires: pip install python-docx python-pptx openpyxl
Start with: python -m mcp.documents.document_mcp
"""

import ipaddress
import json
import logging
import os
import socket
import uuid
from pathlib import Path
from urllib.parse import urlparse
from xml.etree import ElementTree

from mcp.server import MCPServer
from starlette.responses import JSONResponse

from portal.platform.data_loader import load_data
from portal.platform.mcp_host.owui_files import publish_file_sync

port = int(os.getenv("DOCUMENTS_MCP_PORT", "8913"))
mcp = MCPServer("document-tools")

OUTPUT_DIR = Path(os.getenv("OUTPUT_DIR", "data/generated"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
TEMPLATES_DIR = Path(__file__).resolve().parents[4] / "config" / "documents" / "templates"


def _published(output_path: Path, noun: str) -> dict:
    """Publish a freshly generated file through Open WebUI and shape the tool result."""
    pub = publish_file_sync(output_path)
    if "error" in pub:
        return {"success": False, "error": pub["error"]}
    return {
        "success": True,
        "path": str(output_path),
        "filename": pub["filename"],
        "download_url": pub["url"],
        "message": f"{noun} created: {pub['filename']}. [Download]({pub['url']})",
    }


def _with_optional_pdf(result: dict, output_path: Path, also_pdf: bool) -> dict:
    if not also_pdf or not result.get("success"):
        return result
    pdf = export_pdf(str(output_path))
    if pdf.get("success"):
        result["pdf"] = pdf
    else:
        result["pdf_note"] = pdf.get("error", "PDF export unavailable")
    return result


def _load_template(template: str | None) -> dict:
    if not template:
        return {}
    safe = "".join(c for c in template if c.isalnum() or c in ("-", "_"))
    path = TEMPLATES_DIR / f"{safe}.json"
    if not safe or not path.is_file():
        raise ValueError(f"Unknown document template: {template}")
    return json.loads(path.read_text(encoding="utf-8"))


def _add_hyperlink(paragraph, text: str, url: str):
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    relationship = paragraph.part.relate_to(
        url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", True
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship)
    run = OxmlElement("w:r")
    props = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "0563C1")
    props.append(color)
    run.append(props)
    node = OxmlElement("w:t")
    node.text = text
    run.append(node)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _render_inline(paragraph, node) -> None:
    if node.text:
        paragraph.add_run(node.text)
    for child in node:
        text = "".join(child.itertext())
        if child.tag == "a":
            _add_hyperlink(paragraph, text, child.attrib.get("href", ""))
        else:
            run = paragraph.add_run(text)
            run.bold = child.tag in {"strong", "b"}
            run.italic = child.tag in {"em", "i"}
            if child.tag == "code":
                run.font.name = "Courier New"
        if child.tail:
            paragraph.add_run(child.tail)


def _render_markdown(doc, content: str) -> None:
    import markdown
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt

    html = markdown.markdown(content, extensions=["extra", "sane_lists"])
    root = ElementTree.fromstring(f"<root>{html}</root>")

    def render_list(node, level: int = 0, ordered: bool = False) -> None:
        for item in node.findall("li"):
            para = doc.add_paragraph(style="List Number" if ordered else "List Bullet")
            para.paragraph_format.left_indent = Pt(18 * level)
            _render_inline(para, item)
            for child in item:
                if child.tag in {"ul", "ol"}:
                    render_list(child, level + 1, child.tag == "ol")

    for node in root:
        if node.tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            doc.add_heading("".join(node.itertext()), level=int(node.tag[1]))
        elif node.tag == "p":
            _render_inline(doc.add_paragraph(), node)
        elif node.tag in {"ul", "ol"}:
            render_list(node, ordered=node.tag == "ol")
        elif node.tag == "blockquote":
            para = doc.add_paragraph(style="Intense Quote")
            _render_inline(para, node[0] if len(node) else node)
        elif node.tag == "pre":
            para = doc.add_paragraph()
            run = para.add_run("".join(node.itertext()))
            run.font.name = "Courier New"
            shading = OxmlElement("w:shd")
            shading.set(qn("w:fill"), "F2F2F2")
            para._p.get_or_add_pPr().append(shading)
        elif node.tag == "table":
            rows = node.findall(".//tr")
            width = max((len(row) for row in rows), default=1)
            table = doc.add_table(rows=len(rows), cols=width)
            table.style = "Table Grid"
            for row_index, row in enumerate(rows):
                for col_index, cell in enumerate(row):
                    target = table.cell(row_index, col_index)
                    target.text = "".join(cell.itertext())
                    if cell.tag == "th":
                        for run in target.paragraphs[0].runs:
                            run.bold = True
        elif node.tag == "hr":
            para = doc.add_paragraph("―" * 20)
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER


def _safe_remote_image(source: str) -> Path:
    import httpx

    parsed = urlparse(source)
    if parsed.scheme != "https" or not parsed.hostname:
        raise ValueError("Remote images must use https")
    for info in socket.getaddrinfo(parsed.hostname, 443, type=socket.SOCK_STREAM):
        address = ipaddress.ip_address(info[4][0])
        if not address.is_global:
            raise ValueError("Private or local image addresses are not allowed")
    suffix = Path(parsed.path).suffix.lower()
    if suffix not in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}:
        suffix = ".img"
    destination = _unique_path("embedded_image", suffix.lstrip("."))
    with httpx.Client(timeout=30, follow_redirects=False) as client:
        response = client.get(source)
        response.raise_for_status()
        if len(response.content) > 25 * 1024 * 1024:
            raise ValueError("Remote image exceeds 25 MiB")
        destination.write_bytes(response.content)
    return destination


def _resolve_image(source: str) -> Path:
    if source.startswith(("http://", "https://")):
        return _safe_remote_image(source)
    candidate = Path(source)
    if not candidate.is_absolute():
        candidate = OUTPUT_DIR / candidate
    resolved = candidate.resolve()
    root = OUTPUT_DIR.resolve()
    if not resolved.is_file() or not resolved.is_relative_to(root):
        raise ValueError("Local images must be files within the output directory")
    return resolved


@mcp.custom_route("/health", methods=["GET"])
async def health_check(request):
    return JSONResponse({"status": "ok", "service": "documents-mcp"})


# Tool manifest for discovery
TOOLS_MANIFEST = load_data("config/inference", "tools_manifest_document_mcp")


@mcp.custom_route("/tools", methods=["GET"])
async def list_tools(request):
    return JSONResponse({"tools": TOOLS_MANIFEST})


@mcp.custom_route("/tools/{tool_name}", methods=["POST"])
async def invoke_tool(request):
    """REST dispatch endpoint used by portal-pipeline tool_registry."""
    tool_name = request.path_params.get("tool_name", "")
    try:
        body = await request.json()
    except Exception:
        body = {}
    arguments = body.get("arguments", body)

    try:
        fn = globals().get(tool_name)
        if fn is None or not callable(fn):
            return JSONResponse({"error": f"Unknown tool: {tool_name}"}, status_code=404)
        result = fn(**arguments)
        return JSONResponse(result)
    except TypeError as e:
        return JSONResponse({"error": f"Invalid arguments for '{tool_name}': {e}"}, status_code=400)
    except Exception as e:
        logger.error("invoke_tool %s failed: %s", tool_name, e, exc_info=True)
        return JSONResponse({"error": str(e)}, status_code=500)


logger = logging.getLogger(__name__)


def _unique_path(name: str, ext: str) -> Path:
    """Return a unique output path."""
    uid = uuid.uuid4().hex[:8]
    safe = "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in name[:40]).strip("_")
    return OUTPUT_DIR / f"{safe}_{uid}.{ext}"


@mcp.tool()
def create_word_document(
    title: str,
    content: str,
    author: str = "Portal AI",
    images: list[dict] | None = None,
    template: str | None = None,
    also_pdf: bool = False,
) -> dict:
    """
    Create a Word (.docx) document from a title and markdown-style content.

    Content supports:
    - '# Heading' → H1 heading
    - '## Heading' → H2 heading
    - '### Heading' → H3 heading
    - '- item' → bullet list item
    - Regular text → body paragraph

    Args:
        title: Document title (also used as filename base)
        content: Document body; supports basic markdown headings and bullets
        author: Author name for document metadata (default "Portal AI")

    Returns:
        dict with success, path (server path), and filename
    """
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, Pt, RGBColor
    except ImportError:
        return {
            "success": False,
            "error": "python-docx not installed. Run: pip install python-docx",
        }

    try:
        doc = Document()
        theme = _load_template(template)
        doc.core_properties.author = author
        doc.core_properties.title = title

        styles = doc.styles
        if theme.get("body_font"):
            styles["Normal"].font.name = theme["body_font"]
        if theme.get("heading_font"):
            for level in range(1, 7):
                styles[f"Heading {level}"].font.name = theme["heading_font"]
        if theme.get("accent_color"):
            color = RGBColor.from_string(theme["accent_color"].lstrip("#"))
            for level in range(1, 7):
                styles[f"Heading {level}"].font.color.rgb = color

        # Title heading
        heading = doc.add_heading(title, level=0)
        heading.runs[0].font.size = Pt(24)

        _render_markdown(doc, content)

        for image in images or []:
            source = image.get("source") or image.get("path_or_url")
            if not source:
                continue
            doc.add_picture(
                str(_resolve_image(str(source))), width=Inches(float(image.get("width_in", 6)))
            )
            if image.get("caption"):
                caption = doc.add_paragraph(str(image["caption"]))
                caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in caption.runs:
                    run.italic = True
                    run.font.size = Pt(9)

        if theme.get("footer"):
            footer = doc.sections[0].footer.paragraphs[0]
            footer.text = theme["footer"]
            footer.alignment = WD_ALIGN_PARAGRAPH.CENTER

        output_path = _unique_path(title, "docx")
        doc.save(str(output_path))
        return _with_optional_pdf(_published(output_path, "Document"), output_path, also_pdf)
    except Exception as e:
        logger.exception("Word document creation failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def create_powerpoint(
    title: str,
    slides: list[dict],
    author: str = "Portal AI",
    template: str | None = None,
    also_pdf: bool = False,
) -> dict:
    """
    Create a PowerPoint (.pptx) presentation.

    Each slide dict should have:
    - 'title': slide title (str)
    - 'content': slide body text or bullet points (str, newline-separated)
    - 'notes': speaker notes (str, optional)

    Args:
        title: Presentation title (used as filename base)
        slides: List of slide dicts with 'title', 'content', and optional 'notes'
        author: Author name for metadata (default "Portal AI")

    Returns:
        dict with success, path (server path), and filename
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed. Run: pip install python-pptx",
        }

    try:
        prs = Presentation()
        theme = _load_template(template)
        prs.core_properties.author = author
        prs.core_properties.title = title

        # Title slide
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = title

        # Content slides
        content_layout = prs.slide_layouts[1]
        for slide_data in slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            for i, line in enumerate(slide_data.get("content", "").splitlines()):
                stripped = line.strip()
                if not stripped:
                    continue
                if i == 0:
                    tf.text = stripped
                else:
                    tf.add_paragraph().text = stripped

            if "notes" in slide_data and slide_data["notes"]:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data["notes"]

            chart_spec = slide_data.get("chart")
            if chart_spec:
                from pptx.chart.data import CategoryChartData
                from pptx.enum.chart import XL_CHART_TYPE

                chart_types = {
                    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "line": XL_CHART_TYPE.LINE,
                    "pie": XL_CHART_TYPE.PIE,
                }
                chart_data = CategoryChartData()
                chart_data.categories = chart_spec.get("categories", [])
                for series in chart_spec.get("series", []):
                    chart_data.add_series(series.get("name", ""), series.get("values", []))
                graphic_frame = slide.shapes.add_chart(
                    chart_types.get(
                        str(chart_spec.get("type", "bar")).lower(), XL_CHART_TYPE.COLUMN_CLUSTERED
                    ),
                    Inches(1),
                    Inches(2),
                    Inches(8),
                    Inches(4.5),
                    chart_data,
                )
                graphic_frame.chart.has_title = bool(chart_spec.get("title"))
                if chart_spec.get("title"):
                    graphic_frame.chart.chart_title.text_frame.text = str(chart_spec["title"])

            for image in slide_data.get("images", []):
                source = image.get("source") or image.get("path_or_url")
                if source:
                    slide.shapes.add_picture(
                        str(_resolve_image(str(source))),
                        Inches(float(image.get("left_in", 1))),
                        Inches(float(image.get("top_in", 2))),
                        width=Inches(float(image.get("width_in", 6))),
                    )

            if theme.get("accent_color") and slide.shapes.title:
                color = RGBColor.from_string(theme["accent_color"].lstrip("#"))
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = color
                        if theme.get("heading_font"):
                            run.font.name = theme["heading_font"]

        output_path = _unique_path(title, "pptx")
        prs.save(str(output_path))
        return _with_optional_pdf(_published(output_path, "Presentation"), output_path, also_pdf)
    except Exception as e:
        logger.exception("Presentation creation failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def create_excel(
    title: str,
    data: list | None = None,
    sheets: list[dict] | None = None,
    sheet_name: str = "Sheet1",
    charts: list[dict] | None = None,
    also_pdf: bool = False,
) -> dict:
    """
    Create an Excel (.xlsx) spreadsheet.

    Simple usage (flat rows):
        data: List of rows (first row treated as headers), e.g. [["Name","Score"],["Alice",95]]
        sheet_name: Tab name (default "Sheet1")

    Advanced usage (multiple sheets):
        sheets: List of sheet dicts, each with 'name', 'headers', and 'rows'

    Args:
        title: Spreadsheet title (used as filename base)
        data: Simple list of rows (first row is headers)
        sheets: Advanced — list of sheet dicts with 'name', 'headers', 'rows'
        sheet_name: Sheet tab name when using data parameter (default "Sheet1")

    Returns:
        dict with success, path (server path), and filename
    """
    try:
        import openpyxl
        from openpyxl.styles import Font
    except ImportError:
        return {
            "success": False,
            "error": "openpyxl not installed. Run: pip install openpyxl",
        }

    try:
        wb = openpyxl.Workbook()
        wb.properties.title = title

        if data is not None:
            # Simple mode: flat list of rows
            ws = wb.active
            ws.title = sheet_name
            for i, row in enumerate(data):
                ws.append(row)
                if i == 0:
                    # Bold first row (headers)
                    for cell in ws[1]:
                        cell.font = Font(bold=True)
        elif sheets is not None:
            default_sheet = wb.active
            first = True
            for sheet_data in sheets:
                if first:
                    ws = default_sheet
                    ws.title = sheet_data.get("name", "Sheet1")
                    first = False
                else:
                    ws = wb.create_sheet(title=sheet_data.get("name", "Sheet"))

                headers = sheet_data.get("headers", [])
                if headers:
                    ws.append(headers)
                    for cell in ws[1]:
                        cell.font = Font(bold=True)

                for row in sheet_data.get("rows", []):
                    ws.append(row)
        else:
            return {"success": False, "error": "Provide either 'data' or 'sheets' parameter"}

        if charts:
            from openpyxl.chart import BarChart, LineChart, PieChart, Reference

            chart_types = {"bar": BarChart, "line": LineChart, "pie": PieChart}
            for spec in charts:
                target = wb[spec["sheet"]] if spec.get("sheet") in wb.sheetnames else wb.active
                chart = chart_types.get(str(spec.get("type", "bar")).lower(), BarChart)()
                chart.title = spec.get("title", "")
                chart.add_data(
                    Reference(
                        target,
                        min_col=int(spec["min_col"]),
                        max_col=int(spec["max_col"]),
                        min_row=int(spec["min_row"]),
                        max_row=int(spec["max_row"]),
                    ),
                    titles_from_data=True,
                )
                if spec.get("cats_col"):
                    chart.set_categories(
                        Reference(
                            target,
                            min_col=int(spec["cats_col"]),
                            max_col=int(spec["cats_col"]),
                            min_row=int(spec["min_row"]) + 1,
                            max_row=int(spec["max_row"]),
                        )
                    )
                target.add_chart(chart, spec.get("anchor", "H2"))

        output_path = _unique_path(title, "xlsx")
        wb.save(str(output_path))
        return _with_optional_pdf(_published(output_path, "Spreadsheet"), output_path, also_pdf)
    except Exception as e:
        logger.exception("Spreadsheet creation failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def convert_document(
    source_path: str,
    target_format: str,
) -> dict:
    """Copy a document to a new format name.

    Note: True format conversion (e.g., .docx to .pdf) requires LibreOffice
    installed on the host. Without LibreOffice, this tool copies the file with
    the new extension, which is only useful for same-family formats.

    For PDF export, use LibreOffice on the host:
      libreoffice --headless --convert-to pdf <file>

    Args:
        source_path:   Path to the source document
        target_format: Target extension: 'pdf', 'docx', 'pptx', or 'xlsx'
    """
    import shutil
    import subprocess
    from pathlib import Path as _Path

    _allowed_target_formats = frozenset({"pdf", "docx", "pptx", "xlsx"})

    src = _Path(source_path).resolve()
    allowed_root = _Path(OUTPUT_DIR).resolve()
    if not str(src).startswith(str(allowed_root) + os.sep):
        return {"error": "source_path must be a file within the output directory"}
    if not src.exists():
        return {"error": f"Source file not found: {source_path}"}

    target_format = target_format.lower().lstrip(".")
    if target_format not in _allowed_target_formats:
        return {
            "error": (
                f"Unsupported target_format {target_format!r}. "
                f"Allowed: {sorted(_allowed_target_formats)}"
            )
        }
    out_path = _unique_path(src.stem, target_format)

    # Attempt LibreOffice conversion for cross-format (best quality)
    try:
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                target_format,
                "--outdir",
                str(out_path.parent),
                str(src),
            ],
            capture_output=True,
            timeout=60,
        )
        if result.returncode == 0:
            # LibreOffice writes to same dir as source — find the output file
            converted = src.parent / f"{src.stem}.{target_format}"
            if converted.exists():
                shutil.move(str(converted), str(out_path))
                return {**_published(out_path, "Converted document"), "method": "libreoffice"}
    except (FileNotFoundError, subprocess.TimeoutExpired):
        logger.debug(
            "LibreOffice conversion failed (not installed or timed out) — using copy fallback"
        )
    except OSError as e:
        logger.debug("LibreOffice subprocess error: %s", e)

    # Fallback: copy with new extension (only meaningful for same-family formats)
    same_family = {
        frozenset({"docx", "doc"}),
        frozenset({"pptx", "ppt"}),
        frozenset({"xlsx", "xls"}),
    }
    src_ext = src.suffix.lstrip(".")
    is_same_family = any({src_ext, target_format} <= fam for fam in same_family)

    if not is_same_family:
        return {
            "error": (
                f"Cannot convert {src_ext!r} → {target_format!r} without LibreOffice. "
                "Install LibreOffice for PDF and cross-format conversion."
            ),
            "install": "brew install libreoffice  # or apt-get install libreoffice",
        }

    shutil.copy2(str(src), str(out_path))
    return {
        **_published(out_path, "Converted document"),
        "method": "copy",
        "note": f"Copied {src_ext} → {target_format}. "
        "Install LibreOffice for true format conversion.",
    }


@mcp.tool()
def export_pdf(source_path: str) -> dict:
    """Export a generated Office document to PDF using host LibreOffice.

    Returns a clear error when LibreOffice is unavailable; it never renames or
    copies Office bytes into a misleading .pdf file.
    """
    return convert_document(source_path, "pdf")


@mcp.tool()
def prepare_embed_image(
    image_url: str,
    caption: str = "",
    width_in: float = 6.0,
) -> dict:
    """Build an image spec accepted by Word documents and PowerPoint slides."""
    return {"source": image_url, "caption": caption, "width_in": width_in}


@mcp.tool()
def list_generated_files() -> list[dict]:
    """List recently generated documents in the output directory."""
    files = []
    for f in sorted(OUTPUT_DIR.iterdir(), key=lambda x: x.stat().st_mtime, reverse=True)[:20]:
        if f.is_file():
            files.append(
                {
                    "filename": f.name,
                    "path": str(f),
                    "size_bytes": f.stat().st_size,
                    "type": f.suffix.lstrip("."),
                }
            )
    return files


@mcp.tool()
def read_word_document(
    file_path: str,
    include_tables: bool = True,
) -> dict:
    """Extract text content and structure from an existing Word (.docx) file.

    Returns headings, paragraphs, and table data.

    Args:
        file_path:      Absolute path to the .docx file to read
        include_tables: Whether to include table cell content (default True)

    Returns:
        dict with success, metadata, content (list of blocks), and optional tables
    """
    try:
        from docx import Document as _Document
    except ImportError:
        return {
            "success": False,
            "error": "python-docx not installed. Run: pip install python-docx",
        }

    src = Path(file_path).resolve()
    if not src.exists():
        return {"success": False, "error": f"File not found: {file_path}"}
    if src.suffix.lower() != ".docx":
        return {"success": False, "error": f"Expected .docx file, got: {src.suffix}"}

    try:
        doc = _Document(str(src))
        blocks: list[dict] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else "Normal"
            blocks.append({"type": style, "text": text})

        result: dict = {
            "success": True,
            "filename": src.name,
            "author": doc.core_properties.author or "",
            "title": doc.core_properties.title or "",
            "paragraph_count": len(blocks),
            "content": blocks,
        }

        if include_tables and doc.tables:
            tables_data = []
            for i, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                tables_data.append({"table_index": i, "rows": rows})
            result["tables"] = tables_data

        return result
    except Exception as e:
        logger.exception("Word document read failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def read_excel(
    file_path: str,
    max_rows: int = 500,
) -> dict:
    """Extract data from an existing Excel (.xlsx) spreadsheet.

    Returns sheet names and row data for each sheet.

    Args:
        file_path: Absolute path to the .xlsx file to read
        max_rows:  Maximum rows to return per sheet (default 500)

    Returns:
        dict with success, filename, sheet_count, and sheets (list of sheet dicts)
    """
    try:
        import openpyxl as _openpyxl
    except ImportError:
        return {"success": False, "error": "openpyxl not installed. Run: pip install openpyxl"}

    src = Path(file_path).resolve()
    if not src.exists():
        return {"success": False, "error": f"File not found: {file_path}"}
    if src.suffix.lower() not in {".xlsx", ".xlsm"}:
        return {"success": False, "error": f"Expected .xlsx/.xlsm file, got: {src.suffix}"}

    try:
        wb = _openpyxl.load_workbook(str(src), read_only=True, data_only=True)
        sheets_data = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[list] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if max_rows and i >= max_rows:
                    break
                rows.append([str(cell) if cell is not None else "" for cell in row])
            sheets_data.append(
                {
                    "name": sheet_name,
                    "row_count": len(rows),
                    "truncated": max_rows > 0 and len(rows) >= max_rows,
                    "rows": rows,
                }
            )
        wb.close()

        return {
            "success": True,
            "filename": src.name,
            "sheet_count": len(sheets_data),
            "sheets": sheets_data,
        }
    except Exception as e:
        logger.exception("Excel read failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def read_powerpoint(
    file_path: str,
) -> dict:
    """Extract text and speaker notes from an existing PowerPoint (.pptx) file.

    Returns slide titles, content blocks, and notes.

    Args:
        file_path: Absolute path to the .pptx file to read

    Returns:
        dict with success, filename, slide_count, and slides (list of slide dicts)
    """
    try:
        from pptx import Presentation as _Presentation
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed. Run: pip install python-pptx",
        }

    src = Path(file_path).resolve()
    if not src.exists():
        return {"success": False, "error": f"File not found: {file_path}"}
    if src.suffix.lower() != ".pptx":
        return {"success": False, "error": f"Expected .pptx file, got: {src.suffix}"}

    try:
        prs = _Presentation(str(src))
        slides_data = []
        for i, slide in enumerate(prs.slides):
            title = ""
            content_blocks: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                shape_text = shape.text_frame.text.strip()
                if not shape_text:
                    continue
                if shape.shape_id == 1 or (
                    hasattr(slide.shapes, "title") and shape == slide.shapes.title
                ):
                    title = shape_text
                else:
                    content_blocks.append(shape_text)

            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slides_data.append(
                {
                    "slide_number": i + 1,
                    "title": title,
                    "content": content_blocks,
                    "notes": notes,
                }
            )

        return {
            "success": True,
            "filename": src.name,
            "author": prs.core_properties.author or "",
            "slide_count": len(slides_data),
            "slides": slides_data,
        }
    except Exception as e:
        logger.exception("PowerPoint read failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def read_pdf(
    file_path: str,
    max_pages: int = 50,
    include_tables: bool = True,
) -> dict:
    """Extract text content from an existing PDF file page by page.

    Also extracts tables when present. Requires pdfplumber.

    Args:
        file_path:      Absolute path to the .pdf file to read
        max_pages:      Maximum pages to extract (default 50, 0 = all pages)
        include_tables: Whether to extract table data alongside text (default True)

    Returns:
        dict with success, filename, page_count, pages (list of page dicts),
        and optional tables per page
    """
    try:
        import pdfplumber as _pdfplumber
    except ImportError:
        return {
            "success": False,
            "error": "pdfplumber not installed. Run: pip install pdfplumber",
        }

    src = Path(file_path).resolve()
    if not src.exists():
        return {"success": False, "error": f"File not found: {file_path}"}
    if src.suffix.lower() != ".pdf":
        return {"success": False, "error": f"Expected .pdf file, got: {src.suffix}"}

    try:
        pages_data = []
        with _pdfplumber.open(str(src)) as pdf:
            total_pages = len(pdf.pages)
            limit = total_pages if max_pages == 0 else min(max_pages, total_pages)

            for i in range(limit):
                page = pdf.pages[i]
                text = (page.extract_text() or "").strip()
                page_dict: dict = {
                    "page_number": i + 1,
                    "text": text,
                    "char_count": len(text),
                }

                if include_tables:
                    raw_tables = page.extract_tables()
                    if raw_tables:
                        page_dict["tables"] = [
                            [[str(cell) if cell is not None else "" for cell in row] for row in tbl]
                            for tbl in raw_tables
                        ]

                pages_data.append(page_dict)

        return {
            "success": True,
            "filename": src.name,
            "total_pages": total_pages,
            "pages_extracted": len(pages_data),
            "truncated": len(pages_data) < total_pages,
            "pages": pages_data,
        }
    except Exception as e:
        logger.exception("PDF read failed")
        return {"success": False, "error": str(e)}


if __name__ == "__main__":
    mcp.run(transport="streamable-http", host="0.0.0.0", port=port)
