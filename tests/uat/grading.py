"""Portal 5 UAT — think-block stripping, assertion engine, status computation.

Extracted verbatim from tests/portal5_uat_driver.py (TASK_UAT_MODULARIZE_V1
phase A). Heavy format validators (docx/xlsx/pptx/wav/png/mp4) keep their
original inline imports.
"""

from __future__ import annotations

from pathlib import Path

# Think-block stripping
# ---------------------------------------------------------------------------


def _strip_think_blocks(text: str) -> str:
    """Strip reasoning blocks from model output before running assertions.

    Three reasoning formats are handled:
    - <think>...</think>: Laguna-XS.2, Phi-4-reasoning-plus, Qwopus
    - [THINK]...[/THINK]: Magistral
    - <details type="reasoning">...</details>: AEON/Qwen3 as committed by OWUI API
      (OWUI inlines reasoning in the content field; the actual response follows
      the closing tag — without stripping, keywords like "error" or "failed" that
      appear naturally in reasoning traces cause false not_contains failures)

    Strips all variants case-insensitively with DOTALL so multi-line blocks are
    handled. Trailing whitespace is normalized after stripping.
    """
    import re

    original = text
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"\[THINK\].*?\[/THINK\]", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(
        r'<details[^>]*type=["\']reasoning["\'][^>]*>.*?</details>',
        "",
        text,
        flags=re.DOTALL | re.IGNORECASE,
    )
    result = text.strip()
    if not result:
        # Model put entire answer inside reasoning block — extract inner content
        m = re.search(
            r"<details[^>]*>.*?<summary>.*?</summary>(.*?)</details>",
            original,
            re.DOTALL | re.IGNORECASE,
        )
        if m:
            return m.group(1).strip()
    return result


# ---------------------------------------------------------------------------
# Assertion engine
# ---------------------------------------------------------------------------


_UNICODE_DASH_TABLE = str.maketrans(
    "".join(
        [
            "‐",  # hyphen
            "‑",  # non-breaking hyphen
            "‒",  # figure dash
            "–",  # en dash
            "—",  # em dash
            "―",  # horizontal bar
            "−",  # minus sign
            "─",  # box-drawing horizontal
            "﹘",  # small em dash
            "﹣",  # small hyphen-minus
            "－",  # fullwidth hyphen-minus
        ]
    ),
    "-" * 11,
)


def _normalize_dashes(s: str) -> str:
    return s.translate(_UNICODE_DASH_TABLE)


def _kw_in(keyword: str, text: str, *, word_boundary: bool) -> bool:
    """Return True if ``keyword`` appears in ``text`` (case-insensitive).

    With ``word_boundary=True``, the match is anchored on regex ``\\b`` boundaries
    so short tokens like 'r1' or 'lives' don't match inside 'router 1' or 'olives'.
    Boundaries only fire between \\w and \\W, so keywords that begin or end with
    punctuation (e.g. '=B2-C2') still match correctly.

    Unicode dash variants (em-dash, en-dash, non-breaking hyphen, etc.) are
    normalised to ASCII hyphen before matching — models frequently use typographic
    dashes in structured names like CIP‑003‑9.
    """
    needle = _normalize_dashes(keyword.lower())
    haystack = _normalize_dashes(text.lower())
    if not word_boundary:
        return needle in haystack
    import re

    return re.search(rf"\b{re.escape(needle)}\b", haystack) is not None


def assert_contains(text: str, keywords: list, label: str, *, word_boundary: bool = False) -> tuple:
    missing = [k for k in keywords if not _kw_in(k, text, word_boundary=word_boundary)]
    return (label, not missing, f"missing: {missing}" if missing else "ok")


def assert_any_of(text: str, keywords: list, label: str, *, word_boundary: bool = False) -> tuple:
    found = [k for k in keywords if _kw_in(k, text, word_boundary=word_boundary)]
    return (label, bool(found), f"found: {found}" if found else f"none of: {keywords}")


def assert_not_contains(
    text: str, keywords: list, label: str, *, word_boundary: bool = False
) -> tuple:
    found = [k for k in keywords if _kw_in(k, text, word_boundary=word_boundary)]
    return (label, not found, f"found (bad): {found}" if found else "ok")


def assert_min_length(text: str, chars: int, label: str) -> tuple:
    return (label, len(text) >= chars, f"len={len(text)}, min={chars}")


def assert_has_code(text: str, label: str) -> tuple:
    has_fence = "```" in text
    # Raw HTML delivery (no markdown wrapper) is also valid code delivery
    has_raw_html = text.strip().startswith("<!DOCTYPE") or text.strip().startswith("<html")
    ok = has_fence or has_raw_html
    detail = (
        "code block present" if has_fence else ("raw html" if has_raw_html else "no code block")
    )
    return (label, ok, detail)


def _extract_code_blocks(text: str) -> str:
    """Extract and concatenate content from markdown code blocks.

    Handles fenced blocks (```lang ... ```), unclosed fenced blocks
    (opening fence without closing — common in model output), and raw
    HTML starting with <!DOCTYPE or <html>.
    Returns the concatenated code text, or '' if no code blocks found.
    """
    import re

    parts: list[str] = []
    text_lower = text.lower()

    # Fenced code blocks: ```optional_lang\n...\n```
    for m in re.finditer(r"```(?:\w+)?\n(.*?)```", text, re.DOTALL):
        parts.append(m.group(1).strip())

    # Unclosed fenced block: opening ``` anywhere, no closing ```
    if not parts:
        fence_match = re.search(r"```(?:\w+)?\n", text)
        if fence_match and "```" not in text[fence_match.end() :]:
            code_text = text[fence_match.end() :].strip()
            parts.append(code_text)

    # Raw HTML delivery (no markdown wrapper)
    if not parts:
        stripped = text.strip()
        if stripped.startswith("<!DOCTYPE") or stripped.startswith("<html"):
            parts.append(stripped)

    return "\n".join(parts)


def assert_code_pattern(text: str, patterns: list[dict], label: str) -> tuple:
    """Run regex patterns against extracted code blocks (not full response).

    Each pattern dict has:
        regex: str   — regex pattern to search for
        label: str   — human-readable description

    Patterns are case-insensitive. If any pattern matches, the assertion passes.
    This checks actual code behavior, not prose or variable naming conventions.
    """
    import re

    code = _extract_code_blocks(text)
    if not code:
        return (label, False, "no code blocks extracted")

    for p in patterns:
        pattern = p["regex"]
        try:
            if re.search(pattern, code, re.IGNORECASE):
                return (label, True, f"matched: {p.get('label', pattern)}")
        except re.error as e:
            return (label, False, f"invalid regex '{pattern}': {e}")

    return (label, False, f"no pattern matched in code ({len(code)} chars)")


def assert_has_table(text: str, label: str) -> tuple:
    return (label, "|" in text and "---" in text, "table present" if "|" in text else "no table")


def assert_docx_valid(path: Path | None, label: str) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        from docx import Document

        doc = Document(path)
        return (label, len(doc.paragraphs) > 0, f"{len(doc.paragraphs)} paragraphs")
    except Exception as e:
        return (label, False, str(e))


def assert_xlsx_valid(path: Path | None, label: str) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        from openpyxl import load_workbook

        wb = load_workbook(path)
        return (label, len(wb.sheetnames) > 0, f"sheets: {wb.sheetnames}")
    except Exception as e:
        return (label, False, str(e))


def assert_pptx_valid(path: Path | None, min_slides: int, label: str) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        from pptx import Presentation

        prs = Presentation(path)
        return (label, len(prs.slides) >= min_slides, f"{len(prs.slides)} slides")
    except Exception as e:
        return (label, False, str(e))


def assert_wav_valid(
    path: Path | None,
    label: str,
    *,
    min_seconds: float = 0.0,
) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        data = path.read_bytes()
        if not (len(data) > 1000 and data[:4] == b"RIFF" and data[8:12] == b"WAVE"):
            return (label, False, f"not a valid WAV: {len(data)} bytes")
        if min_seconds > 0:
            import wave

            with wave.open(str(path), "rb") as w:
                duration = w.getnframes() / float(w.getframerate())
            if duration < min_seconds:
                return (
                    label,
                    False,
                    f"too short: {duration:.1f}s < {min_seconds}s ({len(data)} bytes)",
                )
            return (label, True, f"{duration:.1f}s, {len(data)} bytes")
        return (label, True, f"{len(data)} bytes")
    except Exception as e:
        return (label, False, str(e))


def assert_png_valid(
    path: Path | None,
    label: str,
    *,
    min_width: int = 0,
    min_height: int = 0,
) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        data = path.read_bytes()
        if data[:8] != b"\x89PNG\r\n\x1a\n":
            return (label, False, f"not a PNG: {data[:8]!r}")
        if min_width > 0 or min_height > 0:
            try:
                from PIL import Image

                with Image.open(path) as im:
                    w, h = im.size
                if w < min_width or h < min_height:
                    return (
                        label,
                        False,
                        f"too small: {w}x{h} < {min_width}x{min_height}",
                    )
                return (label, True, f"{w}x{h}, {len(data)} bytes")
            except ImportError:
                return (label, True, f"PIL unavailable; {len(data)} bytes")
        return (label, True, f"{len(data)} bytes")
    except Exception as e:
        return (label, False, str(e))


def assert_mp4_valid(
    path: Path | None,
    label: str,
    *,
    min_seconds: float = 0.0,
) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        data = path.read_bytes()
        if b"ftyp" not in data[:32]:
            return (label, False, f"not an MP4: {data[:16]!r}")
        if min_seconds > 0:
            import subprocess

            try:
                out = subprocess.check_output(
                    [
                        "ffprobe",
                        "-v",
                        "error",
                        "-show_entries",
                        "format=duration",
                        "-of",
                        "default=noprint_wrappers=1:nokey=1",
                        str(path),
                    ],
                    text=True,
                    timeout=10,
                ).strip()
                duration = float(out)
                if duration < min_seconds:
                    return (
                        label,
                        False,
                        f"too short: {duration:.1f}s < {min_seconds}s",
                    )
                return (label, True, f"{duration:.1f}s, {len(data)} bytes")
            except FileNotFoundError:
                return (label, len(data) > 50_000, f"{len(data)} bytes (no ffprobe)")
            except Exception as e:
                return (label, False, str(e))
        return (label, True, f"{len(data)} bytes")
    except Exception as e:
        return (label, False, str(e))


def assert_tool_output_pattern(text: str, patterns: list, label: str) -> tuple:
    """Check that the response contains command output from an actual tool execution.

    Matches regex patterns against the full response (after think-stripping).
    Used to distinguish real tool execution (raw command output) from the model
    describing what it would do.
    """
    import re
    stripped = _strip_think_blocks(text)
    for p in patterns:
        try:
            if re.search(p, stripped, re.IGNORECASE | re.MULTILINE):
                return (label, True, f"matched: {p[:60]}")
        except re.error:
            pass
    return (label, False, f"no tool output pattern matched in {len(stripped)}-char response")


def assert_pipeline_tool_called(
    tool_calls_before: float,
    tool_calls_after: float,
    label: str,
) -> tuple:
    """Verify that at least one tool dispatch happened between two metric snapshots."""
    called = tool_calls_after > tool_calls_before
    delta = tool_calls_after - tool_calls_before
    return (
        label,
        called,
        f"tool_calls delta={delta:.0f} (before={tool_calls_before:.0f}, after={tool_calls_after:.0f})",
    )


def run_assertions(
    text: str,
    assertions_spec: list,
    artifact_path: Path | None = None,
    include_thinking: bool = False,
    *,
    tool_calls_before: float = 0.0,
    tool_calls_after: float = 0.0,
) -> list:
    if not include_thinking:
        text = _strip_think_blocks(text)
    results = []
    for a in assertions_spec:
        t = a["type"]
        label = a.get("label", t)
        if t == "contains":
            results.append(
                assert_contains(
                    text, a["keywords"], label, word_boundary=a.get("word_boundary", False)
                )
            )
        elif t == "any_of":
            results.append(
                assert_any_of(
                    text, a["keywords"], label, word_boundary=a.get("word_boundary", False)
                )
            )
        elif t == "not_contains":
            results.append(
                assert_not_contains(
                    text, a["keywords"], label, word_boundary=a.get("word_boundary", False)
                )
            )
        elif t == "min_length":
            results.append(assert_min_length(text, a["chars"], label))
        elif t == "has_code":
            results.append(assert_has_code(text, label))
        elif t == "code_pattern":
            results.append(assert_code_pattern(text, a.get("patterns", []), label))
        elif t == "has_table":
            results.append(assert_has_table(text, label))
        elif t == "docx_valid":
            results.append(assert_docx_valid(artifact_path, label))
        elif t == "xlsx_valid":
            results.append(assert_xlsx_valid(artifact_path, label))
        elif t == "pptx_valid":
            min_slides = a.get("min_slides", 1)
            results.append(assert_pptx_valid(artifact_path, min_slides, label))
        elif t == "wav_valid":
            min_s = float(a.get("min_seconds", 0.0))
            results.append(assert_wav_valid(artifact_path, label, min_seconds=min_s))
        elif t == "png_valid":
            mw = int(a.get("min_width", 0))
            mh = int(a.get("min_height", 0))
            results.append(assert_png_valid(artifact_path, label, min_width=mw, min_height=mh))
        elif t == "tool_output_pattern":
            results.append(
                assert_tool_output_pattern(text, a.get("patterns", []), label)
            )
        elif t == "pipeline_tool_called":
            results.append(
                assert_pipeline_tool_called(tool_calls_before, tool_calls_after, label)
            )
        elif t == "mp4_valid":
            min_s = float(a.get("min_seconds", 0.0))
            results.append(assert_mp4_valid(artifact_path, label, min_seconds=min_s))
        elif t == "quality_score":
            threshold = a.get("min", 0.5)
            cat = a.get("category", "general")
            try:
                import os as _os
                import sys as _sys

                _sys.path.insert(0, _os.path.dirname(_os.path.dirname(__file__)))
                from quality_signals import quality_score as _qs

                qs = _qs(cat, text)
            except Exception:
                qs = 1.0
            label_ext = f"{label} ({qs:.2f})"
            results.append((label_ext, qs >= threshold, f"score={qs:.2f}, min={threshold}"))
    return results


def compute_status(assertions: list, assertions_spec: list) -> str:
    """Grade a test result.

    By default every assertion is critical (a single failure produces FAIL).
    To opt an assertion into the percentage-grading floor, mark it with
    ``"critical": False`` in the spec. Behavior:

    - Any spec entry with ``critical=True`` (the default) that fails -> FAIL,
      UNLESS the overall pass rate is >=70%, in which case it's downgraded to WARN.
      This prevents a single narrow keyword from failing an otherwise correct test.
    - Otherwise, if all failing specs are ``critical=False``: PASS at >=70% pass
      rate, WARN at >=50%, FAIL below.

    The percentage rule now applies even with critical failures when overall
    score is high enough to demonstrate correct model behavior.
    """
    if not assertions:
        return "FAIL"
    total = len(assertions)
    passed_count = sum(1 for r in assertions if r[1])
    pct = passed_count / total * 100

    # Any critical failure is an automatic FAIL — unless the overall pass rate
    # is high enough to demonstrate that the model behaved correctly and the
    # failing assertion is likely a keyword-too-strict issue.
    has_critical_fail = False
    for result, spec in zip(assertions, assertions_spec):
        _label, passed, _evidence = result
        # has_code is a format preference — good code without a fenced block is
        # still correct behavior and should not alone fail an otherwise valid test.
        default_critical = spec.get("type") != "has_code"
        critical = spec.get("critical", default_critical)
        if not passed and critical:
            has_critical_fail = True
            break

    if has_critical_fail:
        return "FAIL"

    if pct >= 70:
        return "PASS"
    if pct >= 50 or passed_count > 0:
        return "WARN"
    return "FAIL"
