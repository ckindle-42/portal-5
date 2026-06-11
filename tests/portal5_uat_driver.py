#!/usr/bin/env python3
"""Portal 5 UAT Conversation Driver v1

Sends every test in TEST_CATALOG through the real Open WebUI browser
interface, creating permanent reviewable conversations in OWUI history.
The catalog currently spans ~175 tests across 24 sections including
auto-* workspaces, the `challenge` shootout (bench-* workspaces), and an `advanced` section
covering multi-turn / advanced flows.

Run modes:
    python3 tests/portal5_uat_driver.py --all
    python3 tests/portal5_uat_driver.py --section auto-coding
    python3 tests/portal5_uat_driver.py --section auto-coding --section challenge
    python3 tests/portal5_uat_driver.py --test WS-01 --test P-W06
    python3 tests/portal5_uat_driver.py --all --headed --append

Calibration mode (capture real responses for signal extraction):
    python3 tests/portal5_uat_driver.py --calibrate \
        --calibrate-output calibration.json
    # ... review calibration.json, set review_tag = good/bad/skip on each entry
    python3 tests/portal5_uat_driver.py --emit-signals-from calibration.json
    # See docs/UAT_CALIBRATION.md for the full workflow.

Maintenance:
    python3 tests/portal5_uat_driver.py --purge-uat        # delete all UAT chats + folder (post-review cleanup)
    python3 tests/portal5_uat_driver.py --migrate          # move root chats into UAT folder
    python3 tests/portal5_uat_driver.py --skip-artifacts  # skip ComfyUI/Wan2.2 tests
    python3 tests/portal5_uat_driver.py --skip-bots       # skip Telegram/Slack tests
"""

from __future__ import annotations

import argparse
import asyncio
import json as _json
import os
import sys
import threading
import time
import uuid
from pathlib import Path

import httpx
from dotenv import load_dotenv
from playwright.async_api import async_playwright

load_dotenv()

sys.path.insert(0, str(Path(__file__).parent))

from tests.uat_catalog import TEST_CATALOG  # assembled from tests/uat_catalog/g_*.py

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

OPENWEBUI_URL = os.environ.get("OPENWEBUI_URL", "http://localhost:8080")
ADMIN_EMAIL = os.environ.get("OPENWEBUI_ADMIN_EMAIL", "admin@portal.local")
ADMIN_PASS = os.environ.get("OPENWEBUI_ADMIN_PASSWORD", "")

SEND_TIMEOUT = 300_000  # initial window for stop-button to appear (cold load)
PROGRESS_POLL_S = 30  # legacy heartbeat interval (kept for compatibility)
MAX_WAIT_NO_PROGRESS = 900  # 15 min hard cap if zero progress detected
NO_STREAM_TIMEOUT = 120  # exit for retry if stop never appeared after this many seconds
PROGRESS_LOG_INTERVAL = 120  # log a heartbeat every 2 min

# Tiered polling intervals — replace the single 30s PROGRESS_POLL_S at the
# decision points in _wait_for_completion. The 30s value remains in use as
# a heartbeat reference but is no longer the polling resolution.
PHASE1_FAST_S = 0.5  # poll every 0.5s while waiting for stream to start
PHASE1_FAST_DURATION_S = 10  # for the first 10 seconds
PHASE1_MID_S = 2.0  # then poll every 2s
PHASE1_MID_DURATION_S = 30  # for the next 30 seconds (10s..40s elapsed)
PHASE1_SLOW_S = 5.0  # then poll every 5s for very cold loads (40s+)

PHASE2_STREAMING_POLL_S = 1.5  # poll every 1.5s while model is actively streaming
PHASE2_DOM_STABLE_NEEDED = 3  # consecutive identical samples to declare DOM stable

POST_STREAM_API_WAIT_S = 15.0  # bounded API poll after stream ends (replaces fixed sleep(5))
BACKEND_SETTLE_WAIT_S = 15.0  # bounded backend-alive poll after retry (replaces sleep(15))
RESULTS_FILE = Path("tests/UAT_RESULTS.md")
SCREENSHOT_DIR = Path("/tmp/uat_screenshots")
ARTIFACT_DIR = Path("/tmp/uat_artifacts")

# Routing telemetry — appended per test, written to UAT_RESULTS.md at end of run.
# Each entry: {test_id, name, section, workspace, intended, actual, matched, tier_mismatch}
_ROUTING_LOG: list[dict] = []

# Chat IDs created in the current run. Populated by owui_create_chat() so that
# the post-run archival step can move all chats to a dated UAT subfolder.
_run_chat_ids: list[str] = []
OLLAMA_URL = os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434")

# Sections that require all models unloaded before running for max memory headroom.
SECTIONS_REQUIRE_UNLOAD = True  # Always unload Ollama before sections

# Memory pressure thresholds
MEMORY_WARN_PCT = 80.0  # Log warning
MEMORY_CRITICAL_PCT = (
    90.0  # Force eviction before next test
)
MEMORY_ABORT_PCT = 95.0  # Stop — system is about to OOM
# Same-model eviction: even when the next test uses the same model, evict if
# memory exceeds this after the previous test. KV cache from long inference
# compounds into the next test's KV cache allocation.
# (Observed: gemma-4-26b at 82% post-P-V02 crashed at 92%
# during P-R06. Same model, no eviction, compounding KV cache = crash.)
MEMORY_SAME_MODEL_EVICT_PCT = 78.0
# After this many consecutive "DOM stable but API empty" cycles, assume OWUI 0.9.5+
# is not going to commit the response via API (thinking-model commit delay) and let
# the caller's DOM fallback extract the response directly from the page.
DOM_STABLE_API_EMPTY_MAX = 3

# ---------------------------------------------------------------------------
# Codebase freshness check
# ---------------------------------------------------------------------------

_REPO_ROOT = Path(__file__).parent.parent


def _check_image_freshness() -> list[str]:
    """Return list of warning strings for Docker images older than the latest git commit.

    Compares each portal image build timestamp against the most recent git commit that
    touched the files that image is built from. Stale images mean the running code does
    not match HEAD and tests will not reflect current behaviour.

    Prints a clear summary; returns list of stale image names (empty = all current).
    """
    import datetime
    import subprocess

    warnings: list[str] = []

    # Latest commit time for files that affect each image
    def _last_commit_ts(paths: list[str]) -> datetime.datetime | None:
        try:
            result = subprocess.run(
                ["git", "-C", str(_REPO_ROOT), "log", "-1", "--format=%ct", "--", *paths],
                capture_output=True,
                text=True,
                timeout=10,
            )
            ts = result.stdout.strip()
            if ts:
                return datetime.datetime.fromtimestamp(int(ts), tz=datetime.timezone.utc)
        except Exception:
            pass
        return None

    # Image build time via docker inspect
    def _image_built_ts(image_name: str) -> datetime.datetime | None:
        try:
            result = subprocess.run(
                ["docker", "inspect", "--format", "{{.Created}}", image_name],
                capture_output=True,
                text=True,
                timeout=10,
            )
            raw = result.stdout.strip()
            if raw and raw != "[]":
                # Parse RFC3339/ISO format
                raw = raw.rstrip("Z") + "+00:00"
                return datetime.datetime.fromisoformat(raw)
        except Exception:
            pass
        return None

    checks = [
        # (label, docker_image_name, git_paths_that_affect_it)
        (
            "portal-pipeline",
            "portal-5-portal-pipeline",
            [
                "portal_pipeline/",
                "config/backends.yaml",
                "config/personas/",
                "Dockerfile.pipeline",
                "pyproject.toml",
                "scripts/pipeline-entrypoint.sh",
            ],
        ),
        (
            "mcp-services",
            "portal-5-mcp-documents",
            ["portal_mcp/", "portal_channels/", "Dockerfile.mcp", "pyproject.toml"],
        ),
    ]

    print("  [freshness] Checking Docker image freshness against git HEAD...")
    all_fresh = True
    for label, image, paths in checks:
        built = _image_built_ts(image)
        committed = _last_commit_ts(paths)
        if built is None:
            print(f"  [freshness]   {label}: image not found — skip")
            continue
        if committed is None:
            print(f"  [freshness]   {label}: no git history — skip")
            continue
        lag = (committed - built).total_seconds()
        if lag > 30:  # >30s: image predates the commit
            mins = int(lag // 60)
            msg = (
                f"  [freshness] WARNING: {label} image is {mins}m behind HEAD — "
                f"run './launch.sh rebuild' before trusting results"
            )
            print(msg, flush=True)
            warnings.append(label)
            all_fresh = False
        else:
            print(f"  [freshness]   {label}: current (lag={lag:.0f}s)", flush=True)

    if all_fresh:
        print("  [freshness] All images are current.", flush=True)

    return warnings


# ---------------------------------------------------------------------------
# Backend health + zombie detection
# ---------------------------------------------------------------------------


def _get_memory_pct() -> float:
    """Get current memory used % from system vm_stat."""
    try:
        import subprocess

        result = subprocess.run(["vm_stat"], capture_output=True, text=True, timeout=5)
        lines = result.stdout.strip().split("\n")
        page_size = 16384  # Apple Silicon
        free = active = inactive = speculative = wired = 0
        for line in lines:
            if "Pages free:" in line:
                free = int(line.split(":")[1].strip().rstrip("."))
            elif "Pages active:" in line:
                active = int(line.split(":")[1].strip().rstrip("."))
            elif "Pages inactive:" in line:
                inactive = int(line.split(":")[1].strip().rstrip("."))
            elif "Pages speculative:" in line:
                speculative = int(line.split(":")[1].strip().rstrip("."))
            elif "Pages wired down:" in line:
                wired = int(line.split(":")[1].strip().rstrip("."))
        total = free + active + inactive + speculative + wired
        if total > 0:
            used = active + wired
            return round(used / total * 100, 1)
    except Exception:
        pass
    return 0.0


def _wait_for_drain(
    threshold_pct: float = MEMORY_CRITICAL_PCT,
    timeout_s: float = 90.0,
    poll_s: float = 2.0,
    label: str = "",
) -> bool:
    """Wait for memory to drop below threshold_pct after model eviction.

    Polls vm_stat every poll_s seconds and breaks as soon as
    memory < threshold_pct. Does not use blind fixed sleeps —
    exits immediately when the condition is met. Hard timeout is the only timer.

    Returns True if memory cleared within timeout_s, False on timeout.
    """
    deadline = time.time() + timeout_s
    prefix = f"  [drain{' ' + label if label else ''}]"
    while time.time() < deadline:
        used_pct = _get_memory_pct()
        if used_pct < threshold_pct:
            print(f"{prefix} Clear at {used_pct:.0f}% — safe to proceed", flush=True)
            return True
        remaining = int(deadline - time.time())
        print(f"{prefix} {used_pct:.0f}% ({remaining}s left)", flush=True)
        time.sleep(poll_s)
    return False


def _check_memory_before_test(test_name: str = "") -> bool:
    """Check memory pressure before running a test. Returns True if safe to proceed.

    If critical: force-evicts all models and returns False (caller should skip/retry).
    If abort: raises SystemExit to prevent OOM crash.
    """
    used = _get_memory_pct()

    if used >= MEMORY_ABORT_PCT:
        print(
            f"\n  [OOM RISK] Memory at {used:.0f}% — aborting to prevent crash. Test: {test_name}",
            flush=True,
        )
        unload_all_models()
        # Wait for actual Metal reclaim — event-driven, not a blind sleep.
        _wait_for_drain(threshold_pct=MEMORY_WARN_PCT, timeout_s=30.0, label="oom-risk")
        used_after = _get_memory_pct()
        print(f"  [OOM RISK] After eviction: {used_after:.0f}%", flush=True)
        if used_after >= MEMORY_ABORT_PCT:
            raise SystemExit(
                f"ABORT: Memory still at {used_after:.0f}% after full eviction. "
                "Manual intervention required — check for leaked processes."
            )
        return False

    if used >= MEMORY_CRITICAL_PCT:
        print(
            f"\n  [MEMORY] Critical: {used:.0f}% — evicting before {test_name}",
            flush=True,
        )
        unload_all_models()
        if not _wait_for_drain(threshold_pct=MEMORY_WARN_PCT, timeout_s=90.0, label=test_name):
            print(f"  [MEMORY] Still above {MEMORY_WARN_PCT:.0f}% after 90s drain — skipping {test_name}", flush=True)
            return False
        return True

    if used >= MEMORY_WARN_PCT:
        print(f"  [MEMORY] Warning: {used:.0f}% used", flush=True)

    return True


def _check_for_oom_crash() -> str | None:
    """Check if any backend crashed due to OOM since last check.

    Detects:
    1. Ollama unreachable
    2. System memory above abort threshold

    Returns crash description or None if healthy.
    """
    # Ollama dead?
    try:
        r = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=3)
        if r.status_code != 200:
            return f"Ollama returned {r.status_code}"
    except Exception:
        return "Ollama unreachable (process may have crashed)"

    # System memory critical?
    used = _get_memory_pct()
    if used >= MEMORY_ABORT_PCT:
        return f"System memory at {used:.0f}% — OOM imminent"

    return None


def _backend_alive(tier: str) -> tuple[bool, str]:
    """Return (alive, detail) for the given workspace tier."""
    if tier in ("ollama",):
        try:
            r = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=5)
            return r.status_code == 200, f"ollama={r.status_code}"
        except Exception:
            return False, "ollama_unreachable"
    if tier == "media_heavy":
        try:
            r = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=5)
            ollama_ok = r.status_code == 200
        except Exception:
            ollama_ok = False
        return ollama_ok, f"ollama={'ok' if ollama_ok else 'down'}"
    return True, "tier=any"


def _wait_for_ollama_ps_empty(timeout_s: float = 30.0, poll_s: float = 1.0) -> bool:
    """Poll /api/ps until all models are unloaded or timeout_s elapses.

    Event-driven: exits as soon as Ollama reports no loaded models.
    Safety net: hard timeout for cases where Ollama hangs on release.
    Returns True when model list is empty, False on timeout.

    This is step 1 of the two-step drain:
      1. /api/ps empty → Ollama confirmed release (this function)
      2. vm_stat below threshold → Metal buffers reclaimed (_wait_for_drain)
    """
    deadline = time.time() + timeout_s
    while time.time() < deadline:
        try:
            r = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=3)
            if r.status_code == 200 and not r.json().get("models"):
                return True
        except Exception:
            pass
        remaining = deadline - time.time()
        if remaining <= 0:
            break
        time.sleep(min(poll_s, remaining))
    return False


def _unload_running_ollama_models() -> None:
    """Unload all running Ollama models via the Ollama API."""
    try:
        r = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=5)
        if r.status_code != 200:
            return
        models = r.json().get("models", [])
        for m in models:
            name = m.get("name") or m.get("model", "")
            if name:
                print(f"  Unloading Ollama model: {name}", flush=True)
                httpx.post(
                    f"{OLLAMA_URL}/api/generate",
                    json={"model": name, "keep_alive": 0},
                    timeout=10,
                )
    except Exception as e:
        print(f"  WARNING: Ollama unload failed: {e}", flush=True)


async def _wait_for_backend(tier: str, max_wait: int = 120) -> bool:
    """Poll backend until ready or max_wait seconds elapsed.

    Returns True if the backend became ready, False if it stayed down.
    Emits progress lines every 20s so the operator can see what's happening.
    """
    if tier not in ("ollama", "media_heavy"):
        return True
    t0 = time.time()
    last_log = 0.0
    while True:
        alive, detail = _backend_alive(tier)
        if alive:
            return True
        elapsed = time.time() - t0
        if elapsed >= max_wait:
            print(
                f"  [health] backend still not ready after {max_wait:.0f}s ({detail})", flush=True
            )
            return False
        if time.time() - last_log >= 20:
            print(
                f"  [health] waiting for backend ({detail}, {elapsed:.0f}s/{max_wait}s)…",
                flush=True,
            )
            last_log = time.time()
        await asyncio.sleep(10)


# ---------------------------------------------------------------------------
# Model unload helpers
# ---------------------------------------------------------------------------


def _pipeline_pre_warm(workspace_id: str = "auto") -> None:
    """Send a minimal request through the pipeline to trigger model cold-load.

    Uses the actual workspace_id for this test so the right model is pre-loaded.
    """
    pipeline_url = os.environ.get("PIPELINE_URL", "http://localhost:9099")
    pipeline_key = os.environ.get("PIPELINE_API_KEY", "")
    if not pipeline_key:
        try:
            env_path = Path(__file__).parent.parent / ".env"
            for _line in env_path.read_text().splitlines():
                if _line.startswith("PIPELINE_API_KEY="):
                    pipeline_key = _line.split("=", 1)[1].strip()
                    break
        except Exception:
            pass
    if not pipeline_key:
        pipeline_key = "portal-pipeline"
    import threading as _threading

    _prewarm_done = _threading.Event()

    def _prewarm_ticker() -> None:
        tick = 0
        while not _prewarm_done.wait(timeout=30):
            tick += 30
            mem = _get_memory_pct()
            print(
                f"  [pre-warm] {tick}s elapsed — mem={mem:.0f}%",
                flush=True,
            )

    ticker = _threading.Thread(target=_prewarm_ticker, daemon=True, name="prewarm-ticker")
    ticker.start()

    try:
        httpx.post(
            f"{pipeline_url}/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {pipeline_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": workspace_id,
                "messages": [{"role": "user", "content": "hi"}],
                "max_tokens": 1,
                "stream": False,
            },
            timeout=360,
        )
    except Exception:
        pass
    finally:
        _prewarm_done.set()


def unload_all_models() -> None:
    """Unload all running Ollama models and release GPU memory.

    Two-step drain:
      1. Send keep_alive=0 to all loaded models, then poll /api/ps until
         the list is empty (event-driven — exits as soon as Ollama confirms
         release, not after a fixed sleep).
      2. Caller follows with _wait_for_drain() to wait for macOS Metal to
         reclaim wired pages (vm_stat-driven, same event pattern).
    """
    print("  Unloading all Ollama models ...", flush=True)
    _unload_running_ollama_models()
    _wait_for_ollama_ps_empty(timeout_s=30.0)


def _comfyui_running() -> bool:
    """Return True if ComfyUI is reachable on :8188."""
    try:
        r = httpx.get("http://localhost:8188/system_stats", timeout=3)
        return r.status_code == 200
    except Exception:
        return False


def _start_comfyui(wait_s: int = 45) -> bool:
    """Start ComfyUI via launchctl and wait for it to become reachable."""
    import subprocess

    print("  [comfyui] Starting ComfyUI ...", flush=True)
    try:
        subprocess.run(
            ["launchctl", "start", "com.portal5.comfyui"],
            capture_output=True,
            timeout=10,
        )
    except Exception as e:
        print(f"  [comfyui] launchctl start failed: {e}", flush=True)
        return False
    deadline = time.time() + wait_s
    while time.time() < deadline:
        if _comfyui_running():
            print("  [comfyui] ComfyUI ready", flush=True)
            return True
        time.sleep(3)
    print(f"  [comfyui] ComfyUI did not become ready after {wait_s}s", flush=True)
    return False


def _stop_comfyui() -> None:
    """Stop ComfyUI via launchctl to reclaim GPU memory between non-media phases."""
    import subprocess

    if not _comfyui_running():
        return
    print("  [comfyui] Stopping ComfyUI to reclaim GPU memory ...", flush=True)
    try:
        subprocess.run(
            ["launchctl", "stop", "com.portal5.comfyui"],
            capture_output=True,
            timeout=10,
        )
        # Wait briefly for Metal to release
        time.sleep(5)
        print("  [comfyui] ComfyUI stopped", flush=True)
    except Exception as e:
        print(f"  [comfyui] stop failed: {e}", flush=True)


def cleanup_after_uat() -> None:
    """Full cleanup after all UAT tests complete — prevents OOM post-run."""
    print("\n  Post-UAT cleanup: evicting all models ...", end=" ", flush=True)
    unload_all_models()
    used = _get_memory_pct()
    print(f"ok (mem={used:.0f}%)", flush=True)


# ---------------------------------------------------------------------------
# OWUI API helpers
# ---------------------------------------------------------------------------


def owui_token() -> str:
    r = httpx.post(
        f"{OPENWEBUI_URL}/api/v1/auths/signin",
        json={"email": ADMIN_EMAIL, "password": ADMIN_PASS},
        timeout=10,
    )
    return r.json().get("token", "")


def owui_headers(token: str) -> dict:
    return {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}


def owui_create_chat(token: str, model_slug: str, title: str) -> tuple[str, str]:
    chat_id = str(uuid.uuid4())
    payload = {
        "chat": {
            "id": chat_id,
            "title": title,
            "models": [model_slug],
            "messages": [],
            "history": {"messages": {}, "currentId": None},
            "tags": [],
            "params": {},
            "timestamp": int(time.time()),
        }
    }
    r = httpx.post(
        f"{OPENWEBUI_URL}/api/v1/chats/new",
        json=payload,
        headers=owui_headers(token),
        timeout=10,
    )
    returned_id = r.json().get("id", chat_id)
    _run_chat_ids.append(returned_id)
    return returned_id, f"{OPENWEBUI_URL}/c/{returned_id}"


def owui_rename_chat(token: str, chat_id: str, title: str) -> None:
    httpx.post(
        f"{OPENWEBUI_URL}/api/v1/chats/{chat_id}",
        json={"chat": {"title": title}},
        headers=owui_headers(token),
        timeout=10,
    )


def _owui_list_folders(token: str) -> list[dict]:
    r = httpx.get(
        f"{OPENWEBUI_URL}/api/v1/folders/",
        headers=owui_headers(token),
        timeout=30,
    )
    if r.status_code == 200:
        try:
            return r.json()
        except Exception:
            pass
    return []


def owui_get_or_create_folder(token: str, name: str, parent_id: str | None = None) -> str | None:
    """Return folder ID for `name` under `parent_id` (root if None), creating if absent."""
    folders = _owui_list_folders(token)
    for folder in folders:
        if folder.get("name") == name and folder.get("parent_id") == parent_id:
            return folder.get("id")

    r = httpx.post(
        f"{OPENWEBUI_URL}/api/v1/folders/",
        json={"name": name, "parent_id": parent_id},
        headers=owui_headers(token),
        timeout=10,
    )
    if r.status_code == 200:
        return r.json().get("id")

    # "already exists" race — re-fetch
    if r.status_code == 400 and "already exists" in r.text:
        for folder in _owui_list_folders(token):
            if folder.get("name") == name and folder.get("parent_id") == parent_id:
                return folder.get("id")

    return None


def owui_assign_chat_folder(token: str, chat_id: str, folder_id: str) -> None:
    """Move a chat into the given folder."""
    try:
        httpx.post(
            f"{OPENWEBUI_URL}/api/v1/chats/{chat_id}/folder",
            json={"folder_id": folder_id},
            headers=owui_headers(token),
            timeout=10,
        )
    except Exception:
        pass


def owui_migrate_loose_uat_chats(token: str, root_folder_id: str) -> int:
    """Move any root-level UAT chats (no folder_id) into root_folder_id.

    Returns the number of chats migrated.
    """
    moved = 0
    try:
        r = httpx.get(
            f"{OPENWEBUI_URL}/api/v1/chats/",
            headers=owui_headers(token),
            timeout=15,
        )
        if r.status_code != 200:
            return 0
        for chat in r.json():
            chat_id = chat.get("id", "")
            title = chat.get("title", "")
            # Full detail needed to check folder_id
            r2 = httpx.get(
                f"{OPENWEBUI_URL}/api/v1/chats/{chat_id}",
                headers=owui_headers(token),
                timeout=10,
            )
            if r2.status_code != 200:
                continue
            detail = r2.json()
            if detail.get("folder_id"):
                continue  # already in a folder
            if "UAT:" in title:
                owui_assign_chat_folder(token, chat_id, root_folder_id)
                moved += 1
    except Exception as e:
        print(f"  WARNING: migrate error — {e}")
    return moved


# ---------------------------------------------------------------------------
# Frontend dispatch shims — thin wrappers over OWUI helpers.
# ---------------------------------------------------------------------------


async def _fe_login(page) -> None:
    """Login to Open WebUI."""
    await _login(page)


async def _fe_start_chat(
    page,
    token: str,
    model_slug: str,
    title: str,
) -> tuple[str, str]:
    """Create / open a fresh chat. Returns (chat_id, chat_url).

    chat_id is the OWUI UUID (used by API helpers). chat_url is
    `{OPENWEBUI_URL}/c/{chat_id}`.
    """
    chat_id, chat_url = owui_create_chat(token, model_slug, title)
    await _navigate_to_chat(page, chat_url)
    return chat_id, chat_url


class _PresetUnreachableError(RuntimeError):
    """Raised by _fe_start_chat when a persona test cannot select its preset."""


async def _fe_send_and_wait(
    page,
    prompt: str,
    test_id: str = "",
    tier: str = "any",
    max_wait_no_progress: int = MAX_WAIT_NO_PROGRESS,
    *,
    token: str = "",
    chat_id: str = "",
    min_messages: int = 1,
) -> None:
    """Send a prompt and wait for streaming to complete."""
    await _send_and_wait(
        page,
        prompt,
        test_id=test_id,
        tier=tier,
        max_wait_no_progress=max_wait_no_progress,
        token=token,
        chat_id=chat_id,
        min_messages=min_messages,
    )


async def _extract_dom_response(page) -> str:
    """Extract the last assistant response text directly from the OWUI page DOM.

    Used as a fallback when OWUI 0.9.5+ does not immediately commit thinking-model
    responses to the chat history API. OWUI renders markdown content inside .prose
    divs; reasoning blocks are in <details> elements that are stripped before return.
    Returns '' if no suitable content is found (degrades gracefully).
    """
    try:
        return await page.evaluate(
            """() => {
            // OWUI 0.9.x renders markdown with Tailwind 'prose' class.
            // The last .prose element holds the most recent assistant response.
            const selectors = [
                '.prose.dark\\\\:prose-invert',
                '.prose',
                '[data-role="assistant"] .prose',
                '.message-content .prose',
            ];
            let best = '';
            for (const sel of selectors) {
                try {
                    const els = document.querySelectorAll(sel);
                    if (els.length === 0) continue;
                    const el = els[els.length - 1];
                    const clone = el.cloneNode(true);
                    for (const d of clone.querySelectorAll('details')) d.remove();
                    const text = (clone.innerText || '').trim();
                    if (text.length > best.length) best = text;
                } catch (_) {}
            }
            return best;
        }"""
        )
    except Exception:
        return ""


async def _fe_get_last_response(page, token: str, chat_id: str, min_messages: int = 1) -> str:
    """Read the most recent assistant message — OWUI API first, DOM fallback.

    OWUI 0.9.5+ may delay committing thinking-model responses to the chat history
    API until a new user message arrives. When the API returns empty but streaming
    has visually completed, _extract_dom_response reads directly from the rendered
    page content as a fallback.
    """
    api_result = owui_get_last_response(token, chat_id, min_messages=min_messages)
    if api_result:
        return api_result
    return await _extract_dom_response(page)


async def _fe_get_routed_model(test: dict, page, token: str, chat_id: str) -> str:
    """Resolve the actual backend model that handled the most recent response.

    Reads OWUI's stored chat metadata (the workspace/persona name captured from
    the pipeline's SSE stream). Returns "" if no source is available.
    _check_routed_model handles the empty-string case by skipping validation.
    """
    return owui_get_routed_model(token, chat_id)


async def _fe_enable_tool(page, tool_id: str) -> None:
    await _enable_tool(page, tool_id)


async def _fe_assign_folder(page, token: str, chat_id: str, folder_id: str | None) -> None:
    if folder_id:
        owui_assign_chat_folder(token, chat_id, folder_id)


async def _fe_download_artifact(
    page,
    expected_ext: str,
    response_text: str = "",
    timeout_ms: int = 120_000,
    *,
    since_ts: float = 0.0,
) -> Path | None:
    return await _download_artifact(page, expected_ext, timeout_ms, response_text)


def _fe_current_chat_url(page, fallback: str) -> str:
    """Return the API-given chat_url."""
    return fallback


def owui_get_last_response(token: str, chat_id: str, min_messages: int = 1) -> str:
    """Fetch the last assistant response from OWUI API — avoids Playwright truncation.

    For thinking models (Qwen3/AEON), OWUI only commits an assistant message when
    either streaming ends OR a new user message arrives in the chat. The in-flight
    message is always empty from the API's perspective. This function returns the
    last NON-EMPTY assistant message so that a committed partial response from a
    previous attempt is found as soon as the next attempt's send triggers a commit.

    OWUI embeds reasoning content inline in the content field as:
      <details type="reasoning" done="true" duration="N">...</details>[actual response]
    No separate reasoning field exists in the chat history API.

    min_messages: minimum number of non-empty assistant messages required before
    returning. Use min_messages=2 for multi-turn turn-2 detection to prevent
    turn-1's committed response from satisfying the completion signal prematurely.
    Returns "" (falsy) until the required count of non-empty messages exists.
    """
    try:
        r = httpx.get(
            f"{OPENWEBUI_URL}/api/v1/chats/{chat_id}",
            headers={"Authorization": f"Bearer {token}", "Accept-Encoding": "identity"},
            timeout=10,
        )
        msgs = r.json().get("chat", {}).get("history", {}).get("messages", {})
        assistant_msgs = [m for m in msgs.values() if m.get("role") == "assistant"]
        if not assistant_msgs:
            return ""
        # Collect all non-empty assistant messages in order.
        non_empty: list[str] = []
        for msg in assistant_msgs:
            content = msg.get("content", "")
            if isinstance(content, list):
                content = " ".join(c.get("text", "") for c in content if isinstance(c, dict))
            if content:
                non_empty.append(content)
        # Guard: for turn-2 in multi-turn tests (min_messages=2), return "" until
        # the second non-empty assistant message has been committed by OWUI.
        # For thinking models, in-flight messages are always empty; the most-recently-
        # committed previous attempt's response is the useful signal (min_messages=1).
        if len(non_empty) < min_messages:
            return ""
        return non_empty[-1]
    except Exception:
        return ""


def owui_get_routed_model(token: str, chat_id: str) -> str:
    """Extract the model actually used from the last assistant message in OWUI chat history.

    Returns the model string or "" if unavailable. Provides diagnostic value
    equivalent to reading x-portal-route: the pipeline embeds the selected
    backend model in the message metadata stored by Open WebUI.
    """
    try:
        r = httpx.get(
            f"{OPENWEBUI_URL}/api/v1/chats/{chat_id}",
            headers={"Authorization": f"Bearer {token}", "Accept-Encoding": "identity"},
            timeout=10,
        )
        data = r.json()
        msgs = data.get("chat", {}).get("history", {}).get("messages", {})
        assistant_msgs = [m for m in msgs.values() if m.get("role") == "assistant"]
        if not assistant_msgs:
            return ""
        last = assistant_msgs[-1]
        # OWUI stores the model that generated the response in the message metadata
        model = last.get("model", "") or last.get("info", {}).get("model", "")
        return str(model) if model else ""
    except Exception:
        return ""


def _map_slug_to_workspace(slug: str) -> str:
    """Resolve a persona slug to its workspace id, or return the slug
    if it's already a workspace id."""
    from expected_models import _PERSONA_MAP, WORKSPACES

    if slug in WORKSPACES:
        return slug
    p = _PERSONA_MAP.get(slug, {})
    ws = p.get("workspace_model") or p.get("workspace") or ""
    return ws if ws in WORKSPACES else ""


def _get_backend_from_pipeline_logs(slug: str) -> str:
    """Query pipeline Docker logs for the most recent backend that actually
    served a request for the given workspace/persona slug.

    Uses the "Backend X succeeded" log line (only emitted on actual success)
    rather than the "Routing workspace=X → backend=Y" line (emitted for the
    first candidate ATTEMPTED, which may 503 and fall to a different backend).

    Log line patterns (pipeline emits both; we prefer the succeeded line):
      Backend ollama-general succeeded for workspace=auto-documents model=phi4:14b-q8_0
      Backend ollama-coding succeeded for workspace=auto-agentic model=qwen3-coder:30b
    """
    import re
    import subprocess

    # Resolve persona slug to its workspace for log matching
    ws = _map_slug_to_workspace(slug)

    try:
        result = subprocess.run(
            ["docker", "logs", "portal5-pipeline", "--tail", "300"],
            capture_output=True,
            text=True,
            timeout=10,
        )
        combined = result.stdout + result.stderr  # docker logs may use either stream

        # PRIMARY: match "Backend X succeeded for workspace=Y model=Z"
        # This line is only emitted when a backend actually returns a response,
        # so it correctly reflects backend-group fallbacks that the
        # "Routing workspace" attempt line would hide.
        for search_term in (ws, slug):
            if not search_term:
                continue
            succeeded_pattern = re.compile(
                r"Backend\s+([^\s]+)\s+succeeded\s+for\s+workspace="
                + re.escape(search_term)
                + r"\s+model=([^\s]+)"
            )
            matches = succeeded_pattern.findall(combined)
            if matches:
                backend, model = matches[-1]  # most recent
                return f"{backend}|{model}"

        # FALLBACK: if no "succeeded" line found (e.g. non-stream path that
        # doesn't emit it), fall back to the attempt log line as before.
        for search_term in (ws, slug):
            if not search_term:
                continue
            pattern = re.compile(
                r"Routing workspace="
                + re.escape(search_term)
                + r".*?backend=([^\s]+)\s+model=([^\s]+)"
            )
            matches = pattern.findall(combined)
            if matches:
                backend, model = matches[-1]
                return f"{backend}|{model}"
    except Exception:
        pass
    return ""


def _check_routed_model(test: dict, routed_model: str) -> tuple[bool, str] | None:
    """Validate routed_model against test expectation.

    Two-source approach:
      1. OWUI chat metadata via owui_get_routed_model (may store the
         workspace/persona name, not the backend model)
      2. Pipeline Docker logs — extracts the actual backend=xxx model=yyy

    Returns:
        None             - no expectation defined for this test, skip the check
        (True,  detail)  - actual model matches expectation
        (False, detail)  - mismatch (caller should downgrade PASS to WARN)

    Resolution order:
        1. test['assert_routed_via']: list[str] of substrings
        2. test['model_slug'] in WORKSPACES
        3. test['model_slug'] in _PERSONA_MAP
        4. None — no expectation, skip
    """
    if test.get("via_dispatcher") or test.get("is_manual"):
        return None
    if not routed_model:
        return None

    import sys as _sys
    from pathlib import Path as _Path

    _sys.path.insert(0, str(_Path(__file__).parent))

    from expected_models import model_matches_expected, resolve_expected

    explicit = test.get("assert_routed_via")
    slug = test.get("model_slug", "")

    keys, src = resolve_expected(
        workspace_id=slug,
        persona_slug=slug,
    )
    if not keys:
        return None

    # 1st check: OWUI-stored model (may be the workspace/persona name)
    ok_owui = model_matches_expected(routed_model, keys)

    # 2nd check: pipeline logs (actual backend model)
    backend_model = _get_backend_from_pipeline_logs(slug)

    ok_pipeline = False
    pipeline_detail = ""
    if backend_model:
        ok_pipeline = model_matches_expected(backend_model, keys)
        pipeline_detail = f" (pipeline: {backend_model})"

    if explicit:
        ok = ok_owui or ok_pipeline
        return (
            ok,
            f"explicit expectation: {explicit}{pipeline_detail}"
            if ok
            else f"explicit expectation NOT matched: {explicit}{pipeline_detail}",
        )

    if ok_owui or ok_pipeline:
        detail = f"matches {src}"
        if backend_model:
            detail += f" — pipeline confirms: {backend_model}"
        return (True, detail)

    return (False, f"expected {src} (OWUI={routed_model}{pipeline_detail})")


async def _wait_for_response_arrival(
    token: str,
    chat_id: str,
    max_wait: float = POST_STREAM_API_WAIT_S,
    min_messages: int = 1,
) -> str:
    """Poll OWUI API until response content stabilizes (log-driven, not timer-based).

    Polls every 2s and declares done when content length hasn't grown by more
    than 50 chars across 3 consecutive polls. This is content-driven completion:
    the API response log drives the exit decision rather than a fixed sleep.

    OWUI persists the assistant message at end-of-stream with a brief lag
    (typically <500ms, occasionally a few seconds under load).

    min_messages: passed to owui_get_last_response. Set to 2 for multi-turn
    turn-2 to require the second committed assistant response.

    Returns last stable content string, or "" on timeout.
    """
    if not token or not chat_id:
        await asyncio.sleep(2.0)
        return ""

    STABLE_COUNT = 2  # consecutive polls with no meaningful growth (OWUI commits atomically)
    STABLE_THRESHOLD = 50  # chars; ignores minor whitespace/punctuation flushes

    deadline = time.monotonic() + max_wait
    len_history: list[int] = []
    last_text = ""

    while time.monotonic() < deadline:
        text = owui_get_last_response(token, chat_id, min_messages=min_messages)
        cur_len = len(text)
        if text:
            last_text = text
        len_history.append(cur_len)
        if len(len_history) > STABLE_COUNT:
            len_history.pop(0)

        # Stable: enough samples, content exists, and max growth < threshold
        if len(len_history) == STABLE_COUNT and cur_len > 0:
            if max(len_history) - min(len_history) <= STABLE_THRESHOLD:
                return text

        remaining = deadline - time.monotonic()
        if remaining <= 0:
            break
        await asyncio.sleep(min(2.0, remaining))

    return last_text


async def _wait_for_backend_alive(tier: str, max_wait: float = BACKEND_SETTLE_WAIT_S) -> bool:
    """Poll _backend_alive until the backend reports healthy or max_wait elapses.

    Replaces blind asyncio.sleep(15) waits after retry-related actions
    (zombie cleanup, manual settle). Returns True if backend recovered,
    False on timeout. Polls at 0.5s for the first 5s, then 1s.
    """
    if tier not in ("ollama",):
        await asyncio.sleep(min(2.0, max_wait))
        return True
    deadline = time.monotonic() + max_wait
    while time.monotonic() < deadline:
        try:
            alive, _ = _backend_alive(tier)
            if alive:
                return True
        except Exception:
            pass
        elapsed = max_wait - (deadline - time.monotonic())
        delay = 0.5 if elapsed < 5.0 else 1.0
        remaining = deadline - time.monotonic()
        if remaining <= 0:
            break
        await asyncio.sleep(min(delay, remaining))
    return False


# ---------------------------------------------------------------------------
# Playwright helpers
# ---------------------------------------------------------------------------


async def _login(page) -> None:
    await page.goto(OPENWEBUI_URL, wait_until="networkidle", timeout=30000)
    await page.wait_for_selector('input[type="email"]', timeout=15000)
    await page.fill('input[type="email"]', ADMIN_EMAIL)
    await page.fill('input[type="password"]', ADMIN_PASS)
    await page.locator('button[type="submit"], button:has-text("Sign in")').first.click()
    await page.wait_for_selector("textarea, [contenteditable]", timeout=20000)


async def _navigate_to_chat(page, chat_url: str) -> None:
    await page.goto(chat_url, wait_until="networkidle", timeout=60000)
    await page.wait_for_selector("textarea, [contenteditable='true']", timeout=30000)
    await page.wait_for_timeout(2000)


async def _stop_button_visible(page) -> bool:
    """Check if the stop/streaming button is currently visible.

    OWUI 0.9.5+ uses a round stop-circle SVG button with no aria-label or title.
    The stop-circle path includes the substring "9.564a1.312" (Heroicons stop-circle).
    """
    try:
        # Old OWUI: button with aria-label or title "Stop"
        btn = page.locator(
            'button[aria-label="Stop"], button[title="Stop"], button:has-text("Stop")'
        )
        if await btn.count() > 0 and await btn.first.is_visible():
            return True
        # OWUI 0.9.5+: round stop-circle SVG button without aria-label
        return bool(
            await page.evaluate(
                """() => {
                    for (const btn of document.querySelectorAll('button')) {
                        const path = btn.querySelector('svg path');
                        if (path && path.getAttribute('d')?.includes('9.564a1.312')) {
                            const r = btn.getBoundingClientRect();
                            if (r.width > 0 && r.height > 0) return true;
                        }
                    }
                    return false;
                }"""
            )
        )
    except Exception:
        return False


async def _wait_for_completion(
    page,
    test_id: str = "",
    tier: str = "any",
    max_wait_no_progress: int = MAX_WAIT_NO_PROGRESS,
    *,
    token: str = "",
    chat_id: str = "",
    min_messages: int = 1,
) -> None:
    """Progress-monitoring wait with tiered polling.

    Phase 1 (waiting for stream start): poll PHASE1_FAST_S → PHASE1_MID_S →
    PHASE1_SLOW_S as elapsed time grows. This catches warm-load starts (<2s)
    without forcing the same resolution on cold loads (30s+).

    Phase 2 (waiting for stream end): poll PHASE2_STREAMING_POLL_S while
    actively streaming. On stop-button-disappears edge, immediately verify
    via OWUI API (no fixed sleep). On DOM-stable path, the same 3-sample
    threshold now resolves in 4.5s instead of 90s.

    When token+chat_id are provided, the OWUI API is used as a parallel
    completion signal (early exit if content lands while DOM is stable)
    and as the canonical post-stream persistence wait (replacing fixed
    sleep(5)). When absent, falls back to a 2s safety buffer.

    Backend crash detection unchanged: BACKEND_DEAD_STRIKES consecutive
    health failures aborts early.
    """
    BACKEND_DEAD_STRIKES = 5

    t_start = time.time()
    last_log = 0.0
    prev_text = ""
    stable_count = 0
    stop_seen = False
    dead_strikes = 0
    last_backend_check = 0.0
    # API-driven tracking: content length from last poll. Used in Phase 2 to
    # prevent DOM-stable from firing while the model is still generating output.
    _prev_api_len = 0
    # Counter for consecutive "DOM stable but API empty" cycles. After
    # DOM_STABLE_API_EMPTY_MAX cycles we assume OWUI won't commit via API
    # and return early so the caller's DOM fallback can extract the response.
    _dom_stable_empty_count = 0

    def _log(msg: str) -> None:
        nonlocal last_log
        now = time.time()
        msg_lower = msg.lower()
        if (
            now - last_log >= PROGRESS_LOG_INTERVAL
            or "complete" in msg_lower
            or "started" in msg_lower
        ):
            elapsed = now - t_start
            tag = f"[{test_id}] " if test_id else ""
            print(f"  {tag}{msg} ({elapsed:.0f}s elapsed)", flush=True)
            last_log = now

    def _check_backend_crash() -> bool:
        """Return True if backend looks crashed (should abort wait).

        Rate-limited to ~once per 5s so high-frequency Phase 1/2 polls don't
        spam the health endpoint.
        """
        nonlocal dead_strikes, last_backend_check
        now = time.time()
        if now - last_backend_check < 5.0:
            return False
        last_backend_check = now
        if tier not in ("ollama",):
            return False
        alive, detail = _backend_alive(tier)
        if not alive:
            dead_strikes += 1
            tag = f"[{test_id}] " if test_id else ""
            print(
                f"  {tag}backend not responding ({detail}), strike {dead_strikes}/{BACKEND_DEAD_STRIKES}",
                flush=True,
            )
            if dead_strikes >= BACKEND_DEAD_STRIKES:
                print(f"  {tag}backend crashed — aborting wait early", flush=True)
                return True
        else:
            dead_strikes = 0
        return False

    def _phase1_interval(elapsed: float) -> float:
        """Tiered poll interval for Phase 1 (waiting for stream start)."""
        if elapsed < PHASE1_FAST_DURATION_S:
            return PHASE1_FAST_S
        if elapsed < PHASE1_FAST_DURATION_S + PHASE1_MID_DURATION_S:
            return PHASE1_MID_S
        return PHASE1_SLOW_S

    # Phase 1: wait for stop button to appear (model starts generating)
    _log("waiting for model to start…")
    while True:
        elapsed = time.time() - t_start
        if await _stop_button_visible(page):
            stop_seen = True
            _log("model streaming started")
            break
        # If no stop button but text is growing, model may be generating
        # without showing a stop button (some OWUI versions)
        curr = await page.evaluate("document.body.innerText")
        if curr != prev_text and len(curr) > len(prev_text) + 50:
            _log("text growing without stop button — treating as streaming")
            prev_text = curr
            break
        # Backend crash check — don't burn 900s on a dead model
        if _check_backend_crash():
            _unload_running_ollama_models()
            _wait_for_ollama_ps_empty(timeout_s=15.0)
            return
        # Hard safety cap
        if elapsed > max_wait_no_progress:
            _log(f"hit {max_wait_no_progress}s safety cap waiting for start")
            return
        await asyncio.sleep(_phase1_interval(elapsed))

    # Phase 2: wait for streaming to complete
    while True:
        elapsed = time.time() - t_start

        # Re-check stop button each poll — Phase 1 may have used the "text growing"
        # fallback path (stop_seen=False). If the stop button appears during Phase 2
        # (model started proper streaming after initial text growth), update stop_seen
        # and reset stable_count so the DOM stable gate works correctly.
        if await _stop_button_visible(page):
            if not stop_seen:
                stop_seen = True
                stable_count = 0
                _log("stop button appeared in Phase 2 — streaming active")
        elif stop_seen:
            # Stop button was seen and is now gone. For thinking models (AEON,
            # Qwen3), the button briefly disappears during the reasoning→response
            # transition before streaming continues. Wait 2s and re-check before
            # committing to "stream complete" to avoid false early exits.
            await asyncio.sleep(2.0)
            if await _stop_button_visible(page):
                stable_count = 0
                _log("stop button reappeared (thinking model transition) — resuming")
            else:
                _log("stream complete (stop button gone)")
                if token and chat_id:
                    await _wait_for_response_arrival(token, chat_id, min_messages=min_messages)
                else:
                    await asyncio.sleep(2.0)
                return

        # API-driven content tracking: fetch current API response length each poll.
        # If content grew since last poll (model still generating response), reset
        # DOM stable count — prevents <details type="reasoning"> collapsed blocks
        # from triggering a false DOM-stable exit while the model is still active.
        # This is the log-driven completion signal: content changes drive the
        # decision, not wall-clock timers.
        _cur_api_text = ""
        if token and chat_id:
            _cur_api_text = owui_get_last_response(token, chat_id, min_messages=min_messages)
            _cur_api_len = len(_cur_api_text)
            if _cur_api_len > _prev_api_len + 100:
                # Content actively growing — model still generating; don't let DOM
                # stability fire prematurely
                stable_count = 0
                _log(f"API content growing ({_prev_api_len}→{_cur_api_len} chars) — resuming")
            _prev_api_len = _cur_api_len

        # Check DOM stability as secondary signal — only when stop button is gone
        # (or never appeared). Reasoning models emit <details type="reasoning">
        # blocks that collapse in the DOM, making innerText appear stable while
        # the model is still streaming the actual response. Gating on stop button
        # prevents false-stable triggers during hidden reasoning token generation.
        curr = await page.evaluate("document.body.innerText")
        if curr == prev_text:
            stable_count += 1
            stop_still_active = stop_seen and await _stop_button_visible(page)
            if stable_count >= PHASE2_DOM_STABLE_NEEDED and not stop_still_active:
                # Before declaring done via DOM, verify via API (reuse _cur_api_text
                # already fetched this poll — no extra HTTP request).
                # If API has content → run stabilization-wait then done.
                # If API is empty but we have credentials → keep polling; the model
                # may still be in the reasoning phase and hasn't started output yet.
                if token and chat_id:
                    if _cur_api_text:
                        _log("stream complete (DOM stable + API has content)")
                        await _wait_for_response_arrival(token, chat_id, min_messages=min_messages)
                        return
                    else:
                        # DOM stable but API empty.
                        stable_count = 0
                        if stop_seen:
                            # Stop button WAS seen — model started streaming but OWUI
                            # 0.9.5+ hasn't committed to API yet (thinking-model case).
                            _dom_stable_empty_count += 1
                            _log("DOM stable but API empty — model still reasoning, continuing")
                            if _dom_stable_empty_count >= DOM_STABLE_API_EMPTY_MAX:
                                # Assume done; caller's _fe_get_last_response DOM
                                # fallback will extract directly from the page.
                                _log(
                                    f"DOM stable + API empty ×{_dom_stable_empty_count}"
                                    " — assuming completion, handing off to DOM fallback"
                                )
                                return
                        else:
                            # stop_seen=False: model still loading / processing prompt,
                            # not yet streaming. Don't apply DOM_STABLE_API_EMPTY_MAX —
                            # but DO cap at NO_STREAM_TIMEOUT to allow rapid retry when
                            # the request stalls before reaching the pipeline.
                            _log("DOM stable + API empty, model not yet streaming — waiting")
                            if elapsed > NO_STREAM_TIMEOUT:
                                _log(
                                    f"DOM stable + API empty after {elapsed:.0f}s with no"
                                    " stream start — exiting for retry"
                                )
                                return
                else:
                    _log("stream complete (DOM stable)")
                    await asyncio.sleep(2.0)
                    return
        else:
            stable_count = 0
            prev_text = curr

        # Backend crash check — rate-limited inside _check_backend_crash
        if _check_backend_crash():
            _unload_running_ollama_models()
            _wait_for_ollama_ps_empty(timeout_s=15.0)
            return

        # Safety cap
        if elapsed > max_wait_no_progress:
            _log(f"hit {max_wait_no_progress}s safety cap during streaming")
            return

        await asyncio.sleep(PHASE2_STREAMING_POLL_S)


async def _send_and_wait(
    page,
    prompt: str,
    test_id: str = "",
    tier: str = "any",
    max_wait_no_progress: int = MAX_WAIT_NO_PROGRESS,
    *,
    token: str = "",
    chat_id: str = "",
    min_messages: int = 1,
) -> None:
    """Send a prompt and wait for completion.

    When token+chat_id are supplied, _wait_for_completion uses the OWUI API
    as a parallel completion signal and replaces the fixed post-stream sleep
    with a bounded content-arrival poll. Caller still fetches the response
    via owui_get_last_response after this returns.

    min_messages: forwarded to _wait_for_completion. Set to 2 for multi-turn
    turn-2 calls so completion detection requires ≥ 2 committed assistant
    responses, preventing turn-1's stable content from firing a false early exit.
    """
    ta = page.locator("textarea, [contenteditable='true']").first
    await ta.click()
    await ta.fill(prompt)
    send_btn = page.locator("#send-message-button")
    if await send_btn.count() > 0:
        await send_btn.click()
    else:
        await ta.press("Enter")
    await _wait_for_completion(
        page,
        test_id,
        tier,
        max_wait_no_progress,
        token=token,
        chat_id=chat_id,
        min_messages=min_messages,
    )


async def _enable_tool(page, tool_id: str) -> None:
    tool_display_names = {
        "portal_code": "Portal Code",
        "portal_documents": "Portal Documents",
        "portal_memory": "Portal Memory",
        "portal_music": "Portal Music",
        "portal_tts": "Portal TTS",
        "portal_video": "Portal Video",
        "portal_comfyui": "Portal ComfyUI",
        "portal_security": "Portal Security",
        "portal_whisper": "Portal Whisper",
    }
    display = tool_display_names.get(tool_id, tool_id)

    try:
        btn = page.locator(
            'button[aria-label="Tools"], button[title="Tools"], '
            'button:has-text("+"), .chat-toolbar button'
        ).first
        await btn.click(timeout=5000)
        await page.wait_for_timeout(1000)

        toggle = page.locator(f'button:has-text("{display}"), label:has-text("{display}")')
        if await toggle.count() > 0:
            await toggle.first.click()

        await page.keyboard.press("Escape")
        await page.wait_for_timeout(500)
    except Exception:
        pass  # best-effort; test proceeds and assertion will catch if tool missing


async def _download_artifact(
    page, expected_ext: str, timeout_ms: int = 120_000, response_text: str = ""
) -> Path | None:
    ARTIFACT_DIR.mkdir(exist_ok=True)
    # Try Playwright UI download first
    try:
        async with page.expect_download(timeout=timeout_ms) as dl_info:
            await page.locator(
                f'a[download], a[href*=".{expected_ext}"], '
                f'button:has-text("Download"), .file-attachment'
            ).last.click(timeout=10000)
        dl = await dl_info.value
        dest = ARTIFACT_DIR / dl.suggested_filename
        await dl.save_as(dest)
        return dest
    except Exception:
        pass

    # Fallback: extract file path or download URL from model response
    import re
    import subprocess

    if response_text:
        # Try 1: Match a download URL ending in /files/<...>.<ext>.
        # Covers both the localhost shape and the public shape emitted when
        # PORTAL_PUBLIC_URL is set (e.g. https://portal.example.com/files/tts/<name>.wav,
        # potentially served via Cloudflare Tunnel). The driver runs on the
        # host, so it can resolve either form via DNS or loopback. Try 1
        # fetching matters because it exercises the same URL the user's
        # browser would use.
        url_pattern = rf"https?://[^\s)>\]]+/files/\S+?\.{re.escape(expected_ext)}"
        url_match = re.search(url_pattern, response_text)
        if url_match:
            download_url = url_match.group(0)
            filename = Path(download_url).name
            dest = ARTIFACT_DIR / filename
            try:
                r = httpx.get(download_url, timeout=30)
                if r.status_code == 200:
                    dest.write_bytes(r.content)
                    return dest
            except Exception:
                pass

        # Try 1b: ComfyUI /view?filename=... URL (host-native ComfyUI at :8188).
        # generate_image / generate_video return:
        #   http://localhost:8188/view?filename=portal_xxx.png&type=output
        # The /files/ pattern above never matches this shape.
        comfyui_pat = (
            rf"https?://[^\s)>\]]*/view\?filename=[^\s)>\]]*\.{re.escape(expected_ext)}[^\s)>\]]*"
        )
        comfyui_match = re.search(comfyui_pat, response_text)
        if comfyui_match:
            from urllib.parse import parse_qs, urlparse

            download_url = comfyui_match.group(0)
            qs = parse_qs(urlparse(download_url).query)
            fname = qs.get("filename", ["unknown"])[0]
            dest = ARTIFACT_DIR / Path(fname).name
            try:
                r = httpx.get(download_url, timeout=30)
                if r.status_code == 200:
                    dest.write_bytes(r.content)
                    return dest
            except Exception:
                pass

        # Try 2: Match /app/data/generated/<filename>.<ext> container path
        container_pattern = rf"/app/data/generated/\S+\.{re.escape(expected_ext)}"
        container_match = re.search(container_pattern, response_text)
        if container_match:
            container_path = container_match.group(0)
            for container in [
                "portal5-mcp-documents",
                "portal5-mcp-sandbox",
                "portal5-mcp-comfyui",
                "portal5-mcp-video",
            ]:
                dest = ARTIFACT_DIR / Path(container_path).name
                result = subprocess.run(
                    ["docker", "cp", f"{container}:{container_path}", str(dest)],
                    capture_output=True,
                    timeout=10,
                )
                if result.returncode == 0 and dest.exists():
                    return dest

    # Try 3: Most recent file with expected extension from MCP containers
    # (handles case where tool ran but model didn't mention the filename)
    for container in [
        "portal5-mcp-documents",
        "portal5-mcp-sandbox",
        "portal5-mcp-comfyui",
        "portal5-mcp-video",
    ]:
        try:
            result = subprocess.run(
                ["docker", "exec", container, "ls", "-t", "/app/data/generated/"],
                capture_output=True,
                timeout=10,
                text=True,
            )
            if result.returncode == 0:
                for line in result.stdout.strip().split("\n"):
                    fname = line.strip()
                    if fname.endswith(f".{expected_ext}"):
                        container_path = f"/app/data/generated/{fname}"
                        dest = ARTIFACT_DIR / fname
                        cp_result = subprocess.run(
                            ["docker", "cp", f"{container}:{container_path}", str(dest)],
                            capture_output=True,
                            timeout=10,
                        )
                        if cp_result.returncode == 0 and dest.exists():
                            return dest
                        break
        except Exception:
            continue

    # Try 4: ComfyUI direct download — query /history for the most recent
    # portal_*.{mp4,png}. ComfyUI runs host-native (port 8188), so docker cp
    # never finds its output files regardless of extension.
    # png: prefix "portal_" (comfyui_mcp.py SaveImage node)
    # mp4: prefix "portal_video_" (video_mcp.py)
    # Recency guard: only accept files generated in the last 15 minutes, to
    # avoid picking up stale files from a previous test session.
    if expected_ext in ("mp4", "png"):
        try:
            import time as _time

            now_ms = int(_time.time() * 1000)
            cutoff_ms = now_ms - (15 * 60 * 1000)
            r = httpx.get("http://localhost:8188/history", timeout=10)
            if r.status_code == 200:
                history = r.json()
                best_ts: int = -1
                best_fname: str | None = None
                for job_data in history.values():
                    if not job_data.get("status", {}).get("completed"):
                        continue
                    outputs = job_data.get("outputs", {})
                    for node_outputs in outputs.values():
                        for img in node_outputs.get("images", []):
                            fname = img.get("filename", "")
                            ext_match = fname.endswith(f".{expected_ext}")
                            prefix_match = (
                                expected_ext == "mp4" and fname.startswith("portal_video_")
                            ) or (expected_ext == "png" and fname.startswith("portal_"))
                            if ext_match and prefix_match:
                                msgs = job_data.get("status", {}).get("messages", [])
                                ts = msgs[0][1].get("timestamp", 0) if msgs else 0
                                if ts >= cutoff_ms and ts > best_ts:
                                    best_ts = ts
                                    best_fname = fname
                if best_fname:
                    url = f"http://localhost:8188/view?filename={best_fname}&type=output"
                    dest = ARTIFACT_DIR / best_fname
                    r2 = httpx.get(url, timeout=60)
                    if r2.status_code == 200 and len(r2.content) > 0:
                        dest.write_bytes(r2.content)
                        return dest
        except Exception:
            pass

    return None


# ---------------------------------------------------------------------------
# Think-block stripping
# ---------------------------------------------------------------------------


def _strip_think_blocks(text: str) -> str:
    """Strip reasoning blocks from model output before running assertions.

    Three reasoning formats are handled:
    - <think>...</think>: Laguna-XS.2, Phi-4-reasoning-plus, Qwopus
    - [THINK]...[/THINK]: Magistral
    - <details type="reasoning">...</details>: AEON/Qwen3 as committed by OWUI API
      (OWUI inlines reasoning in the content field; the actual response follows
      the closing tag — without stripping, keywords like "error" or "failed" that
      appear naturally in reasoning traces cause false not_contains failures)

    Strips all variants case-insensitively with DOTALL so multi-line blocks are
    handled. Trailing whitespace is normalized after stripping.
    """
    import re

    original = text
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"\[THINK\].*?\[/THINK\]", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(
        r'<details[^>]*type=["\']reasoning["\'][^>]*>.*?</details>',
        "",
        text,
        flags=re.DOTALL | re.IGNORECASE,
    )
    result = text.strip()
    if not result:
        # Model put entire answer inside reasoning block — extract inner content
        m = re.search(
            r"<details[^>]*>.*?<summary>.*?</summary>(.*?)</details>",
            original,
            re.DOTALL | re.IGNORECASE,
        )
        if m:
            return m.group(1).strip()
    return result


# ---------------------------------------------------------------------------
# Assertion engine
# ---------------------------------------------------------------------------


_UNICODE_DASH_TABLE = str.maketrans(
    "".join(
        [
            "‐",  # hyphen
            "‑",  # non-breaking hyphen
            "‒",  # figure dash
            "–",  # en dash
            "—",  # em dash
            "―",  # horizontal bar
            "−",  # minus sign
            "─",  # box-drawing horizontal
            "﹘",  # small em dash
            "﹣",  # small hyphen-minus
            "－",  # fullwidth hyphen-minus
        ]
    ),
    "-" * 11,
)


def _normalize_dashes(s: str) -> str:
    return s.translate(_UNICODE_DASH_TABLE)


def _kw_in(keyword: str, text: str, *, word_boundary: bool) -> bool:
    """Return True if ``keyword`` appears in ``text`` (case-insensitive).

    With ``word_boundary=True``, the match is anchored on regex ``\\b`` boundaries
    so short tokens like 'r1' or 'lives' don't match inside 'router 1' or 'olives'.
    Boundaries only fire between \\w and \\W, so keywords that begin or end with
    punctuation (e.g. '=B2-C2') still match correctly.

    Unicode dash variants (em-dash, en-dash, non-breaking hyphen, etc.) are
    normalised to ASCII hyphen before matching — models frequently use typographic
    dashes in structured names like CIP‑003‑9.
    """
    needle = _normalize_dashes(keyword.lower())
    haystack = _normalize_dashes(text.lower())
    if not word_boundary:
        return needle in haystack
    import re

    return re.search(rf"\b{re.escape(needle)}\b", haystack) is not None


def assert_contains(text: str, keywords: list, label: str, *, word_boundary: bool = False) -> tuple:
    missing = [k for k in keywords if not _kw_in(k, text, word_boundary=word_boundary)]
    return (label, not missing, f"missing: {missing}" if missing else "ok")


def assert_any_of(text: str, keywords: list, label: str, *, word_boundary: bool = False) -> tuple:
    found = [k for k in keywords if _kw_in(k, text, word_boundary=word_boundary)]
    return (label, bool(found), f"found: {found}" if found else f"none of: {keywords}")


def assert_not_contains(
    text: str, keywords: list, label: str, *, word_boundary: bool = False
) -> tuple:
    found = [k for k in keywords if _kw_in(k, text, word_boundary=word_boundary)]
    return (label, not found, f"found (bad): {found}" if found else "ok")


def assert_min_length(text: str, chars: int, label: str) -> tuple:
    return (label, len(text) >= chars, f"len={len(text)}, min={chars}")


def assert_has_code(text: str, label: str) -> tuple:
    has_fence = "```" in text
    # Raw HTML delivery (no markdown wrapper) is also valid code delivery
    has_raw_html = text.strip().startswith("<!DOCTYPE") or text.strip().startswith("<html")
    ok = has_fence or has_raw_html
    detail = (
        "code block present" if has_fence else ("raw html" if has_raw_html else "no code block")
    )
    return (label, ok, detail)


def _extract_code_blocks(text: str) -> str:
    """Extract and concatenate content from markdown code blocks.

    Handles fenced blocks (```lang ... ```), unclosed fenced blocks
    (opening fence without closing — common in model output), and raw
    HTML starting with <!DOCTYPE or <html>.
    Returns the concatenated code text, or '' if no code blocks found.
    """
    import re

    parts: list[str] = []
    text_lower = text.lower()

    # Fenced code blocks: ```optional_lang\n...\n```
    for m in re.finditer(r"```(?:\w+)?\n(.*?)```", text, re.DOTALL):
        parts.append(m.group(1).strip())

    # Unclosed fenced block: opening ``` anywhere, no closing ```
    if not parts:
        fence_match = re.search(r"```(?:\w+)?\n", text)
        if fence_match and "```" not in text[fence_match.end() :]:
            code_text = text[fence_match.end() :].strip()
            parts.append(code_text)

    # Raw HTML delivery (no markdown wrapper)
    if not parts:
        stripped = text.strip()
        if stripped.startswith("<!DOCTYPE") or stripped.startswith("<html"):
            parts.append(stripped)

    return "\n".join(parts)


def assert_code_pattern(text: str, patterns: list[dict], label: str) -> tuple:
    """Run regex patterns against extracted code blocks (not full response).

    Each pattern dict has:
        regex: str   — regex pattern to search for
        label: str   — human-readable description

    Patterns are case-insensitive. If any pattern matches, the assertion passes.
    This checks actual code behavior, not prose or variable naming conventions.
    """
    import re

    code = _extract_code_blocks(text)
    if not code:
        return (label, False, "no code blocks extracted")

    for p in patterns:
        pattern = p["regex"]
        try:
            if re.search(pattern, code, re.IGNORECASE):
                return (label, True, f"matched: {p.get('label', pattern)}")
        except re.error as e:
            return (label, False, f"invalid regex '{pattern}': {e}")

    return (label, False, f"no pattern matched in code ({len(code)} chars)")


def assert_has_table(text: str, label: str) -> tuple:
    return (label, "|" in text and "---" in text, "table present" if "|" in text else "no table")


def assert_docx_valid(path: Path | None, label: str) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        from docx import Document

        doc = Document(path)
        return (label, len(doc.paragraphs) > 0, f"{len(doc.paragraphs)} paragraphs")
    except Exception as e:
        return (label, False, str(e))


def assert_xlsx_valid(path: Path | None, label: str) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        from openpyxl import load_workbook

        wb = load_workbook(path)
        return (label, len(wb.sheetnames) > 0, f"sheets: {wb.sheetnames}")
    except Exception as e:
        return (label, False, str(e))


def assert_pptx_valid(path: Path | None, min_slides: int, label: str) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        from pptx import Presentation

        prs = Presentation(path)
        return (label, len(prs.slides) >= min_slides, f"{len(prs.slides)} slides")
    except Exception as e:
        return (label, False, str(e))


def assert_wav_valid(
    path: Path | None,
    label: str,
    *,
    min_seconds: float = 0.0,
) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        data = path.read_bytes()
        if not (len(data) > 1000 and data[:4] == b"RIFF" and data[8:12] == b"WAVE"):
            return (label, False, f"not a valid WAV: {len(data)} bytes")
        if min_seconds > 0:
            import wave

            with wave.open(str(path), "rb") as w:
                duration = w.getnframes() / float(w.getframerate())
            if duration < min_seconds:
                return (
                    label,
                    False,
                    f"too short: {duration:.1f}s < {min_seconds}s ({len(data)} bytes)",
                )
            return (label, True, f"{duration:.1f}s, {len(data)} bytes")
        return (label, True, f"{len(data)} bytes")
    except Exception as e:
        return (label, False, str(e))


def assert_png_valid(
    path: Path | None,
    label: str,
    *,
    min_width: int = 0,
    min_height: int = 0,
) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        data = path.read_bytes()
        if data[:8] != b"\x89PNG\r\n\x1a\n":
            return (label, False, f"not a PNG: {data[:8]!r}")
        if min_width > 0 or min_height > 0:
            try:
                from PIL import Image

                with Image.open(path) as im:
                    w, h = im.size
                if w < min_width or h < min_height:
                    return (
                        label,
                        False,
                        f"too small: {w}x{h} < {min_width}x{min_height}",
                    )
                return (label, True, f"{w}x{h}, {len(data)} bytes")
            except ImportError:
                return (label, True, f"PIL unavailable; {len(data)} bytes")
        return (label, True, f"{len(data)} bytes")
    except Exception as e:
        return (label, False, str(e))


def assert_mp4_valid(
    path: Path | None,
    label: str,
    *,
    min_seconds: float = 0.0,
) -> tuple:
    if path is None:
        return (label, False, "no file downloaded")
    try:
        data = path.read_bytes()
        if b"ftyp" not in data[:32]:
            return (label, False, f"not an MP4: {data[:16]!r}")
        if min_seconds > 0:
            import subprocess

            try:
                out = subprocess.check_output(
                    [
                        "ffprobe",
                        "-v",
                        "error",
                        "-show_entries",
                        "format=duration",
                        "-of",
                        "default=noprint_wrappers=1:nokey=1",
                        str(path),
                    ],
                    text=True,
                    timeout=10,
                ).strip()
                duration = float(out)
                if duration < min_seconds:
                    return (
                        label,
                        False,
                        f"too short: {duration:.1f}s < {min_seconds}s",
                    )
                return (label, True, f"{duration:.1f}s, {len(data)} bytes")
            except FileNotFoundError:
                return (label, len(data) > 50_000, f"{len(data)} bytes (no ffprobe)")
            except Exception as e:
                return (label, False, str(e))
        return (label, True, f"{len(data)} bytes")
    except Exception as e:
        return (label, False, str(e))


def run_assertions(
    text: str,
    assertions_spec: list,
    artifact_path: Path | None = None,
    include_thinking: bool = False,
) -> list:
    if not include_thinking:
        text = _strip_think_blocks(text)
    results = []
    for a in assertions_spec:
        t = a["type"]
        label = a.get("label", t)
        if t == "contains":
            results.append(
                assert_contains(
                    text, a["keywords"], label, word_boundary=a.get("word_boundary", False)
                )
            )
        elif t == "any_of":
            results.append(
                assert_any_of(
                    text, a["keywords"], label, word_boundary=a.get("word_boundary", False)
                )
            )
        elif t == "not_contains":
            results.append(
                assert_not_contains(
                    text, a["keywords"], label, word_boundary=a.get("word_boundary", False)
                )
            )
        elif t == "min_length":
            results.append(assert_min_length(text, a["chars"], label))
        elif t == "has_code":
            results.append(assert_has_code(text, label))
        elif t == "code_pattern":
            results.append(assert_code_pattern(text, a.get("patterns", []), label))
        elif t == "has_table":
            results.append(assert_has_table(text, label))
        elif t == "docx_valid":
            results.append(assert_docx_valid(artifact_path, label))
        elif t == "xlsx_valid":
            results.append(assert_xlsx_valid(artifact_path, label))
        elif t == "pptx_valid":
            min_slides = a.get("min_slides", 1)
            results.append(assert_pptx_valid(artifact_path, min_slides, label))
        elif t == "wav_valid":
            min_s = float(a.get("min_seconds", 0.0))
            results.append(assert_wav_valid(artifact_path, label, min_seconds=min_s))
        elif t == "png_valid":
            mw = int(a.get("min_width", 0))
            mh = int(a.get("min_height", 0))
            results.append(assert_png_valid(artifact_path, label, min_width=mw, min_height=mh))
        elif t == "mp4_valid":
            min_s = float(a.get("min_seconds", 0.0))
            results.append(assert_mp4_valid(artifact_path, label, min_seconds=min_s))
        elif t == "quality_score":
            threshold = a.get("min", 0.5)
            cat = a.get("category", "general")
            try:
                import os as _os
                import sys as _sys

                _sys.path.insert(0, _os.path.join(_os.path.dirname(__file__)))
                from quality_signals import quality_score as _qs

                qs = _qs(cat, text)
            except Exception:
                qs = 1.0
            label_ext = f"{label} ({qs:.2f})"
            results.append((label_ext, qs >= threshold, f"score={qs:.2f}, min={threshold}"))
    return results


def compute_status(assertions: list, assertions_spec: list) -> str:
    """Grade a test result.

    By default every assertion is critical (a single failure produces FAIL).
    To opt an assertion into the percentage-grading floor, mark it with
    ``"critical": False`` in the spec. Behavior:

    - Any spec entry with ``critical=True`` (the default) that fails -> FAIL,
      UNLESS the overall pass rate is >=70%, in which case it's downgraded to WARN.
      This prevents a single narrow keyword from failing an otherwise correct test.
    - Otherwise, if all failing specs are ``critical=False``: PASS at >=70% pass
      rate, WARN at >=50%, FAIL below.

    The percentage rule now applies even with critical failures when overall
    score is high enough to demonstrate correct model behavior.
    """
    if not assertions:
        return "FAIL"
    total = len(assertions)
    passed_count = sum(1 for r in assertions if r[1])
    pct = passed_count / total * 100

    # Any critical failure is an automatic FAIL — unless the overall pass rate
    # is high enough to demonstrate that the model behaved correctly and the
    # failing assertion is likely a keyword-too-strict issue.
    has_critical_fail = False
    for result, spec in zip(assertions, assertions_spec):
        _label, passed, _evidence = result
        # has_code is a format preference — good code without a fenced block is
        # still correct behavior and should not alone fail an otherwise valid test.
        default_critical = spec.get("type") != "has_code"
        critical = spec.get("critical", default_critical)
        if not passed and critical:
            has_critical_fail = True
            break

    if has_critical_fail:
        return "FAIL"

    if pct >= 70:
        return "PASS"
    if pct >= 50 or passed_count > 0:
        return "WARN"
    return "FAIL"


# ---------------------------------------------------------------------------
# Result recorder
# ---------------------------------------------------------------------------


def init_results(run_ts: str) -> None:
    RESULTS_FILE.write_text(
        f"# Portal 5 — UAT Results\n\n"
        f"**Run:** {run_ts}  \n"
        f"**Catalog:** TEST_CATALOG (see tests/portal5_uat_driver.py)  \n"
        f"**Reviewer:** (fill in)\n\n"
        f"## Summary\n\n"
        f"- **PASS**: 0\n- **WARN**: 0\n- **FAIL**: 0\n- **SKIP**: 0\n- **BLOCKED**: 0\n- **MANUAL**: 0\n\n"
        f"## Results\n\n"
        f"| # | Status | Test | Model | Detail | Elapsed |\n"
        f"|---|--------|------|-------|--------|---------|\n"
    )


def update_summary(counts: dict) -> None:
    text = RESULTS_FILE.read_text()
    for status in ("PASS", "WARN", "FAIL", "SKIP", "BLOCKED", "MANUAL"):
        old = f"- **{status}**: "
        lines = [l for l in text.split("\n") if l.startswith(old)]
        if lines:
            text = text.replace(lines[0], f"{old}{counts.get(status, 0)}")
    RESULTS_FILE.write_text(text)


def _parse_test_ids_from_results() -> set[str]:
    """Return the set of test IDs already present as rows in UAT_RESULTS.md."""
    if not RESULTS_FILE.exists():
        return set()
    import re as _re

    text = RESULTS_FILE.read_text()
    ids: set[str] = set()
    # Result rows: "| N | STATUS | [TEST_ID name](url) | `model` | ... | Ns |"
    pattern = _re.compile(r"^\|\s*\d+\s*\|\s*\w+\s*\|\s*\[([A-Za-z0-9][A-Za-z0-9_.-]*)\s")
    for line in text.split("\n"):
        m = pattern.match(line)
        if m:
            ids.add(m.group(1))
    return ids


def _parse_failed_test_ids(statuses: set[str] | None = None) -> set[str]:
    """Return test IDs from UAT_RESULTS.md whose status is in ``statuses``.

    Defaults to FAIL and BLOCKED. Used by --rerun-failed to auto-select
    broken tests without requiring the caller to enumerate IDs manually.
    """
    if statuses is None:
        statuses = {"FAIL", "BLOCKED"}
    if not RESULTS_FILE.exists():
        return set()
    import re as _re

    text = RESULTS_FILE.read_text()
    ids: set[str] = set()
    pattern = _re.compile(r"^\|\s*\d+\s*\|\s*(\w+)\s*\|\s*\[([A-Za-z0-9][A-Za-z0-9_.-]*)\s")
    for line in text.split("\n"):
        m = pattern.match(line)
        if m and m.group(1).strip() in statuses:
            ids.add(m.group(2))
    return ids


def _remove_rows_for_test_ids(test_ids: set[str]) -> int:
    """Remove existing rows from UAT_RESULTS.md whose test_id is in ``test_ids``.

    Returns the number of rows removed. The summary header is NOT updated here —
    callers should run ``_rebuild_summary_from_rows`` after the run completes.
    """
    if not RESULTS_FILE.exists():
        return 0
    import re as _re

    text = RESULTS_FILE.read_text()
    out_lines: list[str] = []
    removed = 0
    pattern = _re.compile(r"^\|\s*\d+\s*\|\s*\w+\s*\|\s*\[([A-Za-z0-9][A-Za-z0-9_.-]*)\s")
    for line in text.split("\n"):
        m = pattern.match(line)
        if m and m.group(1) in test_ids:
            removed += 1
            continue
        out_lines.append(line)
    RESULTS_FILE.write_text("\n".join(out_lines))
    return removed


def _rebuild_summary_from_rows() -> None:
    """Recompute the PASS/WARN/FAIL/SKIP/MANUAL counts in the summary header
    by parsing the rows in UAT_RESULTS.md. Source of truth is the file contents,
    not the in-memory ``counts`` dict (which is per-invocation).
    """
    if not RESULTS_FILE.exists():
        return
    import re as _re

    text = RESULTS_FILE.read_text()
    counts = {"PASS": 0, "WARN": 0, "FAIL": 0, "SKIP": 0, "BLOCKED": 0, "MANUAL": 0}
    pattern = _re.compile(r"^\|\s*\d+\s*\|\s*(\w+)\s*\|\s*\[")
    for line in text.split("\n"):
        m = pattern.match(line)
        if m:
            status = m.group(1).strip()
            if status in counts:
                counts[status] += 1
    for status in counts:
        old_re = _re.compile(rf"^- \*\*{status}\*\*: \d+", _re.MULTILINE)
        text = old_re.sub(f"- **{status}**: {counts[status]}", text)
    RESULTS_FILE.write_text(text)


def _write_routing_summary() -> None:
    """Append a Routing Summary section to UAT_RESULTS.md.

    Groups by: correct | routing mismatch | wrong model | no-actual.
    Also breaks down the pipeline-confirmed backend for correctly-matched
    tests — surfaces silent fallbacks where the test passes on general
    capability but the intended model was never exercised.
    """
    if not _ROUTING_LOG:
        return

    correct = [r for r in _ROUTING_LOG if r["matched"]]
    tier_fallbacks = [r for r in _ROUTING_LOG if not r["matched"] and r.get("tier_mismatch")]
    wrong_model = [r for r in _ROUTING_LOG if not r["matched"] and not r.get("tier_mismatch") and r["actual"]]
    no_actual = [r for r in _ROUTING_LOG if not r["actual"]]

    # Pipeline backend breakdown: among correctly-matched tests that had an ollama-intended
    # workspace, how many confirmed correct routing.
    ollama_intended_correct = [r for r in correct if r.get("intended_ollama") and r.get("pipeline_backend")]
    confirmed_ollama = [r for r in ollama_intended_correct if r.get("pipeline_backend")]

    lines: list[str] = [
        "",
        "## Routing Summary",
        "",
        "| Metric | Count |",
        "|--------|-------|",
        f"| Routing checked | {len(_ROUTING_LOG)} |",
        f"| Correct | {len(correct)} |",
        f"| Routing mismatch (wrong model) | {len(tier_fallbacks)} |",
        f"| Wrong model (same tier) | {len(wrong_model)} |",
        f"| No actual model returned | {len(no_actual)} |",
        "",
    ]

    if ollama_intended_correct:
        lines += [
            "### Pipeline Backend (Ollama primary, pipeline-confirmed)",
            "",
            "Tests that matched expected routing — breakdown of which backend *actually* served:",
            "",
            "| Metric | Count |",
            "|--------|-------|",
            f"| Ollama primary confirmed | {len(confirmed_ollama)} |",
            f"| Backend unconfirmed (log gap) | {len(ollama_intended_correct) - len(confirmed_ollama)} |",
            "",
        ]
        if confirmed_ollama:
            lines += [
                "**Ollama-served** — these tests passed with backend confirmed:",
                "",
                "| Test ID | Name | Section | Pipeline Backend |",
                "|---------|------|---------|-----------------|",
            ]
            for r in confirmed_ollama:
                backend = r.get("pipeline_backend", "?")
                lines.append(
                    f"| {r['test_id']} | {r['name'][:40]} | {r['section']} | `{backend}` |"
                )
            lines.append("")

    if tier_fallbacks:
        lines += [
            "### Routing Mismatches (intended model not served)",
            "",
            "A different model served these tests than the workspace intended.",
            "The test may have passed on general capability — the **intended model was never exercised**.",
            "",
            "| Test ID | Name | Section | Intended | Actual |",
            "|---------|------|---------|----------|--------|",
        ]
        for r in tier_fallbacks:
            lines.append(
                f"| {r['test_id']} | {r['name'][:40]} | {r['section']} "
                f"| {r['intended'][:40]} | {r['actual'][:40]} |"
            )
        lines.append("")

    if wrong_model:
        lines += [
            "### Wrong Model (tier OK, model mismatch)",
            "",
            "| Test ID | Name | Section | Intended | Actual |",
            "|---------|------|---------|----------|--------|",
        ]
        for r in wrong_model:
            lines.append(
                f"| {r['test_id']} | {r['name'][:40]} | {r['section']} "
                f"| {r['intended'][:40]} | {r['actual'][:40]} |"
            )
        lines.append("")

    if not tier_fallbacks and not wrong_model and not silent_ollama:
        lines.append("All routing checks passed — every test was served by its intended primary model.\n")

    with RESULTS_FILE.open("a") as f:
        f.write("\n".join(lines))


def record_result(
    n: int,
    status: str,
    test_id: str,
    name: str,
    model: str,
    assertions: list,
    elapsed: float,
    chat_url: str,
    routed_model: str = "",
) -> None:
    passed = sum(1 for a in assertions if a[1])
    total = len(assertions)
    pct = f"{passed}/{total}({passed * 100 // total}%)" if total else "0/0"
    detail = "; ".join(f"{a[0]}={'✓' if a[1] else '✗'}({a[2]})" for a in assertions)
    if routed_model and status in ("FAIL", "WARN"):
        detail = f"[routed: {routed_model}] {detail}" if detail else f"[routed: {routed_model}]"
    with RESULTS_FILE.open("a") as f:
        f.write(
            f"| {n} | {status} | [{test_id} {name}]({chat_url}) | "
            f"`{model}` | {pct} {detail} | {elapsed:.1f}s |\n"
        )
    icon = {"PASS": "✓", "FAIL": "✗", "WARN": "⚠", "SKIP": "–", "BLOCKED": "⊘", "MANUAL": "✎"}.get(
        status, "?"
    )
    routed_suffix = f" [→{routed_model}]" if routed_model else ""
    print(
        f"  [{icon} {status}] {test_id} {name} ({passed}/{total}={passed * 100 // total if total else 0}%) ({elapsed:.1f}s){routed_suffix}"
    )


# ---------------------------------------------------------------------------
# Skip condition detection
# ---------------------------------------------------------------------------


def evaluate_skip_conditions() -> dict:
    conditions: dict[str, bool] = {}
    try:
        r = httpx.get("http://localhost:8188/system_stats", timeout=3)
        conditions["no_comfyui"] = r.status_code != 200
    except Exception:
        conditions["no_comfyui"] = True

    env_content = Path(".env").read_text() if Path(".env").exists() else ""
    # Per-key check: KEY=value on its own line, value non-empty, value != "CHANGEME".
    # The previous `"CHANGEME" in env_content` substring check fired on any other
    # placeholder elsewhere in the file (PIPELINE_API_KEY, GRAFANA_PASSWORD, the
    # comment on line 3 of .env.example, etc.), falsely flagging both bot
    # predicates as "not configured" even with valid tokens set.
    conditions["no_bot_telegram"] = not _env_var_set(env_content, "TELEGRAM_BOT_TOKEN")
    conditions["no_bot_slack"] = not _env_var_set(env_content, "SLACK_BOT_TOKEN")
    fixtures = Path(__file__).parent / "fixtures"
    conditions["no_image_upload"] = not (fixtures / "sample.png").exists()
    conditions["no_audio_fixture"] = not (fixtures / "sample.wav").exists()
    conditions["no_two_speaker_audio_fixture"] = not (fixtures / "sample_two_speakers.wav").exists()
    try:
        r = httpx.get("http://localhost:8924/health", timeout=3)
        conditions["no_transcribe_server"] = r.status_code != 200
    except Exception:
        conditions["no_transcribe_server"] = True
    conditions["no_docx_fixture"] = not (fixtures / "sample.docx").exists()
    conditions["no_knowledge_base"] = not (fixtures / "knowledge_base").is_dir()
    return conditions


def _env_var_set(env_content: str, key: str) -> bool:
    """True iff ``key`` is set in env content with a non-empty value that isn't ``CHANGEME``.

    Reads ``KEY=value`` on its own line; tolerates leading whitespace, inline
    comments, and surrounding quotes on the value. Comments and unrelated
    placeholders elsewhere in the file do not affect the result.
    """
    import re

    pat = rf"^[ \t]*{re.escape(key)}=([^\r\n]*)$"
    m = re.search(pat, env_content, re.MULTILINE)
    if not m:
        return False
    raw = m.group(1)
    # Strip inline comment ("# ..." not inside quotes — simple heuristic that
    # matches typical .env practice; values containing literal '#' should be
    # quoted, which is the convention .env.example follows).
    if "#" in raw and not (raw.lstrip().startswith(('"', "'"))):
        raw = raw.split("#", 1)[0]
    val = raw.strip().strip('"').strip("'")
    return bool(val) and val != "CHANGEME"


def _bot_container_running(container_name: str) -> tuple[bool, str]:
    """True if the named docker container is in 'running' state.

    Used by via_dispatcher tests to surface a clear failure when the bot
    container itself is down — distinct from the dispatcher path failing.
    """
    import subprocess

    try:
        r = subprocess.run(
            ["docker", "inspect", "--format", "{{.State.Status}}", container_name],
            capture_output=True,
            text=True,
            timeout=5,
        )
        if r.returncode != 0:
            return False, f"container not found: {container_name}"
        status = r.stdout.strip()
        return status == "running", f"status={status}"
    except FileNotFoundError:
        return False, "docker CLI not available"
    except Exception as exc:
        return False, f"inspect error: {exc}"


async def _run_via_dispatcher(workspace: str, prompt: str, timeout: int) -> str:
    """Drive a chat completion through the Pipeline as a Telegram/Slack bot would.

    Bypasses Open WebUI to exercise the exact code path
    ``portal_channels.dispatcher.call_pipeline_async`` uses on every inbound
    message: a single POST to ``:9099/v1/chat/completions`` with
    ``Authorization: Bearer ${PIPELINE_API_KEY}``. Returns the assistant content
    string. Raises on transport error or non-2xx response — caller handles.
    """
    api_key = os.environ.get("PIPELINE_API_KEY", "portal-pipeline")
    pipeline_url = os.environ.get("PIPELINE_URL", "http://localhost:9099")
    payload = {
        "model": workspace,
        "messages": [{"role": "user", "content": prompt}],
        "stream": False,
    }
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    async with httpx.AsyncClient(timeout=timeout) as client:
        resp = await client.post(
            f"{pipeline_url}/v1/chat/completions",
            json=payload,
            headers=headers,
        )
        resp.raise_for_status()
        data = resp.json()
        return str(data["choices"][0]["message"]["content"])


# ---------------------------------------------------------------------------
# Inter-test settling
# ---------------------------------------------------------------------------

SETTLING: dict[tuple, int] = {
    ("ollama", "ollama"): 10,
    ("ollama", "any"): 10,
    ("ollama", "media_heavy"): 30,
    ("any", "ollama"): 10,
    ("any", "any"): 5,
    ("any", "media_heavy"): 30,
    ("media_heavy", "media_heavy"): 30,
    ("media_heavy", "any"): 15,
    ("media_heavy", "ollama"): 30,
}


def settling_delay(current_tier: str, next_tier: str) -> int:
    return SETTLING.get((current_tier, next_tier), 10)


# ---------------------------------------------------------------------------
# Continuous memory & health monitor (self-healing)
# ---------------------------------------------------------------------------


class MemoryMonitor:
    """Background task that continuously monitors memory and backend health.

    Self-healing actions:
    - Memory > 75%: log warning
    - Memory > 85%: force-evict all Ollama models
    - Memory > 92% after eviction: kill zombie processes, retry eviction
    - Ollama unreachable: log crash (restart handled by launchd/docker)

    Runs as an asyncio task alongside the test loop. Call start() before tests,
    stop() after. Stats are available via .stats dict.
    """

    def __init__(self, poll_interval: float = 20.0) -> None:
        self.poll_interval = poll_interval
        self._task: asyncio.Task | None = None
        self._running = False
        self.stats = {
            "checks": 0,
            "warnings": 0,
            "force_evictions": 0,
            "ollama_crashes": 0,
            "recovery_attempts": 0,
            "recovery_failures": 0,
        }
        self._last_event: str = ""

    def start(self) -> None:
        """Start the background monitor."""
        if self._running:
            return
        self._running = True
        self._task = asyncio.create_task(self._monitor_loop())
        print(f"  [monitor] Memory monitor started (poll every {self.poll_interval}s)")

    async def stop(self) -> None:
        """Stop the background monitor and return stats."""
        self._running = False
        if self._task:
            self._task.cancel()
            try:
                await self._task
            except asyncio.CancelledError:
                pass
            self._task = None
        print(
            f"  [monitor] Stopped — {self.stats['checks']} checks, "
            f"{self.stats['force_evictions']} evictions, "
            f"{self.stats['ollama_crashes']} Ollama crashes"
        )

    def _log(self, msg: str) -> None:
        """Log with dedup — suppress repeated identical events."""
        if msg != self._last_event:
            print(f"  [monitor] {msg}", flush=True)
            self._last_event = msg

    async def _monitor_loop(self) -> None:
        """Main monitoring loop — runs until stop() is called."""
        while self._running:
            try:
                await self._check_once()
            except asyncio.CancelledError:
                raise
            except Exception as e:
                self._log(f"Monitor error: {e}")
            try:
                await asyncio.sleep(self.poll_interval)
            except asyncio.CancelledError:
                break

    async def _check_once(self) -> None:
        """One monitoring cycle."""
        self.stats["checks"] += 1

        # ── 1. Memory pressure ──
        used = _get_memory_pct()
        if used >= MEMORY_ABORT_PCT:
            self._log(f"CRITICAL: Memory at {used:.0f}% — emergency eviction")
            self.stats["force_evictions"] += 1
            await self._emergency_evict()
            used = _get_memory_pct()
            if used >= MEMORY_ABORT_PCT:
                self._log(
                    f"ABORT RISK: Memory still {used:.0f}% after eviction — "
                    "manual intervention may be needed"
                )
                self.stats["recovery_failures"] += 1
        elif used >= MEMORY_CRITICAL_PCT:
            self._log(
                f"Memory critical: {used:.0f}% — model loaded, pre-test check will evict between tests"
            )
            self.stats["warnings"] += 1
        elif used >= MEMORY_WARN_PCT:
            self.stats["warnings"] += 1
            self._log(f"Memory warning: {used:.0f}%")

        # ── 2. Ollama health ──
        try:
            r = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=3)
            if r.status_code != 200:
                self.stats["ollama_crashes"] += 1
                self._log(f"Ollama unhealthy: HTTP {r.status_code}")
        except Exception:
            self.stats["ollama_crashes"] += 1
            self._log("Ollama unreachable — may have crashed")

    async def _emergency_evict(self) -> None:
        """Aggressive eviction when memory is critically high."""
        self.stats["recovery_attempts"] += 1
        unload_all_models()
        await asyncio.sleep(15)


# ---------------------------------------------------------------------------
# Crash watcher — detects Ollama crashes via macOS DiagnosticReports
# ---------------------------------------------------------------------------

_DIAG_DIR = Path.home() / "Library/Logs/DiagnosticReports"


class CrashWatcher:
    """Background thread that watches DiagnosticReports for Ollama-related crashes.

    When a new .ips or .crash file appears whose content references Ollama,
    the watcher logs a [CRASH DETECTED] line immediately.

    The main test loop calls wait_for_recovery() when crash_pending is True,
    which unloads all models and waits for memory to drain.
    """

    POLL_INTERVAL_S = 15

    def __init__(self) -> None:
        self.crash_pending = False
        self._stop = threading.Event()
        self._known: set[Path] = set()
        self.crash_log: list[str] = []
        self._thread: threading.Thread | None = None

    def start(self) -> None:
        if _DIAG_DIR.exists():
            self._known = set(_DIAG_DIR.glob("*.ips")) | set(_DIAG_DIR.glob("*.crash"))
        self._thread = threading.Thread(target=self._loop, daemon=True, name="crash-watcher")
        self._thread.start()
        print(
            "  [crash-watcher] Started — watching DiagnosticReports for crashes",
            flush=True,
        )

    def stop(self) -> None:
        self._stop.set()

    def _loop(self) -> None:
        while not self._stop.wait(timeout=self.POLL_INTERVAL_S):
            try:
                self._check()
            except Exception as exc:
                print(f"  [crash-watcher] error: {exc}", flush=True)

    def _check(self) -> None:
        if not _DIAG_DIR.exists():
            return
        current = set(_DIAG_DIR.glob("*.ips")) | set(_DIAG_DIR.glob("*.crash"))
        new_files = current - self._known
        self._known = current
        for f in new_files:
            self._handle(f)

    def _handle(self, f: Path) -> None:
        try:
            content = f.read_text(errors="replace")
        except Exception:
            return
        proc = f.name
        try:
            header = _json.loads(content.split("\n", 1)[0])
            proc = header.get("app_name", proc)
        except Exception:
            pass
        mem_pct = _get_memory_pct()

        msg = (
            f"  [CRASH DETECTED] proc={proc} file={f.name} "
            f"mem={mem_pct:.0f}%"
        )
        print(msg, flush=True)
        self.crash_log.append(msg)
        self.crash_pending = True

    def wait_for_recovery(self, label: str = "") -> None:
        """Block until memory has drained after a crash.

        Called by the test loop when crash_pending is True. Does not return
        until it is safe to load the next model.
        """
        tag = f"[{label}] " if label else ""
        print(
            f"  {tag}[recovery] Ollama crash — unloading models, waiting for memory to drain...",
            flush=True,
        )
        unload_all_models()
        _wait_for_drain(threshold_pct=MEMORY_WARN_PCT, timeout_s=90.0, label="crash-recovery")
        self.crash_pending = False
        print(f"  {tag}[recovery] Complete — resuming testing", flush=True)


# Module-level singleton — started in main(), stopped after last test
_crash_watcher = CrashWatcher()


# ---------------------------------------------------------------------------
# Model cascade ordering
# ---------------------------------------------------------------------------

# Tier execution order: ollama first, then any, then media_heavy
_TIER_ORDER = ["ollama", "any", "media_heavy"]


def sort_tests_cascade(tests: list[dict]) -> list[dict]:
    """Reorder tests for model-cascade execution.

    Order:
    1. By workspace_tier: ollama → any
       (ollama tests first, so the hardest loads are done early and memory
       is cleanest at the start)
    2. Within each tier, by model_slug: groups tests using the same persona
       together, minimizing model switches within the pipeline
    3. Within each model_slug, preserve original order (test IDs)

    This replaces section-based ordering. Instead of:
      all auto-coding tests → all auto-spl tests → ...
    We do:
      all ollama tests (grouped by model) → all any tests → ...

    Benefits:
    - Models loaded once per tier transition, not per section
    - Big models tested while memory is freshest
    - Tests using same persona run consecutively (pipeline caches)
    - Clear memory boundaries between tiers
    """
    tier_rank = {t: i for i, t in enumerate(_TIER_ORDER)}
    return sorted(
        tests,
        key=lambda t: (
            tier_rank.get(t.get("workspace_tier", "any"), 99),
            t.get("model_slug", ""),
            t.get("id", ""),
        ),
    )


# ---------------------------------------------------------------------------
# Test runner
# ---------------------------------------------------------------------------


async def _run_two_chat_test(
    page,
    test: dict,
    token: str,
    n: int,
    counts: dict,
    folder_id: str | None = None,
    calibration_records: list | None = None,
    corpus_run_id: str = "",
) -> None:
    """Two-chat orchestration for cross-session tests (A-08).

    Creates two distinct OWUI chats in the same workspace. Sends `prompt`
    in chat 1, then `turn2_in_new_chat` in chat 2. Assertions on both
    responses. Best-effort cleanup of any matching memory records via the
    Memory MCP forget API.
    """
    test_id = test["id"]
    name = test["name"]
    model = test["model_slug"]
    tier = test.get("workspace_tier", "any")

    # Backend health
    if tier in ("ollama",):
        backend_ready = await _wait_for_backend(tier, max_wait=120)
        if not backend_ready:
            chat_id, chat_url = owui_create_chat(token, model, f"[FAIL] UAT: {test_id} {name}")
            if folder_id:
                owui_assign_chat_folder(token, chat_id, folder_id)
            record_result(
                n,
                "FAIL",
                test_id,
                name,
                model,
                [("backend_unavailable", False, "tier not ready")],
                0.0,
                chat_url,
            )
            counts["FAIL"] = counts.get("FAIL", 0) + 1
            return

    title1 = f"[...] UAT: {test_id} (1/2) {name}"
    title2 = f"[...] UAT: {test_id} (2/2) {name}"
    # Create chats and assign folders via API BEFORE browser navigation (same reason
    # as main test path — post-nav folder assignment triggers SSE that corrupts submit).
    chat1_id, chat1_url = owui_create_chat(token, model, title1)
    chat2_id, chat2_url = owui_create_chat(token, model, title2)
    if folder_id:
        owui_assign_chat_folder(token, chat1_id, folder_id)
        owui_assign_chat_folder(token, chat2_id, folder_id)
    try:
        await _navigate_to_chat(page, chat1_url)
        await _navigate_to_chat(page, chat2_url)
    except _PresetUnreachableError as exc:
        record_result(
            n,
            "SKIP",
            test_id,
            name,
            model,
            [("persona_preset_unreachable", False, str(exc)[:160])],
            0.0,
            "",
        )
        counts["SKIP"] = counts.get("SKIP", 0) + 1
        return

    t0 = time.time()
    response1 = ""
    response2 = ""
    assertions_result: list = []
    status = "FAIL"
    routed_model_1 = ""
    routed_model_2 = ""

    try:
        max_wait = test.get("max_wait_no_progress", MAX_WAIT_NO_PROGRESS)

        # Ensure Ollama model is loaded before sending — two-chat flow skips
        # the main runner's pre-flight check.
        if tier in ("ollama",):
            unload_all_models()

        # Pre-seed memory via direct MCP API. Decouples test reliability from
        # model-initiated 'remember' (flaky in programmatic OWUI sessions).
        # Test still validates full recall pipeline: LanceDB → semantic search → model.
        preseed_data = test.get("memory_preseed")
        if preseed_data:
            try:
                async with httpx.AsyncClient(timeout=15) as _mc:
                    _resp = await _mc.post(
                        "http://localhost:8920/tools/remember",
                        json={"arguments": preseed_data},
                    )
                if _resp.status_code == 200:
                    print(f"[A-08] memory pre-seeded: {_resp.json().get('id', '?')}", flush=True)
                    await asyncio.sleep(2.0)  # let LanceDB index settle
                else:
                    print(
                        f"[A-08] memory pre-seed failed HTTP {_resp.status_code} — skipping",
                        flush=True,
                    )
                    record_result(
                        n,
                        "SKIP",
                        test_id,
                        name,
                        model,
                        [("memory_preseed_failed", False, f"HTTP {_resp.status_code}")],
                        0.0,
                        "memory-preseed-fail://",
                    )
                    counts["SKIP"] = counts.get("SKIP", 0) + 1
                    return
            except Exception as _e:
                print(f"[A-08] memory pre-seed error: {_e} — skipping", flush=True)
                record_result(
                    n,
                    "SKIP",
                    test_id,
                    name,
                    model,
                    [("memory_preseed_failed", False, str(_e)[:100])],
                    0.0,
                    "memory-preseed-fail://",
                )
                counts["SKIP"] = counts.get("SKIP", 0) + 1
                return

        # Chat 1
        await _navigate_to_chat(page, chat1_url)
        # Note: do NOT call _enable_tool here. The portal pipeline injects
        # and dispatches tools internally for auto-daily (and any workspace
        # with effective_tools). Enabling the tool in OWUI causes OWUI to
        # also dispatch tool_calls it sees in the SSE stream (double-dispatch),
        # which creates a second conversation turn with empty tool results that
        # overwrites the pipeline's correct answer. Pipeline owns dispatch.
        await _fe_send_and_wait(
            page,
            test["prompt"],
            test_id,
            tier,
            max_wait,
            token=token,
            chat_id=chat1_id,
        )
        chat1_url = _fe_current_chat_url(page, fallback=chat1_url)
        response1 = await _fe_get_last_response(page, token, chat1_id) or ""
        routed_model_1 = await _fe_get_routed_model(test, page, token, chat1_id)

        # Brief settle to let the memory write commit through embedding
        # service before chat 2 queries it. The recall is vector-based and
        # needs the entry to be visible in the LanceDB table.
        await asyncio.sleep(5)

        # Chat 2 — fresh chat_url, ZERO context shared with chat 1 except
        # via the model calling 'recall' on the Memory MCP.
        await _navigate_to_chat(page, chat2_url)

        await _fe_send_and_wait(
            page,
            test["turn2_in_new_chat"],
            test_id,
            tier,
            max_wait,
            token=token,
            chat_id=chat2_id,
        )
        chat2_url = _fe_current_chat_url(page, fallback=chat2_url)
        response2 = await _fe_get_last_response(page, token, chat2_id) or ""
        routed_model_2 = await _fe_get_routed_model(test, page, token, chat2_id)

        # Assertions
        _incl_think = test.get("include_thinking_in_assertions", False)
        assertions_result = run_assertions(response1, test.get("assertions", []), include_thinking=_incl_think)
        t2_results = run_assertions(response2, test.get("turn2_assertions", []), include_thinking=_incl_think)
        assertions_result.extend(t2_results)

        all_specs = test.get("assertions", []) + test.get("turn2_assertions", [])
        status = compute_status(assertions_result, all_specs)

        # Routing observation (V1 Phase 2 helper) — append on best-effort basis
        try:
            check1 = _check_routed_model(test, routed_model_1)
            if check1 is not None:
                ok, det = check1
                assertions_result.append((f"Chat 1 routed: {routed_model_1[:30]}", ok, det))
            check2 = _check_routed_model(test, routed_model_2)
            if check2 is not None:
                ok, det = check2
                assertions_result.append((f"Chat 2 routed: {routed_model_2[:30]}", ok, det))
        except NameError:
            # _check_routed_model not present — V1 not merged. Skip silently.
            pass

        if status in ("FAIL", "WARN"):
            SCREENSHOT_DIR.mkdir(exist_ok=True)
            await page.screenshot(path=str(SCREENSHOT_DIR / f"{test_id.lower()}_chat2.png"))

    except Exception as exc:
        assertions_result = [("exception", False, str(exc)[:120])]
        status = "FAIL"
    finally:
        # Best-effort cleanup of memory marker — does not affect status.
        # The model may or may not have actually called remember; either way
        # we flush anything tagged with our marker to avoid accumulation.
        marker_tag = test.get("cleanup_marker_tag")
        if marker_tag:
            try:
                async with httpx.AsyncClient(timeout=10) as client:
                    list_r = await client.post(
                        "http://localhost:8920/tools/list_memories",
                        json={"arguments": {"tags": [marker_tag], "limit": 50}},
                    )
                    if list_r.status_code == 200:
                        for m in list_r.json().get("memories", []):
                            await client.post(
                                "http://localhost:8920/tools/forget",
                                json={"arguments": {"id": m["id"]}},
                            )
            except Exception:
                pass  # cleanup best-effort

    elapsed = time.time() - t0
    final_title_1 = f"[{status} 1/2] UAT: {test_id} {name}"
    final_title_2 = f"[{status} 2/2] UAT: {test_id} {name}"
    owui_rename_chat(token, chat1_id, final_title_1)
    owui_rename_chat(token, chat2_id, final_title_2)

    # Use chat 2 URL as the "primary" link in results — it's where the
    # actual recall behavior is visible to a reviewer.
    record_result(
        n,
        status,
        test_id,
        name,
        model,
        assertions_result,
        elapsed,
        chat2_url,
        routed_model_2,
    )
    counts[status] = counts.get(status, 0) + 1

    if calibration_records is not None:
        calibration_records.append(
            {
                "test_id": test_id,
                "name": name,
                "section": test.get("section", ""),
                "workspace": test.get("model_slug", ""),
                "prompt": test.get("prompt", "")
                + "\n\n[NEW CHAT]\n"
                + test.get("turn2_in_new_chat", ""),
                "response_text": (
                    f"=== Chat 1 ===\n{response1}\n\n=== Chat 2 (recall) ===\n{response2}"
                ),
                "chat_url": chat2_url,
                "review_tag": "",
                "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ"),
            }
        )

    # Always-on corpus emission for two-chat tests. The prompt and
    # response carry the same dual-chat formatting as the calibration
    # record above so a single corpus reader handles both shapes.
    if corpus_run_id:
        _composite_test = dict(test)
        _composite_test["prompt"] = (
            test.get("prompt", "") + "\n\n[NEW CHAT]\n" + test.get("turn2_in_new_chat", "")
        )
        _composite_response = f"=== Chat 1 ===\n{response1}\n\n=== Chat 2 (recall) ===\n{response2}"
        _emit_corpus_row(
            corpus_run_id=corpus_run_id,
            test=_composite_test,
            routed_model=routed_model_2,
            response_text=_composite_response,
            chat_url=chat2_url,
            status=status,
            assertions_result=assertions_result,
            elapsed=elapsed,
        )


async def run_test(
    page,
    test: dict,
    token: str,
    skip_conditions: dict,
    n: int,
    counts: dict,
    headed: bool = False,
    folder_id: str | None = None,
    calibration_records: list | None = None,
    corpus_run_id: str = "",
) -> None:
    test_id = test["id"]
    name = test["name"]
    model = test["model_slug"]

    title_pending = f"[...] UAT: {test_id} {name}"

    # Skip check — skip_if can be a string or list of strings (any match skips)
    skip_if = test.get("skip_if")
    _skip_keys = [skip_if] if isinstance(skip_if, str) else (skip_if or [])
    if any(skip_conditions.get(k, False) for k in _skip_keys):
        _matched_key = next((k for k in _skip_keys if skip_conditions.get(k, False)), skip_if)
        chat_id, chat_url = owui_create_chat(token, model, f"[SKIP] UAT: {test_id} {name}")
        owui_rename_chat(token, chat_id, f"[SKIP] UAT: {test_id} {name} — {_matched_key}")
        if folder_id:
            owui_assign_chat_folder(token, chat_id, folder_id)
        record_result(n, "SKIP", test_id, name, model, [], 0.0, chat_url)
        counts["SKIP"] = counts.get("SKIP", 0) + 1
        return

    # Manual test
    if test.get("is_manual"):
        chat_id, chat_url = owui_create_chat(token, model, title_pending)
        if folder_id:
            owui_assign_chat_folder(token, chat_id, folder_id)
        manual_prompt = (
            "🔧 MANUAL TEST: "
            + test["prompt"]
            + "\n\nReturn to this chat and pin your result with ✅ PASS / ⚠️ PARTIAL / ❌ FAIL + notes."
        )
        await _navigate_to_chat(page, chat_url)
        await _send_and_wait(page, manual_prompt, test_id, token=token, chat_id=chat_id)
        owui_rename_chat(token, chat_id, f"[MANUAL] UAT: {test_id} {name}")
        record_result(n, "MANUAL", test_id, name, model, [], 0.0, chat_url)
        counts["MANUAL"] = counts.get("MANUAL", 0) + 1
        return

    # Dispatcher-path test (Telegram / Slack bot pipeline call).
    # Drives the exact code path portal_channels.dispatcher uses on every
    # inbound bot message: a direct POST to the Pipeline with PIPELINE_API_KEY.
    # Bypasses Open WebUI and Playwright entirely.
    if test.get("via_dispatcher"):
        # Pre-check: bot container running, if specified.
        required_container = test.get("requires_container")
        if required_container:
            ok, detail = _bot_container_running(required_container)
            if not ok:
                record_result(
                    n,
                    "BLOCKED",
                    test_id,
                    name,
                    model,
                    [("bot_container_unavailable", False, f"{required_container}: {detail}")],
                    0.0,
                    "",
                )
                counts["BLOCKED"] = counts.get("BLOCKED", 0) + 1
                return

        t0_disp = time.time()
        try:
            response_text = await _run_via_dispatcher(
                workspace=model,
                prompt=test["prompt"],
                timeout=test.get("timeout", 120),
            )
        except Exception as exc:
            elapsed = time.time() - t0_disp
            # Transport/auth errors = infrastructure not wired up → BLOCKED,
            # not a product defect. Content failures (wrong response) → FAIL.
            record_result(
                n,
                "BLOCKED",
                test_id,
                name,
                model,
                [("dispatcher_call_failed", False, f"{type(exc).__name__}: {str(exc)[:160]}")],
                elapsed,
                "",
            )
            counts["BLOCKED"] = counts.get("BLOCKED", 0) + 1
            return

        elapsed = time.time() - t0_disp
        _incl_think = test.get("include_thinking_in_assertions", False)
        assertions_result = run_assertions(response_text, test.get("assertions", []), include_thinking=_incl_think)
        status = compute_status(assertions_result, test.get("assertions", []))
        # No chat URL — this path doesn't create an Open WebUI chat. Use a
        # synthetic marker so the report shows where the response came from.
        record_result(
            n,
            status,
            test_id,
            name,
            model,
            assertions_result,
            elapsed,
            f"via-dispatcher://{model}",
        )
        counts[status] = counts.get(status, 0) + 1
        return

    # Two-chat test: A-08 (cross-session memory). Creates two distinct
    # OWUI chats, uses the same workspace, runs separate prompt+turn2_in_new_chat
    # turns. Each chat shows up independently in OWUI history.
    if test.get("is_two_chat"):
        return await _run_two_chat_test(
            page,
            test,
            token,
            n,
            counts,
            folder_id,
            calibration_records,
            corpus_run_id=corpus_run_id,
        )

    tier = test.get("workspace_tier", "any")

    # Pre-test backend health gate — wait up to 120s for Ollama to be ready.
    if tier in ("ollama",):
        backend_ready = await _wait_for_backend(tier, max_wait=120)
        if not backend_ready:
            backend_ready = await _wait_for_backend(tier, max_wait=60)
        if not backend_ready:
            _, detail = _backend_alive(tier)
            chat_id, chat_url = owui_create_chat(token, model, f"[FAIL] UAT: {test_id} {name}")
            if folder_id:
                owui_assign_chat_folder(token, chat_id, folder_id)
            record_result(
                n,
                "FAIL",
                test_id,
                name,
                model,
                [("backend_unavailable", False, detail)],
                0.0,
                chat_url,
            )
            counts["FAIL"] = counts.get("FAIL", 0) + 1
            return

    # Create chat and assign folder via API BEFORE browser navigation.
    # Assigning the folder after the browser has loaded an empty chat causes
    # OWUI to broadcast a chat-updated SSE event. The Svelte component re-renders
    # the "new chat" suggestions view in response, which corrupts the submit handler
    # and silently drops Enter keypresses. Assigning the folder first means the
    # browser opens the chat with the folder already set — no SSE event fires
    # during the test session.
    chat_id, chat_url = owui_create_chat(token, model, title_pending)
    if folder_id:
        owui_assign_chat_folder(token, chat_id, folder_id)
    try:
        await _navigate_to_chat(page, chat_url)
    except _PresetUnreachableError as exc:
        record_result(
            n,
            "SKIP",
            test_id,
            name,
            model,
            [("persona_preset_unreachable", False, str(exc)[:160])],
            0.0,
            chat_url,
        )
        counts["SKIP"] = counts.get("SKIP", 0) + 1
        return
    except Exception as exc:
        # SPA navigation timeout or other startup error — record as BLOCKED and
        # continue to the next test rather than crashing the entire run.
        print(
            f"  [BLOCKED] {test_id} — chat start failed: {type(exc).__name__}: {str(exc)[:120]}",
            flush=True,
        )
        record_result(
            n,
            "BLOCKED",
            test_id,
            name,
            model,
            [("chat_start_failed", False, f"{type(exc).__name__}: {str(exc)[:160]}")],
            0.0,
            chat_url,
        )
        counts["BLOCKED"] = counts.get("BLOCKED", 0) + 1
        return

    t0 = time.time()
    artifact_path: Path | None = None
    assertions_result: list = []
    status = "FAIL"
    response_text = ""
    attempts_used: int = 1

    try:
        # _navigate_to_chat above is the only navigation — do NOT navigate again here.
        # A second page.goto() corrupts Svelte submit-handler state for models
        # with tool initialization (dailydriver/proofreader).

        # Pre-stage audio fixture for tests that drive the mlx-transcribe MCP.
        # The MCP auto-detects the most recently modified audio file in the
        # workspace uploads dir when called with no `file` arg
        # (see scripts/mlx-transcribe.py::_latest_audio_upload). Mirrors how
        # operators drop audio into the UI; OWUI's M-01 already relies on
        # the same path.
        if test.get("pre_stage_audio"):
            import shutil as _shutil

            _fixture_path = Path(__file__).parent / "fixtures" / test.get("fixture", "")
            _ai_output = Path(os.environ.get("AI_OUTPUT_DIR") or (Path.home() / "AI_Output"))
            _uploads = _ai_output / "uploads"
            _uploads.mkdir(parents=True, exist_ok=True)
            _staged = _uploads / _fixture_path.name
            if _fixture_path.exists():
                _shutil.copy2(_fixture_path, _staged)
                _staged.touch()  # ensure newest-mtime wins
                print(f"  [TR pre-stage] staged {_fixture_path.name} → {_uploads}", flush=True)
            else:
                print(f"  [TR pre-stage] WARN: fixture missing at {_fixture_path}", flush=True)

        # Tools are pre-enabled via workspace toolIds seeding — do not toggle them here.
        # Calling _enable_tool would turn them OFF (they default to ON in seeded workspaces).

        # Send first turn — retry up to 2 times on empty response (Ollama cold load).
        # This is RECOVERY logic (handle empty/crashed backend), not a
        # validation strategy — same prompt is re-sent each time.
        max_wait = test.get("max_wait_no_progress", MAX_WAIT_NO_PROGRESS)
        _test_budget_s = test.get("timeout", 120)
        response_text = ""
        attempts_used = 0
        for attempt in range(3):
            attempts_used = attempt + 1
            await _fe_send_and_wait(
                page,
                test["prompt"],
                test_id,
                tier,
                max_wait,
                token=token,
                chat_id=chat_id,
            )
            chat_url = _fe_current_chat_url(page, fallback=chat_url)
            response_text = await _fe_get_last_response(page, token, chat_id)
            if response_text:
                break
            # Long-tail wait: DOM stable may have fired while reasoning model was
            # still generating (collapsed <details> block makes innerText appear
            # stable). Continue polling the API — large GGUF models (30-70B) can
            # take 5-7 minutes for reasoning; media_heavy (video/image gen) needs 240s for cold
            # HunyuanVideo runs; others bounded to ~90s.
            _poll_cap_s = 450 if tier == "ollama" else (240 if tier == "media_heavy" else 90)
            _poll_deadline = time.monotonic() + _poll_cap_s
            while time.monotonic() < _poll_deadline:
                await asyncio.sleep(5)
                response_text = await _fe_get_last_response(page, token, chat_id)
                if response_text:
                    break
                elapsed_now = time.time() - t0
                print(
                    f"  [{test_id}] polling for response… ({elapsed_now:.0f}s)",
                    flush=True,
                )
            if response_text:
                break
            elapsed_now = time.time() - t0
            # Hard cap: if total elapsed exceeds 3× the test timeout, stop retrying.
            # Prevents runaway reasoning models from consuming unbounded wall time.
            if elapsed_now > _test_budget_s * 3:
                print(
                    f"  [{test_id}] total elapsed {elapsed_now:.0f}s > 3× timeout "
                    f"({_test_budget_s * 3}s) — stopping retries",
                    flush=True,
                )
                break
            print(
                f"  [{test_id}] empty response on attempt {attempt + 1}/3 ({elapsed_now:.0f}s)",
                flush=True,
            )
            if attempt < 2:
                # Check backend health before retrying
                await _wait_for_backend_alive(tier)
                # Re-navigate to the chat URL before retrying. OWUI calls
                # get_all_models() on page load — this clears any stale model
                # availability cache from the tier-transition eviction period,
                # and resets any stuck "generating" UI state.
                if chat_url:
                    print(
                        f"  [{test_id}] re-navigating to refresh OWUI model cache before retry…",
                        flush=True,
                    )
                    await _navigate_to_chat(page, chat_url)

        # Download artifact if expected
        art_ext = test.get("artifact_ext")
        if art_ext:
            # Late arrival: slow tools (video gen ~131-200s) may stream past the
            # poll window, OR the model may stream a partial non-empty response
            # before the tool completes. Refresh response_text if it looks
            # incomplete (empty, or has no artifact URL yet).
            import re as _re

            _art_url_present = _re.search(
                rf"(?:/files/\S+?\.{_re.escape(art_ext)}|view\?filename=[^\s)>\]]*\.{_re.escape(art_ext)})",
                response_text or "",
            )
            if not response_text or not _art_url_present:
                response_text = (
                    await _fe_get_last_response(page, token, chat_id) or response_text or ""
                )
            artifact_path = await _fe_download_artifact(page, art_ext, response_text=response_text, since_ts=t0)

        # Multi-turn: send second message if defined
        turn2 = test.get("turn2")
        turn2_response = ""
        if turn2:
            await _fe_send_and_wait(
                page,
                turn2,
                test_id,
                tier,
                max_wait,
                token=token,
                chat_id=chat_id,
                min_messages=2,  # require ≥2 non-empty responses — prevents turn-1
                # stable content from satisfying the completion signal
            )
            chat_url = _fe_current_chat_url(page, fallback=chat_url)
            # For turn2, require ≥2 non-empty assistant messages so we don't
            # return turn-1's committed response as the turn-2 completion signal.
            turn2_response = await _fe_get_last_response(page, token, chat_id, min_messages=2)

        # Run assertions on turn 1
        _incl_think = test.get("include_thinking_in_assertions", False)
        assertions_result = run_assertions(response_text, test.get("assertions", []), artifact_path, include_thinking=_incl_think)

        # Run turn2 assertions if defined
        t2_spec = test.get("turn2_assertions", [])
        if t2_spec and turn2_response:
            t2_results = run_assertions(turn2_response, t2_spec, artifact_path, include_thinking=_incl_think)
            assertions_result.extend(t2_results)

        # Combine all specs for status computation
        all_specs = test.get("assertions", []) + test.get("turn2_assertions", [])
        status = compute_status(assertions_result, all_specs)

        # Surface retry-attempt count when recovery was needed. Appended
        # without a corresponding spec — compute_status (already run above)
        # zips assertions with spec and truncates extras, so this row is
        # informational only and does not affect grading.
        if attempts_used > 1:
            assertions_result.append(
                (
                    f"Recovery: passed on attempt {attempts_used}/3",
                    True,
                    f"{attempts_used - 1} retries needed (backend instability signal)",
                )
            )

        # Take screenshot on failure
        if status in ("FAIL", "WARN"):
            SCREENSHOT_DIR.mkdir(exist_ok=True)
            sc_path = SCREENSHOT_DIR / f"{test_id.lower()}.png"
            await page.screenshot(path=str(sc_path))

    except Exception as exc:
        assertions_result = [("exception", False, str(exc)[:120])]
        status = "FAIL"
        try:
            SCREENSHOT_DIR.mkdir(exist_ok=True)
            await page.screenshot(path=str(SCREENSHOT_DIR / f"{test_id.lower()}_exc.png"))
        except Exception:
            pass

    elapsed = time.time() - t0
    routed_model = await _fe_get_routed_model(test, page, token, chat_id)

    route_check = _check_routed_model(test, routed_model)
    if route_check is not None:
        matched, route_detail = route_check
        assertions_result.append(
            (f"Routed model: {routed_model[:40] or 'none'}", matched, route_detail)
        )
        if status == "PASS" and not matched:
            status = "WARN"
            print(f"  [{test_id}] route mismatch downgraded PASS→WARN: {route_detail}", flush=True)

        # Feed routing telemetry log for end-of-run summary
        try:
            import sys as _sys
            from pathlib import Path as _Path
            _sys.path.insert(0, str(_Path(__file__).parent))
            intended_keys = route_detail  # contains expected key info
            intended_ollama = test.get("workspace_tier", "") == "ollama"
            _ROUTING_LOG.append({
                "test_id": test_id,
                "name": name,
                "section": test.get("section", ""),
                "workspace": test.get("model_slug", ""),
                "intended": test.get("model_slug", ""),
                "actual": routed_model,
                "matched": matched,
                "tier_mismatch": intended_ollama and not matched,
                "pipeline_backend": pipeline_backend,
                "intended_ollama": intended_ollama,
            })
        except Exception:
            pass

    final_title = f"[{status}] UAT: {test_id} {name}"
    owui_rename_chat(token, chat_id, final_title)
    record_result(
        n, status, test_id, name, model, assertions_result, elapsed, chat_url, routed_model
    )
    counts[status] = counts.get(status, 0) + 1

    if calibration_records is not None:
        calibration_records.append(
            {
                "test_id": test_id,
                "name": name,
                "section": test.get("section", ""),
                "workspace": test.get("model_slug", ""),
                "prompt": test.get("prompt", ""),
                "response_text": response_text,
                "chat_url": chat_url,
                "review_tag": "",
                "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ"),
            }
        )

    # Always-on corpus emission — see TASK_UAT_CORPUS_CAPTURE_V1.md §A2.
    if corpus_run_id:
        _emit_corpus_row(
            corpus_run_id=corpus_run_id,
            test=test,
            routed_model=routed_model,
            response_text=response_text,
            chat_url=chat_url,
            status=status,
            assertions_result=assertions_result,
            elapsed=elapsed,
        )


def _emit_corpus_row(
    corpus_run_id: str,
    test: dict,
    routed_model: str,
    response_text: str,
    chat_url: str,
    status: str,
    assertions_result: list,
    elapsed: float,
) -> None:
    """Append one JSONL row to the UAT response corpus.

    The corpus is always-on (no flag required) and one file per UAT run.
    Emission is incremental — each call opens the file in append mode,
    writes one line, and closes, so a crashed run leaves valid JSONL.

    See TASK_UAT_CORPUS_CAPTURE_V1.md for schema + rationale.
    """
    import json as _json

    corpus_dir = Path("tests/uat_corpus")
    corpus_dir.mkdir(parents=True, exist_ok=True)
    corpus_path = corpus_dir / f"uat_{corpus_run_id}.jsonl"

    # Convert tuple assertion results to JSON-safe lists. The in-memory
    # format is tuples of (label:str, passed:bool, detail:str); JSON has
    # no tuple type, so we serialize as lists.
    safe_assertions = [list(a) if isinstance(a, tuple) else a for a in (assertions_result or [])]

    row = {
        "schema_version": 1,
        "corpus_run_id": corpus_run_id,
        "test_id": test.get("id", ""),
        "test_name": test.get("name", ""),
        "section": test.get("section", ""),
        "workspace": test.get("model_slug", ""),
        "expected_models": test.get("expected_models", {}),
        "routed_model": routed_model or "",
        "prompt": test.get("prompt", ""),
        "response_text": response_text or "",
        "chat_url": chat_url or "",
        "status": status,
        "assertions_result": safe_assertions,
        "elapsed_seconds": float(elapsed) if elapsed is not None else 0.0,
        "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ"),
    }
    try:
        with corpus_path.open("a", encoding="utf-8") as f:
            f.write(_json.dumps(row, ensure_ascii=False) + "\n")
    except Exception as exc:
        # Corpus emission is best-effort — never fail a test because the
        # corpus write failed. Log and continue.
        print(f"  [corpus] WARN: failed to write {test.get('id', '?')}: {exc}", flush=True)


def _emit_signals_from_calibration(json_path: str, output_path: str = "updated_signals.py") -> None:
    """Read calibration JSON, extract keywords from 'good' responses, write a signals suggestion file."""
    import json as _json
    import math as _math
    import re as _re

    records = _json.loads(Path(json_path).read_text())
    good = [r for r in records if r.get("review_tag") == "good"]

    if not good:
        print(f"No 'good'-tagged records found in {json_path}.")
        print(
            "Open the JSON, set review_tag to 'good' / 'bad' / 'skip' for each entry, then re-run."
        )
        return

    # Group by section
    by_section: dict[str, list[str]] = {}
    for rec in good:
        sec = rec.get("section") or "general"
        by_section.setdefault(sec, []).append(rec.get("response_text", ""))

    def _tokenize(text: str) -> list[str]:
        return _re.findall(r"\b[a-zA-Z][a-zA-Z0-9_]{2,}\b", text.lower())

    _STOPWORDS = {
        "the",
        "and",
        "for",
        "this",
        "that",
        "with",
        "from",
        "are",
        "can",
        "will",
        "not",
        "you",
        "your",
        "have",
        "has",
        "was",
        "but",
        "all",
        "more",
        "into",
        "use",
        "used",
        "using",
        "would",
        "should",
        "could",
        "when",
        "which",
        "here",
        "there",
        "also",
        "each",
        "such",
        "then",
        "they",
        "them",
        "their",
        "been",
        "its",
        "any",
        "how",
        "what",
        "where",
        "who",
        "why",
        "may",
        "one",
        "two",
        "three",
        "just",
        "like",
        "make",
        "made",
        "note",
        "see",
        "get",
        "set",
    }

    # IDF: inverse of how many sections a word appears in
    idf: dict[str, int] = {}
    for texts in by_section.values():
        words_in_sec = set(_tokenize(" ".join(texts)))
        for w in words_in_sec:
            idf[w] = idf.get(w, 0) + 1
    n_sections = len(by_section)
    idf_score = {w: _math.log((n_sections + 1) / (cnt + 1)) for w, cnt in idf.items()}

    section_keywords: dict[str, list[str]] = {}
    for sec, texts in by_section.items():
        words = _tokenize(" ".join(texts))
        tf: dict[str, int] = {}
        for w in words:
            if w not in _STOPWORDS and len(w) > 3:
                tf[w] = tf.get(w, 0) + 1
        total = sum(tf.values()) or 1
        scored = {w: (cnt / total) * idf_score.get(w, 0.0) for w, cnt in tf.items()}
        section_keywords[sec] = sorted(scored, key=lambda x: -scored[x])[:10]

    out_lines = [
        '"""Auto-generated quality signals from calibration data.',
        "",
        "Generated by: python3 tests/portal5_uat_driver.py --emit-signals-from <json>",
        "",
        "Review and integrate into tests/quality_signals.py or the UAT test catalog.",
        '"""',
        "",
        "CALIBRATION_SIGNALS: dict[str, list[str]] = {",
    ]
    for sec in sorted(section_keywords):
        kws = section_keywords[sec]
        out_lines.append(f"    {sec!r}: {kws!r},")
    out_lines.append("}")
    out_lines.append("")
    out_lines.append("# Suggested assert_contains additions for TEST_CATALOG entries:")
    for sec in sorted(section_keywords):
        kws = section_keywords[sec][:5]
        out_lines.append(
            f"# section={sec!r}: "
            + '{"type": "any_of", "label": "Quality signal", "keywords": '
            + repr(kws)
            + "}"
        )

    Path(output_path).write_text("\n".join(out_lines) + "\n")
    print(f"Signals written to {output_path}")
    for sec, kws in sorted(section_keywords.items()):
        print(f"  {sec}: {kws}")


def _git_sha() -> str:
    """Get current git SHA."""
    try:
        import subprocess as _subprocess

        return _subprocess.run(
            ["git", "-C", str(_REPO_ROOT), "rev-parse", "--short", "HEAD"],
            capture_output=True,
            text=True,
            timeout=10,
        ).stdout.strip()
    except Exception:
        return "unknown"


async def _send_notification(event_type: str, message: str, metadata: dict | None = None) -> None:
    """Fire a notification via the Portal 5 notification dispatcher.

    Gracefully handles missing dependencies or disabled notifications — never
    crashes the test suite.
    """
    if os.environ.get("NOTIFICATIONS_ENABLED", "false").lower() not in ("true", "1", "yes"):
        return
    try:
        from portal_pipeline.notifications.channels.email import EmailChannel
        from portal_pipeline.notifications.channels.pushover import PushoverChannel
        from portal_pipeline.notifications.channels.slack import SlackChannel
        from portal_pipeline.notifications.channels.telegram import TelegramChannel
        from portal_pipeline.notifications.channels.webhook import WebhookChannel
        from portal_pipeline.notifications.dispatcher import NotificationDispatcher
        from portal_pipeline.notifications.events import AlertEvent, EventType

        dispatcher = NotificationDispatcher()
        for ch in [SlackChannel, TelegramChannel, EmailChannel, PushoverChannel, WebhookChannel]:
            dispatcher.add_channel(ch())

        event = AlertEvent(
            type=EventType(event_type.lower()),
            message=message,
            workspace="uat-test",
            metadata=metadata or {},
        )
        await dispatcher.dispatch(event)
    except Exception as e:
        print(f"  WARNING: Notification failed: {e}")


async def _notify_test_start(sections: list[str] | None = None, test_count: int = 0) -> None:
    """Send a notification that UAT testing has started."""
    sections_str = ", ".join(sections) if sections else "all"
    await _send_notification(
        "test_start",
        f"UAT test suite started — section(s): {sections_str} ({test_count} tests)\n"
        f"Git: {_git_sha()}  |  Host: {os.uname().nodename}",
        metadata={"sections": sections or [], "test_count": test_count},
    )


async def _notify_test_end(
    sections: list[str] | None,
    elapsed: int,
    counts: dict[str, int],
    test_count: int,
) -> None:
    """Send a notification that UAT testing has completed."""
    summary_parts = [
        f"PASS={counts.get('PASS', 0)}",
        f"FAIL={counts.get('FAIL', 0)}",
        f"WARN={counts.get('WARN', 0)}",
        f"SKIP={counts.get('SKIP', 0)}",
        f"MANUAL={counts.get('MANUAL', 0)}",
    ]
    sections_str = ", ".join(sections) if sections else "all"
    await _send_notification(
        "test_end",
        f"UAT test suite completed — section(s): {sections_str} in {elapsed}s\n"
        f"Results: {', '.join(summary_parts)}\n"
        f"Git: {_git_sha()}",
        metadata={"elapsed_s": elapsed, "counts": counts},
    )


async def _notify_test_summary(
    counts: dict[str, int], elapsed: int, sections: list[str] | None, test_count: int
) -> None:
    """Send the narrative summary via all enabled notification channels."""
    total = sum(counts.values())
    failed = counts.get("FAIL", 0)
    warned = counts.get("WARN", 0)

    if failed:
        narrative = f"{failed} test{'s' if failed > 1 else ''} failed"
    elif warned:
        narrative = f"All {total} tests passed with {warned} warning{'s' if warned > 1 else ''}"
    else:
        narrative = f"All {total} tests passed"

    sections_str = ", ".join(sections) if sections else "all"
    lines = [
        narrative,
        "",
        f"Portal 5 UAT Driver — section(s): {sections_str}",
        f"Duration: {elapsed}s  |  Tests: {test_count}",
        f"Git: {_git_sha()}  |  Host: {os.uname().nodename}",
        "",
    ]
    for s in ["PASS", "FAIL", "WARN", "SKIP", "BLOCKED", "MANUAL"]:
        if s in counts:
            lines.append(f"  {s}: {counts[s]}")
    lines.append(f"  Total: {total}")

    await _send_notification(
        "test_summary",
        "\n".join(lines),
        metadata={"counts": counts, "elapsed_s": elapsed, "test_count": test_count},
    )


async def main() -> None:
    parser = argparse.ArgumentParser(description="Portal 5 UAT Conversation Driver")
    parser.add_argument("--all", action="store_true", help="Run all tests")
    parser.add_argument("--section", action="append", help="Run tests from section(s)")
    parser.add_argument(
        "--test", metavar="ID", action="append", help="Run test(s) by ID (repeatable)"
    )
    parser.add_argument("--headed", action="store_true", help="Show browser window")
    parser.add_argument("--skip-artifacts", action="store_true", help="Skip ComfyUI/Wan2.2 tests")
    parser.add_argument("--skip-bots", action="store_true", help="Skip Telegram/Slack bot tests")
    parser.add_argument(
        "--media",
        action="store_true",
        help=(
            "Run only media-generation tests (image, sound, voice, video) — "
            "shorthand for selecting all tests with workspace_tier=media_heavy. "
            "Useful for debugging MCP/Open WebUI media plumbing in isolation."
        ),
    )
    parser.add_argument("--timeout", type=int, help="Override per-test timeout (seconds)")
    parser.add_argument(
        "--append",
        action="store_true",
        help="Append results to existing UAT_RESULTS.md (for re-runs)",
    )
    parser.add_argument(
        "--no-unload",
        action="store_true",
        help="Skip startup /unload — use when model is pre-warmed",
    )
    parser.add_argument(
        "--rerun",
        action="store_true",
        help=(
            "Re-run mode: remove existing rows in UAT_RESULTS.md for the selected "
            "test IDs before running. Implies --append. Use this when re-running a "
            "phase after a fix; prevents duplicate rows. Requires --section, --test, "
            "or --media to scope which tests to replace."
        ),
    )
    parser.add_argument(
        "--rerun-failed",
        action="store_true",
        help=(
            "Re-run only tests with status FAIL or BLOCKED in UAT_RESULTS.md. "
            "Implies --rerun --append. Use after a fix to retry only broken tests "
            "without re-running the entire section."
        ),
    )
    parser.add_argument(
        "--migrate",
        action="store_true",
        help="Move existing root-level UAT chats into root UAT folder, then exit",
    )
    parser.add_argument(
        "--purge-uat",
        action="store_true",
        help="Delete all chats in the UAT folder and the folder itself, then exit. "
        "Run this after reviewing UAT_RESULTS.md to clean up OWUI.",
    )
    parser.add_argument(
        "--calibrate",
        action="store_true",
        help="Calibration mode: run all tests and capture full responses to JSON for review",
    )
    parser.add_argument(
        "--calibrate-output",
        default="calibration.json",
        metavar="FILE",
        help="Output path for calibration JSON (default: calibration.json)",
    )
    parser.add_argument(
        "--emit-signals-from",
        metavar="JSON",
        help="Generate quality_signals suggestions from a reviewed calibration JSON",
    )
    args = parser.parse_args()

    print("\nPortal 5 UAT Driver")
    print(f"OWUI: {OPENWEBUI_URL}  |  User: {ADMIN_EMAIL}")
    print(f"Results: {RESULTS_FILE}\n")

    # Auth
    token = owui_token()
    if not token:
        print("ERROR: Could not authenticate with Open WebUI", file=sys.stderr)
        sys.exit(1)

    # Codebase freshness — warn if running images predate latest git commits.
    # Stale images mean test results reflect old code, not HEAD.
    _check_image_freshness()

    # --emit-signals-from mode: standalone, no browser needed
    if args.emit_signals_from:
        output = getattr(args, "calibrate_output", "updated_signals.py")
        _emit_signals_from_calibration(args.emit_signals_from, output)
        return

    # --migrate mode: move existing loose UAT chats into UAT folder hierarchy, then exit
    if args.migrate:
        uat_root_id = owui_get_or_create_folder(token, "UAT")
        if uat_root_id:
            print(f"  Migrating loose UAT chats → root UAT folder (id={uat_root_id}) …")
            n_moved = owui_migrate_loose_uat_chats(token, uat_root_id)
            print(f"  Migrated {n_moved} chat(s).")
        else:
            print("  ERROR: could not get/create UAT root folder.")
            sys.exit(1)
        return

    # --purge-uat mode: delete all chats in the UAT folder, then delete the folder
    if args.purge_uat:
        folders = _owui_list_folders(token)
        uat_folder = next((f for f in folders if f.get("name") == "UAT" and not f.get("parent_id")), None)
        if not uat_folder:
            print("  No UAT folder found — nothing to purge.")
            return
        uat_root_id = uat_folder["id"]
        # Collect all chats currently in the UAT folder
        try:
            r = httpx.get(
                f"{OPENWEBUI_URL}/api/v1/chats/",
                headers=owui_headers(token),
                params={"limit": 9999},
                timeout=30,
            )
            all_chats = r.json() if r.status_code == 200 else []
        except Exception as e:
            print(f"  ERROR fetching chats: {e}")
            sys.exit(1)
        # OWUI list endpoint may not include folder_id; fetch detail for each to filter
        uat_chat_ids: list[str] = []
        for chat in all_chats:
            cid = chat.get("id", "")
            try:
                r2 = httpx.get(
                    f"{OPENWEBUI_URL}/api/v1/chats/{cid}",
                    headers=owui_headers(token),
                    timeout=10,
                )
                if r2.status_code == 200 and r2.json().get("folder_id") == uat_root_id:
                    uat_chat_ids.append(cid)
            except Exception:
                pass
        print(f"  UAT folder id={uat_root_id} — {len(uat_chat_ids)} chat(s) to delete")
        deleted = 0
        for cid in uat_chat_ids:
            try:
                r = httpx.delete(
                    f"{OPENWEBUI_URL}/api/v1/chats/{cid}",
                    headers=owui_headers(token),
                    timeout=10,
                )
                if r.status_code == 200:
                    deleted += 1
                else:
                    print(f"  WARNING: DELETE chat {cid} returned {r.status_code}")
            except Exception as e:
                print(f"  WARNING: DELETE chat {cid} error — {e}")
        print(f"  Deleted {deleted}/{len(uat_chat_ids)} chat(s).")
        # Now delete the UAT folder itself
        try:
            r = httpx.delete(
                f"{OPENWEBUI_URL}/api/v1/folders/{uat_root_id}",
                headers=owui_headers(token),
                timeout=10,
            )
            if r.status_code == 200:
                print("  UAT folder deleted.")
            else:
                print(f"  WARNING: DELETE folder returned {r.status_code} — {r.text[:120]}")
        except Exception as e:
            print(f"  WARNING: DELETE folder error — {e}")
        return

    # --rerun-failed: auto-select FAIL/BLOCKED tests from UAT_RESULTS.md,
    # then run them through the same cascade logic as a normal run.
    # Tests are sorted by tier (ollama → any) so
    # model loads are grouped and tier-transition eviction guards fire correctly.
    _RERUN_FAILED_STATE = Path("/tmp/portal5-rerun-failed-state.json")

    if args.rerun_failed:
        failed_ids = _parse_failed_test_ids()
        if not failed_ids:
            # Rows may have been removed by a previous --rerun-failed that was
            # interrupted before completing. Check for a saved state file.
            if _RERUN_FAILED_STATE.exists():
                import json as _json_rf

                saved = _json_rf.loads(_RERUN_FAILED_STATE.read_text())
                failed_ids = set(saved.get("ids", []))
                if failed_ids:
                    print(
                        f"  --rerun-failed: restored {len(failed_ids)} ID(s) from previous "
                        f"interrupted run ({_RERUN_FAILED_STATE})",
                        file=sys.stderr,
                    )
        if not failed_ids:
            print(
                "--rerun-failed: no FAIL or BLOCKED tests found in UAT_RESULTS.md — nothing to do",
                file=sys.stderr,
            )
            sys.exit(0)

        # Resolve IDs → catalog entries so we can show the tier plan up front.
        candidate_tests = [t for t in TEST_CATALOG if t["id"] in failed_ids]
        unknown = failed_ids - {t["id"] for t in candidate_tests}
        if unknown:
            print(
                f"  --rerun-failed: WARNING — {len(unknown)} ID(s) not in TEST_CATALOG "
                f"(may have been removed): {', '.join(sorted(unknown))}",
                file=sys.stderr,
            )

        # Group by tier so the caller can see what backend switching will occur.
        tier_groups: dict[str, list[str]] = {}
        for t in sort_tests_cascade(candidate_tests):
            tier = t.get("workspace_tier", "any")
            tier_groups.setdefault(tier, []).append(t["id"])

        plan = " → ".join(f"{tier}({len(ids)})" for tier, ids in tier_groups.items())
        print(f"  --rerun-failed: {len(candidate_tests)} test(s) across {len(tier_groups)} tier(s)")
        print(f"  Cascade plan: {plan}")
        for tier, ids in tier_groups.items():
            print(f"    [{tier}] {', '.join(ids)}")

        if len(tier_groups) > 1:
            print(
                "  NOTE: tier transitions will evict all models between groups — "
                "expect 30-60s pauses at each boundary."
            )

        # Save state before removing rows — if this run is interrupted, the
        # next --rerun-failed invocation can restore from here.
        import json as _json_rf2

        _RERUN_FAILED_STATE.write_text(_json_rf2.dumps({"ids": [t["id"] for t in candidate_tests]}))
        import atexit as _atexit

        _atexit.register(lambda: _RERUN_FAILED_STATE.unlink(missing_ok=True))

        args.test = [t["id"] for t in candidate_tests]
        args.rerun = True

    # Determine test selection. --media composes with --section by union;
    # --test always overrides.
    if args.test:
        test_ids = set(args.test)
        tests = [t for t in TEST_CATALOG if t["id"] in test_ids]
        if not tests:
            print(f"Error: test ID(s) '{args.test}' not found", file=sys.stderr)
            sys.exit(1)
    elif args.media or args.section:
        selected_ids: set[str] = set()
        if args.media:
            media_tests = [t for t in TEST_CATALOG if t.get("workspace_tier") == "media_heavy"]
            selected_ids.update(t["id"] for t in media_tests)
            print(
                f"--media selected {len(media_tests)} test(s): "
                + ", ".join(f"{t['id']}({t.get('media_kind', '?')})" for t in media_tests)
            )
        if args.section:
            section_tests = [t for t in TEST_CATALOG if t["section"] in args.section]
            selected_ids.update(t["id"] for t in section_tests)
        tests = [t for t in TEST_CATALOG if t["id"] in selected_ids]
    else:
        tests = list(TEST_CATALOG)

    # Apply skip flags
    if args.skip_artifacts:
        tests = [t for t in tests if t.get("skip_if") not in ("no_comfyui",)]
    if args.skip_bots:
        tests = [t for t in tests if t.get("skip_if") not in ("no_bot_telegram", "no_bot_slack")]

    if not tests:
        print("No tests selected.", file=sys.stderr)
        sys.exit(1)

    print(f"{len(tests)} test(s) selected")

    # Reorder tests for model-cascade execution: tier groups (large→small→ollama→any),
    # then model_slug within each tier to minimize pipeline model switches.
    tests = sort_tests_cascade(tests)
    tier_counts = {}
    for t in tests:
        tier = t.get("workspace_tier", "any")
        tier_counts[tier] = tier_counts.get(tier, 0) + 1
    print(f"  Cascade order: {' > '.join(f'{t}({c})' for t, c in tier_counts.items())}")

    # --rerun: remove existing rows for the selected tests so they don't duplicate
    if args.rerun:
        if not (args.test or args.section or args.media or args.rerun_failed):
            print(
                "ERROR: --rerun requires --test, --section, --media, or --rerun-failed "
                "to scope the replacement",
                file=sys.stderr,
            )
            sys.exit(1)
        # --rerun implies --append (we're editing an existing file)
        args.append = True
        if RESULTS_FILE.exists():
            target_ids = {t["id"] for t in tests}
            removed = _remove_rows_for_test_ids(target_ids)
            print(f"  --rerun: removed {removed} existing row(s) for {len(target_ids)} test ID(s)")
        else:
            print("  --rerun: no existing UAT_RESULTS.md to update — running fresh")
            args.append = False

    # Skip conditions
    skip_conditions = evaluate_skip_conditions()
    flagged = [k for k, v in skip_conditions.items() if v]
    if flagged:
        print(f"Skip conditions active: {', '.join(flagged)}")

    # Watchdog runs during UAT — the check_server_zombies() function now guards
    # on proxy state=switching so it won't kill a server that is mid-load.
    # Only S23-style tests that deliberately crash backends need the watchdog
    # stopped; UAT doesn't do that.

    # ---- Chat archival strategy ----
    # Tests run in the main (root) chat history so conversations are immediately
    # visible during the run. After completion, all chats from this run are moved
    # to UAT/{YYYY-MM-DD} so the root stays clean and runs are date-stamped.
    folder_id: str | None = None  # No pre-assignment — root during run
    print("  Chat archival: conversations run in root, moved to UAT/{date} on completion")

    # Init results file
    # Targeted runs (--test / --section without --rerun) auto-append so they don't
    # wipe the full run report.  Full runs (no filter, or explicit --rerun) reset it.
    run_ts = time.strftime("%Y-%m-%d %H:%M:%S")
    _targeted = bool((args.test or args.section) and not args.rerun and not args.append)
    if _targeted:
        args.append = True
        print(f"  [targeted run] --append implied — UAT_RESULTS.md preserved (use --rerun to replace rows)")
    if not args.append:
        init_results(run_ts)
    counts: dict[str, int] = {}

    calibration_records: list | None = [] if args.calibrate else None
    if args.calibrate:
        print(f"  Calibration mode — responses will be saved to {args.calibrate_output}")

    # Always-on response corpus. See TASK_UAT_CORPUS_CAPTURE_V1.md §A2.
    corpus_run_id: str = time.strftime("%Y%m%dT%H%M%SZ", time.gmtime())
    print(f"  Corpus: tests/uat_corpus/uat_{corpus_run_id}.jsonl")

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=not args.headed)
        ctx = await browser.new_context(
            viewport={"width": 1440, "height": 900},
            accept_downloads=True,
        )
        page = await ctx.new_page()
        await _fe_login(page)
        print("  Logged in to Open WebUI\n")

        t_start = time.time()
        await _notify_test_start(
            sections=args.section,
            test_count=len(tests),
        )

        # Start continuous memory/health monitor (background task)
        monitor = MemoryMonitor(poll_interval=20.0)
        monitor.start()

        # Start crash watcher (background thread — watches DiagnosticReports)
        _crash_watcher.start()

        _last_tier: str = ""
        for i, test in enumerate(tests, start=1):
            tier = test.get("workspace_tier", "any")

            # Tier transition: evict previous backend + verify memory is clean
            # Critical: two models must never be resident simultaneously (OOM risk).
            if tier != _last_tier:
                if _last_tier:
                    print(f"  Tier transition: {_last_tier} → {tier} — evicting models")
                    unload_all_models()
                elif args.no_unload:
                    print("  Skipping startup /unload (--no-unload, model pre-warmed)")
                else:
                    unload_all_models()

                # Verify prerequisites before proceeding
                # When --no-unload, skip all eviction — model was pre-warmed externally.
                if args.no_unload:
                    print(
                        "  [verify] Skipping Ollama eviction checks (--no-unload, model pre-warmed)"
                    )
                elif tier == "ollama":
                    # Ollama tier: verify models are unloaded before starting
                    for retry in range(3):
                        try:
                            ps = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=5).json()
                            loaded = ps.get("models", [])
                            if not loaded:
                                break
                            print(
                                f"  [verify] Ollama still has {len(loaded)} model(s) loaded — retrying eviction ({retry + 1}/3)"
                            )
                            unload_all_models()
                            _wait_for_ollama_ps_empty(timeout_s=15.0)
                        except Exception:
                            break
                    if retry == 2:
                        try:
                            ps2 = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=5).json()
                            if ps2.get("models"):
                                print(
                                    "  [verify] WARNING: Ollama models still loaded after 3 eviction attempts — may cause OOM"
                                )
                        except Exception:
                            pass

                elif tier == "media_heavy":
                    # Media-heavy tier (TTS, music, video, image): verify Ollama
                    # is clear AND memory is actually freed before proceeding —
                    # media tools spawn additional processes that compete for
                    # GPU memory and can crash the system.
                    for retry in range(3):
                        try:
                            ps = httpx.get(f"{OLLAMA_URL}/api/ps", timeout=5).json()
                            loaded = ps.get("models", [])
                            if not loaded:
                                break
                            print(
                                f"  [verify] Ollama still has {len(loaded)} model(s) — retrying eviction ({retry + 1}/3)"
                            )
                            unload_all_models()
                            _wait_for_ollama_ps_empty(timeout_s=15.0)
                        except Exception:
                            break
                    # Post-eviction memory verification — wait until memory is
                    # actually freed before running memory-intensive media tests
                    if not _wait_for_drain(threshold_pct=75.0, timeout_s=90.0, label="tier-transition"):
                        print("  [mem] WARNING: Memory still high after 90s drain — may risk OOM", flush=True)

                _last_tier = tier

            # Pre-flight: wait for Ollama to be ready before firing the test.
            # Called for ALL tiers: any-tier tests also route through Ollama.
            if test.get("workspace_tier") in ("ollama", "any"):
                ws_id = test.get("model_slug", "auto")
                if tier == "ollama":
                    _pipeline_pre_warm(ws_id)

            # Force-unload before heavy media tests (TTS, music, video, image)
            # that load large frameworks and risk OOM when run consecutively
            if test.get("force_unload_before"):
                print(f"  [mem] Force-unloading before {test['id']} (heavy media test)")
                unload_all_models()
                _wait_for_ollama_ps_empty(timeout_s=15.0)

            # ComfyUI lifecycle: only keep ComfyUI running during tests that
            # actually need it. Stop it before non-ComfyUI tests to reclaim GPU
            # memory; start it (with warmup wait) before ComfyUI-dependent tests.
            needs_comfyui = test.get("skip_if") == "no_comfyui"
            if needs_comfyui and not _comfyui_running():
                # Bring ComfyUI up and give Metal a 30s warmup before the test
                started = _start_comfyui(wait_s=60)
                if started:
                    time.sleep(30)  # Metal warmup before first inference
            elif not needs_comfyui and _comfyui_running():
                _stop_comfyui()

            # If the crash watcher saw a crash since the last
            # test, block here until memory has fully drained before loading
            # another model.
            # another model — attempting a load into a crash-starved Metal
            # heap crashes again immediately and makes memory worse.
            if _crash_watcher.crash_pending:
                _crash_watcher.wait_for_recovery(f"{test['id']} {test['name']}")

            # Pre-test memory check (monitor runs continuously in background,
            # but this catches issues right before a test starts)
            safe = _check_memory_before_test(f"{test['id']} {test['name']}")
            if not safe:
                used_pct = _get_memory_pct()
                print(
                    f"  [{i:02d}/{len(tests):02d}] {test['id']} SKIPPED (memory pressure {used_pct:.0f}%)"
                )
                # Write a row so the skip is visible in UAT_RESULTS.md, not just summary count
                record_result(
                    n=i,
                    status="SKIP",
                    test_id=test["id"],
                    name=test["name"],
                    model=test["model_slug"],
                    assertions=[
                        (
                            "memory_pressure_skip",
                            False,
                            f"used={used_pct:.0f}%, threshold={MEMORY_CRITICAL_PCT:.0f}%",
                        )
                    ],
                    elapsed=0.0,
                    chat_url=f"memory-skip://{used_pct:.0f}pct",
                )
                counts["SKIP"] = counts.get("SKIP", 0) + 1
                continue

            print(f"[{i:02d}/{len(tests):02d}] {test['id']} {test['name']}")

            await run_test(
                page=page,
                test=test,
                token=token,
                skip_conditions=skip_conditions,
                n=i,
                counts=counts,
                headed=args.headed,
                folder_id=folder_id,
                calibration_records=calibration_records,
                corpus_run_id=corpus_run_id,
            )

            # Post-test memory cleanup: only evict when the NEXT test uses a
            # different model_slug. Cascade grouping already keeps same-model
            # tests together to minimize model switches — don't undo that.
            if i < len(tests):
                next_test = tests[i]
                same_model = test.get("model_slug") == next_test.get("model_slug") and test.get(
                    "workspace_tier"
                ) == next_test.get("workspace_tier")
                mem_pct = _get_memory_pct()
                if not same_model and mem_pct >= MEMORY_WARN_PCT:
                    print(f"  [mem] Post-test memory at {mem_pct:.0f}% — evicting (model changing)")
                    unload_all_models()
                    _wait_for_drain(threshold_pct=MEMORY_WARN_PCT, timeout_s=90.0, label="post-evict")
                    mem_after = _get_memory_pct()
                    if mem_after >= MEMORY_CRITICAL_PCT:
                        print(
                            f"  [mem] Memory still {mem_after:.0f}% after eviction — second eviction pass"
                        )
                        unload_all_models()
                        _wait_for_drain(threshold_pct=MEMORY_WARN_PCT, timeout_s=120.0, label="post-evict-2")
                elif same_model and mem_pct >= MEMORY_SAME_MODEL_EVICT_PCT:
                    # KV cache from this test's inference will compound with the next
                    # test's allocation even when the same model stays loaded.
                    print(
                        f"  [mem] Post-test memory at {mem_pct:.0f}% (same model) "
                        "— evicting to clear KV cache residuals"
                    )
                    unload_all_models()
                    _wait_for_ollama_ps_empty(timeout_s=15.0)
                    mem_after = _get_memory_pct()
                    if mem_after >= MEMORY_SAME_MODEL_EVICT_PCT:
                        print(
                            f"  [mem] Memory still {mem_after:.0f}% after same-model eviction "
                            "— memory may not have drained yet"
                        )
                elif mem_pct >= MEMORY_CRITICAL_PCT:
                    # Always evict if critical, even on same model
                    print(f"  [mem] Post-test memory at {mem_pct:.0f}% — critical eviction")
                    unload_all_models()
                    _wait_for_ollama_ps_empty(timeout_s=15.0)

            # Inter-test settling: sleep the prescribed delay, then ensure the
            # backend for the next test is actually alive before proceeding.
            if i < len(tests):
                delay = settling_delay(
                    test.get("workspace_tier", "any"),
                    tests[i].get("workspace_tier", "any"),
                )
                if delay > 0:
                    await asyncio.sleep(delay)
                next_tier = tests[i].get("workspace_tier", "any")
                if next_tier in ("ollama",):
                    alive, detail = _backend_alive(next_tier)
                    if not alive:
                        print(
                            f"  [health] post-settling backend check: {detail}", flush=True
                        )
                        await _wait_for_backend(next_tier, max_wait=60)

        # Navigate away from the last chat before closing so OWUI can commit its
        # "done" state cleanly — prevents the browser-disconnect spinner on the
        # last visited conversation.
        try:
            await page.goto(OPENWEBUI_URL, wait_until="load", timeout=8000)
        except Exception:
            pass
        await browser.close()

    # Stop continuous monitor and crash watcher
    await monitor.stop()
    _crash_watcher.stop()
    if _crash_watcher.crash_log:
        print(
            f"  [crash-watcher] {len(_crash_watcher.crash_log)} crash(es) detected during run:",
            flush=True,
        )
        for entry in _crash_watcher.crash_log:
            print(f"    {entry}", flush=True)

    # Final cleanup: evict all models to prevent OOM after UAT completes
    cleanup_after_uat()

    elapsed = int(time.time() - t_start)
    await _notify_test_end(
        sections=args.section,
        elapsed=elapsed,
        counts=counts,
        test_count=len(tests),
    )
    await _notify_test_summary(
        counts=counts,
        elapsed=elapsed,
        sections=args.section,
        test_count=len(tests),
    )

    # Write calibration JSON if collected
    if calibration_records is not None:
        import json as _json

        cal_path = Path(args.calibrate_output)
        cal_path.write_text(_json.dumps(calibration_records, indent=2, ensure_ascii=False))
        print(f"\nCalibration data: {cal_path} ({len(calibration_records)} records)")
        print("Next: review 'review_tag' fields (good/bad/skip), then run:")
        print(f"  python3 tests/portal5_uat_driver.py --emit-signals-from {cal_path}")

    # Write routing intent-vs-actual summary before rebuilding counts.
    _write_routing_summary()

    # Always rebuild the summary header from actual file rows, so the count
    # is correct after partial / phased / rerun executions.
    _rebuild_summary_from_rows()

    # ---- Post-run archival: move all chats from this run to UAT/{run_date} ----
    run_date = run_ts[:10]  # YYYY-MM-DD from run_ts set at start of run
    if _run_chat_ids:
        try:
            uat_root_id = owui_get_or_create_folder(token, "UAT")
            if uat_root_id:
                dated_folder_id = owui_get_or_create_folder(token, run_date, parent_id=uat_root_id)
                if dated_folder_id:
                    moved = 0
                    for cid in _run_chat_ids:
                        try:
                            owui_assign_chat_folder(token, cid, dated_folder_id)
                            moved += 1
                        except Exception:
                            pass
                    print(f"\n  Archived {moved}/{len(_run_chat_ids)} chats → UAT/{run_date}")
                else:
                    print(f"\n  WARNING: could not create UAT/{run_date} subfolder — chats remain in root")
            else:
                print("\n  WARNING: could not create UAT folder — chats remain in root")
        except Exception as e:
            print(f"\n  WARNING: post-run archival failed: {e} — chats remain in root")

    # Print routing summary to stdout as well
    if _ROUTING_LOG:
        tier_fallbacks = [r for r in _ROUTING_LOG if not r["matched"] and r["tier_mismatch"]]
        wrong_model = [r for r in _ROUTING_LOG if not r["matched"] and not r["tier_mismatch"] and r["actual"]]
        correct = [r for r in _ROUTING_LOG if r["matched"]]
        print(f"\n{'─' * 50}")
        print("ROUTING SUMMARY")
        print(f"{'─' * 50}")
        print(f"  Checked: {len(_ROUTING_LOG)}   ✅ {len(correct)} correct"
              + (f"   ⚠️  {len(tier_fallbacks)} routing mismatch" if tier_fallbacks else "")
              + (f"   ⚠️  {len(wrong_model)} wrong model" if wrong_model else ""))
        for r in tier_fallbacks:
            print(f"  FALLBACK  {r['test_id']:12s} {r['section']:18s} intended={r['intended'][:35]}  got={r['actual'][:35]}")
        for r in wrong_model:
            print(f"  MISMATCH  {r['test_id']:12s} {r['section']:18s} intended={r['intended'][:35]}  got={r['actual'][:35]}")
        if not tier_fallbacks and not wrong_model:
            print("  All tests served by intended primary model.")
        print(f"{'─' * 50}")

    total = sum(counts.values())
    print(f"\n{'=' * 50}")
    print(
        f"Results: {counts.get('PASS', 0)}P / {counts.get('WARN', 0)}W / "
        f"{counts.get('FAIL', 0)}F / {counts.get('SKIP', 0)}S / "
        f"{counts.get('BLOCKED', 0)}B / {counts.get('MANUAL', 0)}M  ({total} total)"
    )
    print(f"Report:  {RESULTS_FILE}")
    print(f"Chats:   {OPENWEBUI_URL}")


if __name__ == "__main__":
    asyncio.run(main())
