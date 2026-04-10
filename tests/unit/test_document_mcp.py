"""tests/unit/test_document_mcp.py

Unit tests for document MCP read tools added in commit 57b2f2c:
  - read_word_document
  - read_excel
  - read_powerpoint
  - read_pdf

All tests run without network or Docker. File I/O tests use tmp_path fixtures.
Tests skip gracefully when heavy deps (python-docx, openpyxl, pptx, pdfplumber)
are not installed — consistent with existing MCP test conventions.
"""

from __future__ import annotations

import sys

import pytest

sys.path.insert(0, ".")

# Guard: skip entire module if portal_mcp.documents is not importable.
pytest.importorskip(
    "portal_mcp.documents.document_mcp",
    reason="portal_mcp.documents not importable — run: pip install -e '.[dev,mcp]'",
)


class TestToolsManifest:
    """TOOLS_MANIFEST must include all four new read tools."""

    def test_manifest_includes_read_word_document(self):
        from portal_mcp.documents.document_mcp import TOOLS_MANIFEST

        names = [t["name"] for t in TOOLS_MANIFEST]
        assert "read_word_document" in names, (
            "TOOLS_MANIFEST missing 'read_word_document' — Open WebUI won't expose this tool"
        )

    def test_manifest_includes_read_excel(self):
        from portal_mcp.documents.document_mcp import TOOLS_MANIFEST

        names = [t["name"] for t in TOOLS_MANIFEST]
        assert "read_excel" in names, "TOOLS_MANIFEST missing 'read_excel'"

    def test_manifest_includes_read_powerpoint(self):
        from portal_mcp.documents.document_mcp import TOOLS_MANIFEST

        names = [t["name"] for t in TOOLS_MANIFEST]
        assert "read_powerpoint" in names, "TOOLS_MANIFEST missing 'read_powerpoint'"

    def test_manifest_includes_read_pdf(self):
        from portal_mcp.documents.document_mcp import TOOLS_MANIFEST

        names = [t["name"] for t in TOOLS_MANIFEST]
        assert "read_pdf" in names, "TOOLS_MANIFEST missing 'read_pdf'"

    def test_manifest_total_count(self):
        """Manifest must have at least 9 tools (5 original + 4 read tools)."""
        from portal_mcp.documents.document_mcp import TOOLS_MANIFEST

        assert len(TOOLS_MANIFEST) >= 9, (
            f"Expected at least 9 tools in TOOLS_MANIFEST, got {len(TOOLS_MANIFEST)}"
        )

    def test_read_tools_have_file_path_parameter(self):
        """All read tools must declare file_path as a required parameter."""
        from portal_mcp.documents.document_mcp import TOOLS_MANIFEST

        read_tools = [t for t in TOOLS_MANIFEST if t["name"].startswith("read_")]
        assert len(read_tools) == 4, f"Expected 4 read_ tools, got {len(read_tools)}"
        for tool in read_tools:
            props = tool.get("parameters", {}).get("properties", {})
            assert "file_path" in props, (
                f"Tool '{tool['name']}' missing 'file_path' in parameters"
            )


class TestReadWordDocument:
    """Tests for read_word_document function."""

    def test_missing_file_returns_error(self, tmp_path):
        from portal_mcp.documents.document_mcp import read_word_document

        result = read_word_document(str(tmp_path / "nonexistent.docx"))
        assert result["success"] is False
        assert "error" in result

    def test_wrong_extension_returns_error(self, tmp_path):
        from portal_mcp.documents.document_mcp import read_word_document

        f = tmp_path / "doc.txt"
        f.write_text("hello")
        result = read_word_document(str(f))
        assert result["success"] is False
        assert "error" in result

    def test_valid_docx_returns_success(self, tmp_path):
        """Test with a real .docx if python-docx is installed, skip otherwise."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        from portal_mcp.documents.document_mcp import read_word_document

        doc = Document()
        doc.add_paragraph("Hello Portal 5")
        path = tmp_path / "test.docx"
        doc.save(str(path))

        result = read_word_document(str(path))
        assert result["success"] is True
        assert result["filename"] == "test.docx"
        assert "content" in result
        assert any("Hello Portal 5" in block["text"] for block in result["content"])

    def test_missing_dep_returns_graceful_error(self, tmp_path, monkeypatch):
        """Simulate python-docx not installed — must return error dict, not raise."""
        import builtins

        real_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "docx":
                raise ImportError("mocked missing python-docx")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)

        # Create a valid .docx path so we reach the import check
        f = tmp_path / "doc.docx"
        f.touch()

        from portal_mcp.documents import document_mcp
        # Call the raw function to bypass any module-level caching
        result = document_mcp.read_word_document(str(f))
        assert result["success"] is False
        assert "error" in result


class TestReadExcel:
    """Tests for read_excel function."""

    def test_missing_file_returns_error(self, tmp_path):
        from portal_mcp.documents.document_mcp import read_excel

        result = read_excel(str(tmp_path / "nonexistent.xlsx"))
        assert result["success"] is False
        assert "error" in result

    def test_wrong_extension_returns_error(self, tmp_path):
        from portal_mcp.documents.document_mcp import read_excel

        f = tmp_path / "file.csv"
        f.write_text("a,b,c")
        result = read_excel(str(f))
        assert result["success"] is False
        assert "error" in result

    def test_valid_xlsx_returns_success(self, tmp_path):
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl not installed")

        from portal_mcp.documents.document_mcp import read_excel

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["Name", "Value"])
        ws.append(["Portal5", "42"])
        path = tmp_path / "test.xlsx"
        wb.save(str(path))

        result = read_excel(str(path))
        assert result["success"] is True
        assert result["filename"] == "test.xlsx"
        assert result["sheet_count"] >= 1
        assert "sheets" in result

    def test_max_rows_limits_output(self, tmp_path):
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl not installed")

        from portal_mcp.documents.document_mcp import read_excel

        wb = openpyxl.Workbook()
        ws = wb.active
        for i in range(100):
            ws.append([i, i * 2])
        path = tmp_path / "big.xlsx"
        wb.save(str(path))

        result = read_excel(str(path), max_rows=10)
        assert result["success"] is True
        sheet = result["sheets"][0]
        assert sheet["row_count"] <= 10
        assert sheet["truncated"] is True


class TestReadPowerPoint:
    """Tests for read_powerpoint function."""

    def test_missing_file_returns_error(self, tmp_path):
        from portal_mcp.documents.document_mcp import read_powerpoint

        result = read_powerpoint(str(tmp_path / "nonexistent.pptx"))
        assert result["success"] is False
        assert "error" in result

    def test_wrong_extension_returns_error(self, tmp_path):
        from portal_mcp.documents.document_mcp import read_powerpoint

        f = tmp_path / "slides.pdf"
        f.touch()
        result = read_powerpoint(str(f))
        assert result["success"] is False
        assert "error" in result

    def test_valid_pptx_returns_success(self, tmp_path):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        from portal_mcp.documents.document_mcp import read_powerpoint

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Test Slide"
        path = tmp_path / "test.pptx"
        prs.save(str(path))

        result = read_powerpoint(str(path))
        assert result["success"] is True
        assert result["filename"] == "test.pptx"
        assert result["slide_count"] >= 1
        assert "slides" in result


class TestReadPDF:
    """Tests for read_pdf function."""

    def test_missing_file_returns_error(self, tmp_path):
        from portal_mcp.documents.document_mcp import read_pdf

        result = read_pdf(str(tmp_path / "nonexistent.pdf"))
        assert result["success"] is False
        assert "error" in result

    def test_wrong_extension_returns_error(self, tmp_path):
        from portal_mcp.documents.document_mcp import read_pdf

        f = tmp_path / "doc.docx"
        f.touch()
        result = read_pdf(str(f))
        assert result["success"] is False
        assert "error" in result

    def test_missing_pdfplumber_returns_graceful_error(self, tmp_path):
        """pdfplumber may not be installed — must return error dict, never raise."""
        import builtins

        real_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "pdfplumber":
                raise ImportError("mocked missing pdfplumber")
            return real_import(name, *args, **kwargs)

        f = tmp_path / "doc.pdf"
        f.touch()


        import unittest.mock as mock

        import portal_mcp.documents.document_mcp as doc_mcp

        with mock.patch.dict("sys.modules", {"pdfplumber": None}):
            result = doc_mcp.read_pdf(str(f))
        assert result["success"] is False
        assert "error" in result

    def test_valid_pdf_returns_success(self, tmp_path):
        """Test with a real PDF if pdfplumber is installed, skip otherwise."""
        try:
            import pdfplumber
        except ImportError:
            pytest.skip("pdfplumber not installed")

        # Create a minimal valid PDF using reportlab if available, else skip
        try:
            from reportlab.pdfgen import canvas as rl_canvas
        except ImportError:
            pytest.skip("reportlab not installed (needed to create test PDF)")

        from portal_mcp.documents.document_mcp import read_pdf

        path = tmp_path / "test.pdf"
        c = rl_canvas.Canvas(str(path))
        c.drawString(100, 750, "Hello from Portal 5")
        c.save()

        result = read_pdf(str(path))
        assert result["success"] is True
        assert result["filename"] == "test.pdf"
        assert "pages" in result
