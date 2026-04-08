"""
Document Tools MCP Server
Exposes Word, PowerPoint, and Excel document creation as MCP tools.
Generated files are saved to OUTPUT_DIR with unique IDs.

Requires: pip install python-docx python-pptx openpyxl
Start with: python -m mcp.documents.document_mcp
"""

import logging
import os
import uuid
from pathlib import Path

from starlette.responses import JSONResponse

from portal_mcp.mcp_server.fastmcp import FastMCP

port = int(os.getenv("DOCUMENTS_MCP_PORT", "8913"))
mcp = FastMCP("document-tools", port=port)


@mcp.custom_route("/health", methods=["GET"])
async def health_check(request):
    return JSONResponse({"status": "ok", "service": "documents-mcp"})


# Tool manifest for discovery
TOOLS_MANIFEST = [
    {
        "name": "create_word_document",
        "description": "Create a Word document",
        "parameters": {
            "type": "object",
            "properties": {
                "content": {
                    "type": "string",
                    "description": "Document content (markdown supported)",
                },
                "title": {"type": "string", "description": "Document title"},
            },
            "required": ["content"],
        },
    },
    {
        "name": "create_powerpoint",
        "description": "Create a PowerPoint presentation",
        "parameters": {
            "type": "object",
            "properties": {
                "slides": {"type": "array", "description": "List of slide content"},
                "title": {"type": "string", "description": "Presentation title"},
            },
            "required": ["slides"],
        },
    },
    {
        "name": "create_excel",
        "description": "Create an Excel spreadsheet",
        "parameters": {
            "type": "object",
            "properties": {
                "data": {"type": "array", "description": "Array of rows"},
                "sheet_name": {"type": "string", "description": "Sheet name"},
            },
            "required": ["data"],
        },
    },
    {
        "name": "convert_document",
        "description": "Convert between document formats using pandoc",
        "parameters": {
            "type": "object",
            "properties": {
                "source_path": {
                    "type": "string",
                    "description": "Path to source file (.docx, .pptx, .xlsx, or .pdf)",
                },
                "target_format": {
                    "type": "string",
                    "description": "Target format: 'pdf', 'docx', 'pptx', or 'xlsx'",
                },
            },
            "required": ["source_path", "target_format"],
        },
    },
    {
        "name": "list_generated_files",
        "description": "List recently generated documents in the output directory. "
        "Use this to find files created by other tools.",
        "parameters": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "read_word_document",
        "description": "Extract text content and structure from an existing Word (.docx) file. "
        "Returns headings, paragraphs, and table data.",
        "parameters": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Absolute path to the .docx file to read",
                },
                "include_tables": {
                    "type": "boolean",
                    "description": "Whether to include table cell content (default true)",
                },
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "read_excel",
        "description": "Extract data from an existing Excel (.xlsx) spreadsheet. "
        "Returns sheet names and row data for each sheet.",
        "parameters": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Absolute path to the .xlsx file to read",
                },
                "max_rows": {
                    "type": "integer",
                    "description": "Maximum rows to return per sheet (default 500)",
                },
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "read_powerpoint",
        "description": "Extract text and speaker notes from an existing PowerPoint (.pptx) file. "
        "Returns slide titles, content, and notes.",
        "parameters": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Absolute path to the .pptx file to read",
                },
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "read_pdf",
        "description": "Extract text content from an existing PDF file page by page. "
        "Also extracts tables when present. Requires pdfplumber.",
        "parameters": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Absolute path to the .pdf file to read",
                },
                "max_pages": {
                    "type": "integer",
                    "description": "Maximum pages to extract (default 50, 0 = all)",
                },
                "include_tables": {
                    "type": "boolean",
                    "description": "Whether to extract table data alongside text (default true)",
                },
            },
            "required": ["file_path"],
        },
    },
]


@mcp.custom_route("/tools", methods=["GET"])
async def list_tools(request):
    return JSONResponse({"tools": TOOLS_MANIFEST})


logger = logging.getLogger(__name__)

OUTPUT_DIR = Path(os.getenv("OUTPUT_DIR", "data/generated"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def _unique_path(name: str, ext: str) -> Path:
    """Return a unique output path."""
    uid = uuid.uuid4().hex[:8]
    safe = "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in name[:40]).strip("_")
    return OUTPUT_DIR / f"{safe}_{uid}.{ext}"


@mcp.tool()
def create_word_document(
    title: str,
    content: str,
    author: str = "Portal AI",
) -> dict:
    """
    Create a Word (.docx) document from a title and markdown-style content.

    Content supports:
    - '# Heading' → H1 heading
    - '## Heading' → H2 heading
    - '### Heading' → H3 heading
    - '- item' → bullet list item
    - Regular text → body paragraph

    Args:
        title: Document title (also used as filename base)
        content: Document body; supports basic markdown headings and bullets
        author: Author name for document metadata (default "Portal AI")

    Returns:
        dict with success, path (server path), and filename
    """
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        return {
            "success": False,
            "error": "python-docx not installed. Run: pip install python-docx",
        }

    try:
        doc = Document()
        doc.core_properties.author = author
        doc.core_properties.title = title

        # Title heading
        heading = doc.add_heading(title, level=0)
        heading.runs[0].font.size = Pt(24)

        for line in content.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                para = doc.add_paragraph(stripped[2:], style="List Bullet")
                para.paragraph_format.space_before = Pt(0)
            else:
                doc.add_paragraph(stripped)

        output_path = _unique_path(title, "docx")
        doc.save(str(output_path))
        return {"success": True, "path": str(output_path), "filename": output_path.name}
    except Exception as e:
        logger.exception("Word document creation failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def create_powerpoint(
    title: str,
    slides: list[dict],
    author: str = "Portal AI",
) -> dict:
    """
    Create a PowerPoint (.pptx) presentation.

    Each slide dict should have:
    - 'title': slide title (str)
    - 'content': slide body text or bullet points (str, newline-separated)
    - 'notes': speaker notes (str, optional)

    Args:
        title: Presentation title (used as filename base)
        slides: List of slide dicts with 'title', 'content', and optional 'notes'
        author: Author name for metadata (default "Portal AI")

    Returns:
        dict with success, path (server path), and filename
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed. Run: pip install python-pptx",
        }

    try:
        prs = Presentation()
        prs.core_properties.author = author
        prs.core_properties.title = title

        # Title slide
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = title

        # Content slides
        content_layout = prs.slide_layouts[1]
        for slide_data in slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            for i, line in enumerate(slide_data.get("content", "").splitlines()):
                stripped = line.strip()
                if not stripped:
                    continue
                if i == 0:
                    tf.text = stripped
                else:
                    tf.add_paragraph().text = stripped

            if "notes" in slide_data and slide_data["notes"]:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data["notes"]

        output_path = _unique_path(title, "pptx")
        prs.save(str(output_path))
        return {"success": True, "path": str(output_path), "filename": output_path.name}
    except Exception as e:
        logger.exception("Presentation creation failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def create_excel(
    title: str,
    data: list | None = None,
    sheets: list[dict] | None = None,
    sheet_name: str = "Sheet1",
) -> dict:
    """
    Create an Excel (.xlsx) spreadsheet.

    Simple usage (flat rows):
        data: List of rows (first row treated as headers), e.g. [["Name","Score"],["Alice",95]]
        sheet_name: Tab name (default "Sheet1")

    Advanced usage (multiple sheets):
        sheets: List of sheet dicts, each with 'name', 'headers', and 'rows'

    Args:
        title: Spreadsheet title (used as filename base)
        data: Simple list of rows (first row is headers)
        sheets: Advanced — list of sheet dicts with 'name', 'headers', 'rows'
        sheet_name: Sheet tab name when using data parameter (default "Sheet1")

    Returns:
        dict with success, path (server path), and filename
    """
    try:
        import openpyxl
        from openpyxl.styles import Font
    except ImportError:
        return {
            "success": False,
            "error": "openpyxl not installed. Run: pip install openpyxl",
        }

    try:
        wb = openpyxl.Workbook()
        wb.properties.title = title

        if data is not None:
            # Simple mode: flat list of rows
            ws = wb.active
            ws.title = sheet_name
            for i, row in enumerate(data):
                ws.append(row)
                if i == 0:
                    # Bold first row (headers)
                    for cell in ws[1]:
                        cell.font = Font(bold=True)
        elif sheets is not None:
            default_sheet = wb.active
            first = True
            for sheet_data in sheets:
                if first:
                    ws = default_sheet
                    ws.title = sheet_data.get("name", "Sheet1")
                    first = False
                else:
                    ws = wb.create_sheet(title=sheet_data.get("name", "Sheet"))

                headers = sheet_data.get("headers", [])
                if headers:
                    ws.append(headers)
                    for cell in ws[1]:
                        cell.font = Font(bold=True)

                for row in sheet_data.get("rows", []):
                    ws.append(row)
        else:
            return {"success": False, "error": "Provide either 'data' or 'sheets' parameter"}

        output_path = _unique_path(title, "xlsx")
        wb.save(str(output_path))
        return {"success": True, "path": str(output_path), "filename": output_path.name}
    except Exception as e:
        logger.exception("Spreadsheet creation failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def convert_document(
    source_path: str,
    target_format: str,
) -> dict:
    """Copy a document to a new format name.

    Note: True format conversion (e.g., .docx to .pdf) requires LibreOffice
    installed on the host. Without LibreOffice, this tool copies the file with
    the new extension, which is only useful for same-family formats.

    For PDF export, use LibreOffice on the host:
      libreoffice --headless --convert-to pdf <file>

    Args:
        source_path:   Path to the source document
        target_format: Target extension: 'pdf', 'docx', 'pptx', or 'xlsx'
    """
    import shutil
    import subprocess
    from pathlib import Path as _Path

    _allowed_target_formats = frozenset({"pdf", "docx", "pptx", "xlsx"})

    src = _Path(source_path).resolve()
    allowed_root = _Path(OUTPUT_DIR).resolve()
    if not str(src).startswith(str(allowed_root) + os.sep):
        return {"error": "source_path must be a file within the output directory"}
    if not src.exists():
        return {"error": f"Source file not found: {source_path}"}

    target_format = target_format.lower().lstrip(".")
    if target_format not in _allowed_target_formats:
        return {
            "error": (
                f"Unsupported target_format {target_format!r}. "
                f"Allowed: {sorted(_allowed_target_formats)}"
            )
        }
    out_path = _unique_path(src.stem, target_format)

    # Attempt LibreOffice conversion for cross-format (best quality)
    try:
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                target_format,
                "--outdir",
                str(out_path.parent),
                str(src),
            ],
            capture_output=True,
            timeout=60,
        )
        if result.returncode == 0:
            # LibreOffice writes to same dir as source — find the output file
            converted = src.parent / f"{src.stem}.{target_format}"
            if converted.exists():
                shutil.move(str(converted), str(out_path))
                return {"success": True, "path": str(out_path), "method": "libreoffice"}
    except (FileNotFoundError, subprocess.TimeoutExpired):
        logger.debug(
            "LibreOffice conversion failed (not installed or timed out) — using copy fallback"
        )
    except OSError as e:
        logger.debug("LibreOffice subprocess error: %s", e)

    # Fallback: copy with new extension (only meaningful for same-family formats)
    same_family = {
        frozenset({"docx", "doc"}),
        frozenset({"pptx", "ppt"}),
        frozenset({"xlsx", "xls"}),
    }
    src_ext = src.suffix.lstrip(".")
    is_same_family = any({src_ext, target_format} <= fam for fam in same_family)

    if not is_same_family:
        return {
            "error": (
                f"Cannot convert {src_ext!r} → {target_format!r} without LibreOffice. "
                "Install LibreOffice for PDF and cross-format conversion."
            ),
            "install": "brew install libreoffice  # or apt-get install libreoffice",
        }

    shutil.copy2(str(src), str(out_path))
    return {
        "success": True,
        "path": str(out_path),
        "method": "copy",
        "note": f"Copied {src_ext} → {target_format}. "
        "Install LibreOffice for true format conversion.",
    }


@mcp.tool()
def list_generated_files() -> list[dict]:
    """List recently generated documents in the output directory."""
    files = []
    for f in sorted(OUTPUT_DIR.iterdir(), key=lambda x: x.stat().st_mtime, reverse=True)[:20]:
        if f.is_file():
            files.append(
                {
                    "filename": f.name,
                    "path": str(f),
                    "size_bytes": f.stat().st_size,
                    "type": f.suffix.lstrip("."),
                }
            )
    return files


@mcp.tool()
def read_word_document(
    file_path: str,
    include_tables: bool = True,
) -> dict:
    """Extract text content and structure from an existing Word (.docx) file.

    Returns headings, paragraphs, and table data.

    Args:
        file_path:      Absolute path to the .docx file to read
        include_tables: Whether to include table cell content (default True)

    Returns:
        dict with success, metadata, content (list of blocks), and optional tables
    """
    try:
        from docx import Document as _Document
    except ImportError:
        return {
            "success": False,
            "error": "python-docx not installed. Run: pip install python-docx",
        }

    src = Path(file_path).resolve()
    if not src.exists():
        return {"success": False, "error": f"File not found: {file_path}"}
    if src.suffix.lower() != ".docx":
        return {"success": False, "error": f"Expected .docx file, got: {src.suffix}"}

    try:
        doc = _Document(str(src))
        blocks: list[dict] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else "Normal"
            blocks.append({"type": style, "text": text})

        result: dict = {
            "success": True,
            "filename": src.name,
            "author": doc.core_properties.author or "",
            "title": doc.core_properties.title or "",
            "paragraph_count": len(blocks),
            "content": blocks,
        }

        if include_tables and doc.tables:
            tables_data = []
            for i, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                tables_data.append({"table_index": i, "rows": rows})
            result["tables"] = tables_data

        return result
    except Exception as e:
        logger.exception("Word document read failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def read_excel(
    file_path: str,
    max_rows: int = 500,
) -> dict:
    """Extract data from an existing Excel (.xlsx) spreadsheet.

    Returns sheet names and row data for each sheet.

    Args:
        file_path: Absolute path to the .xlsx file to read
        max_rows:  Maximum rows to return per sheet (default 500)

    Returns:
        dict with success, filename, sheet_count, and sheets (list of sheet dicts)
    """
    try:
        import openpyxl as _openpyxl
    except ImportError:
        return {"success": False, "error": "openpyxl not installed. Run: pip install openpyxl"}

    src = Path(file_path).resolve()
    if not src.exists():
        return {"success": False, "error": f"File not found: {file_path}"}
    if src.suffix.lower() not in {".xlsx", ".xlsm"}:
        return {"success": False, "error": f"Expected .xlsx/.xlsm file, got: {src.suffix}"}

    try:
        wb = _openpyxl.load_workbook(str(src), read_only=True, data_only=True)
        sheets_data = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[list] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if max_rows and i >= max_rows:
                    break
                rows.append([str(cell) if cell is not None else "" for cell in row])
            sheets_data.append(
                {
                    "name": sheet_name,
                    "row_count": len(rows),
                    "truncated": max_rows > 0 and len(rows) >= max_rows,
                    "rows": rows,
                }
            )
        wb.close()

        return {
            "success": True,
            "filename": src.name,
            "sheet_count": len(sheets_data),
            "sheets": sheets_data,
        }
    except Exception as e:
        logger.exception("Excel read failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def read_powerpoint(
    file_path: str,
) -> dict:
    """Extract text and speaker notes from an existing PowerPoint (.pptx) file.

    Returns slide titles, content blocks, and notes.

    Args:
        file_path: Absolute path to the .pptx file to read

    Returns:
        dict with success, filename, slide_count, and slides (list of slide dicts)
    """
    try:
        from pptx import Presentation as _Presentation
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed. Run: pip install python-pptx",
        }

    src = Path(file_path).resolve()
    if not src.exists():
        return {"success": False, "error": f"File not found: {file_path}"}
    if src.suffix.lower() != ".pptx":
        return {"success": False, "error": f"Expected .pptx file, got: {src.suffix}"}

    try:
        prs = _Presentation(str(src))
        slides_data = []
        for i, slide in enumerate(prs.slides):
            title = ""
            content_blocks: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                shape_text = shape.text_frame.text.strip()
                if not shape_text:
                    continue
                if shape.shape_id == 1 or (
                    hasattr(slide.shapes, "title") and shape == slide.shapes.title
                ):
                    title = shape_text
                else:
                    content_blocks.append(shape_text)

            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slides_data.append(
                {
                    "slide_number": i + 1,
                    "title": title,
                    "content": content_blocks,
                    "notes": notes,
                }
            )

        return {
            "success": True,
            "filename": src.name,
            "author": prs.core_properties.author or "",
            "slide_count": len(slides_data),
            "slides": slides_data,
        }
    except Exception as e:
        logger.exception("PowerPoint read failed")
        return {"success": False, "error": str(e)}


@mcp.tool()
def read_pdf(
    file_path: str,
    max_pages: int = 50,
    include_tables: bool = True,
) -> dict:
    """Extract text content from an existing PDF file page by page.

    Also extracts tables when present. Requires pdfplumber.

    Args:
        file_path:      Absolute path to the .pdf file to read
        max_pages:      Maximum pages to extract (default 50, 0 = all pages)
        include_tables: Whether to extract table data alongside text (default True)

    Returns:
        dict with success, filename, page_count, pages (list of page dicts),
        and optional tables per page
    """
    try:
        import pdfplumber as _pdfplumber
    except ImportError:
        return {
            "success": False,
            "error": "pdfplumber not installed. Run: pip install pdfplumber",
        }

    src = Path(file_path).resolve()
    if not src.exists():
        return {"success": False, "error": f"File not found: {file_path}"}
    if src.suffix.lower() != ".pdf":
        return {"success": False, "error": f"Expected .pdf file, got: {src.suffix}"}

    try:
        pages_data = []
        with _pdfplumber.open(str(src)) as pdf:
            total_pages = len(pdf.pages)
            limit = total_pages if max_pages == 0 else min(max_pages, total_pages)

            for i in range(limit):
                page = pdf.pages[i]
                text = (page.extract_text() or "").strip()
                page_dict: dict = {
                    "page_number": i + 1,
                    "text": text,
                    "char_count": len(text),
                }

                if include_tables:
                    raw_tables = page.extract_tables()
                    if raw_tables:
                        page_dict["tables"] = [
                            [[str(cell) if cell is not None else "" for cell in row] for row in tbl]
                            for tbl in raw_tables
                        ]

                pages_data.append(page_dict)

        return {
            "success": True,
            "filename": src.name,
            "total_pages": total_pages,
            "pages_extracted": len(pages_data),
            "truncated": len(pages_data) < total_pages,
            "pages": pages_data,
        }
    except Exception as e:
        logger.exception("PDF read failed")
        return {"success": False, "error": str(e)}


if __name__ == "__main__":
    mcp.settings.host = "0.0.0.0"
    mcp.run(transport="streamable-http")
